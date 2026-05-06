from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Callable

import pandas as pd
from fastapi.testclient import TestClient
from pptx import Presentation

from nps_lens.api.app import create_app
from nps_lens.domain.helix_links import build_helix_incident_url_lookup, enrich_helix_incident_links
from nps_lens.domain.models import UploadContext
from nps_lens.reports import BusinessPptResult
from nps_lens.settings import Settings
from nps_lens.testing.fixtures import fixture_excel


def _settings(tmp_path: Path) -> Settings:
    return Settings(
        data_dir=tmp_path / "data",
        database_path=tmp_path / "data" / "dashboard.sqlite3",
        frontend_dist_dir=tmp_path / "frontend-dist",
        frontend_public_dir=tmp_path / "frontend-public",
        api_host="127.0.0.1",
        api_port=8000,
        default_service_origin="BBVA México",
        default_service_origin_n1="Senda",
        allowed_service_origins=["BBVA México"],
        allowed_service_origin_n1={"BBVA México": ["Senda"]},
        log_level="INFO",
    )


def _upload_nps_march(client: TestClient) -> dict[str, object]:
    march = fixture_excel("NPS Térmico Senda - 03Marzo.xlsx")
    with march.open("rb") as handle:
        response = client.post(
            "/api/uploads/nps",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    march.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "completed"
    return payload


def _upload_nps_jan_feb(client: TestClient) -> dict[str, object]:
    jan_feb = fixture_excel("NPS Térmico Senda - 01Enero-02Febrero.xlsx")
    with jan_feb.open("rb") as handle:
        response = client.post(
            "/api/uploads/nps",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    jan_feb.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "completed"
    return payload


def _build_helix_fixture(path: Path) -> Path:
    pd.DataFrame(
        {
            "BBVA_SourceServiceCompany": ["BBVA México", "BBVA México", "BBVA España"],
            "BBVA_SourceServiceN1": ["Senda", "Senda", "Senda"],
            "BBVA_SourceServiceN2": ["", "", ""],
            "CreatedDate": ["2026-03-01", "2026-03-03", "2026-03-04"],
            "Incident Number": ["INC-1", "INC-2", "INC-3"],
            "Record ID": ["RID-1", "RID-2", "RID-3"],
            "Detailed Description": [
                "Cliente no puede acceder al portal",
                "Fallo en autenticacion web",
                "Contexto ajeno",
            ],
            "Short Description": ["Acceso", "Autenticacion", "Otro"],
        }
    ).to_excel(path, index=False)
    return path


def _build_helix_out_of_period_fixture(path: Path) -> Path:
    pd.DataFrame(
        {
            "BBVA_SourceServiceCompany": ["BBVA México", "BBVA México"],
            "BBVA_SourceServiceN1": ["Senda", "Senda"],
            "BBVA_SourceServiceN2": ["", ""],
            "CreatedDate": ["2026-02-01", "2026-02-03"],
            "Incident Number": ["INC-FEB-1", "INC-FEB-2"],
            "Record ID": ["RID-FEB-1", "RID-FEB-2"],
            "Detailed Description": [
                "Cliente no puede acceder al portal en febrero",
                "Fallo en autenticacion web en febrero",
            ],
            "Short Description": ["Acceso febrero", "Autenticacion febrero"],
        }
    ).to_excel(path, index=False)
    return path


def _persist_report_in_tmp(tmp_path: Path) -> Callable[[BusinessPptResult], Path]:
    def _persist(report: BusinessPptResult) -> Path:
        target = tmp_path / report.file_name
        target.write_bytes(report.content)
        return target

    return _persist


def _ppt_texts(content: bytes) -> list[str]:
    prs = Presentation(BytesIO(content))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")
    return texts


def test_dashboard_context_nps_and_dataset_views_are_restored(tmp_path: Path) -> None:
    app = create_app(_settings(tmp_path))
    client = TestClient(app)
    upload = _upload_nps_march(client)

    context_response = client.get(
        "/api/dashboard/context",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
        },
    )
    assert context_response.status_code == 200
    context_payload = context_response.json()
    assert "2026" in context_payload["available_years"]
    assert "03" in context_payload["available_months_by_year"]["2026"]
    assert "Web" in context_payload["score_channels"]
    assert context_payload["nps_dataset"]["available"] is True
    assert context_payload["nps_dataset"]["rows"] == upload["inserted_rows"]
    assert "Downloads" in context_payload["preferences"]["downloads_path"]
    assert (
        context_payload["preferences"]["helix_base_url"]
        == "https://itsmhelixbbva-smartit.onbmc.com/smartit/app/#/incidentPV/"
    )
    assert context_payload["preferences"]["report_dimension_analysis"] == "palanca"
    assert context_payload["preferences"]["touchpoint_source"] == "palanca_touchpoint"
    assert any(
        option["value"] == "executive_journeys"
        for option in context_payload["causal_method_options"]
    )

    dashboard_response = client.get(
        "/api/dashboard/nps",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "nps_group": "Todos",
        },
    )
    assert dashboard_response.status_code == 200
    dashboard_payload = dashboard_response.json()
    assert dashboard_payload["context_label"]
    assert dashboard_payload["kpis"]["samples"] > 0
    assert dashboard_payload["kpis"]["neutral_rate"] is not None
    assert (
        abs(
            dashboard_payload["gaps"]["overall_nps"]
            - dashboard_payload["scope"]["period"]["kpis"]["classic_nps"]
        )
        < 1e-9
    )
    assert dashboard_payload["scope"]["cumulative"]["label"].startswith("Datos acumulados hasta")
    assert dashboard_payload["overview"]["daily_volume_mix_figure"] is not None
    assert dashboard_payload["overview"]["topics_table"] is not None
    assert dashboard_payload["controls"]["dimensions"] == [
        "Palanca",
        "Subpalanca",
        "Canal",
        "UsuarioDecisión",
    ]
    assert "report_markdown" not in dashboard_payload

    records = app.state.repository.load_records_df(
        UploadContext(
            service_origin="BBVA México",
            service_origin_n1="Senda",
            service_origin_n2="",
        )
    )
    scope_records = app.state.dashboard_service._apply_population_filters(
        records,
        "2026",
        "03",
    )
    scores = pd.to_numeric(scope_records["NPS"], errors="coerce").dropna()
    expected_neutral_rate = float(((scores >= 7) & (scores <= 8)).mean())
    assert abs(dashboard_payload["kpis"]["neutral_rate"] - expected_neutral_rate) < 1e-9

    filtered_dashboard_response = client.get(
        "/api/dashboard/nps",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "score_channel": "Web",
            "nps_group": "Detractores",
        },
    )
    assert filtered_dashboard_response.status_code == 200
    filtered_dashboard_payload = filtered_dashboard_response.json()
    assert filtered_dashboard_payload["kpis"] == dashboard_payload["kpis"]

    data_response = client.get(
        "/api/dashboard/data/nps",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "nps_group": "Todos",
            "limit": 5,
        },
    )
    assert data_response.status_code == 200
    data_payload = data_response.json()
    assert data_payload["dataset_kind"] == "nps"
    assert data_payload["total_rows"] > 0
    assert "Browser" in data_payload["columns"]
    assert len(data_payload["rows"]) == 5


def test_generate_ppt_report_with_valid_nps_and_no_helix_omits_causal_section(
    tmp_path: Path,
    monkeypatch,
) -> None:
    app = create_app(_settings(tmp_path))
    client = TestClient(app)
    _upload_nps_march(client)
    service = app.state.dashboard_service
    monkeypatch.setattr(service, "_persist_report_copy", _persist_report_in_tmp(tmp_path))

    report = service.generate_ppt_report(
        context=UploadContext(
            service_origin="BBVA México",
            service_origin_n1="Senda",
            service_origin_n2="",
        ),
        pop_year="2026",
        pop_month="03",
        nps_group="Todos",
        score_channel="Web",
    )

    assert report.content
    assert report.slide_count > 0
    texts = _ppt_texts(report.content)
    assert not any("Análisis causal no concluyente" in text for text in texts)
    assert not any("Journeys de detracción" in text for text in texts)


def test_generate_ppt_report_with_helix_outside_period_omits_causal_section(
    tmp_path: Path,
    monkeypatch,
) -> None:
    app = create_app(_settings(tmp_path))
    client = TestClient(app)
    _upload_nps_march(client)
    helix_fixture = _build_helix_out_of_period_fixture(tmp_path / "helix-february.xlsx")
    with helix_fixture.open("rb") as handle:
        upload_response = client.post(
            "/api/uploads/helix",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    helix_fixture.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )
    assert upload_response.status_code == 200
    service = app.state.dashboard_service
    monkeypatch.setattr(service, "_persist_report_copy", _persist_report_in_tmp(tmp_path))

    report = service.generate_ppt_report(
        context=UploadContext(
            service_origin="BBVA México",
            service_origin_n1="Senda",
            service_origin_n2="",
        ),
        pop_year="2026",
        pop_month="03",
        nps_group="Todos",
        score_channel="Web",
    )

    assert report.content
    assert report.slide_count > 0
    texts = _ppt_texts(report.content)
    assert not any("Análisis causal no concluyente" in text for text in texts)
    assert not any("Journeys de detracción" in text for text in texts)


def test_dashboard_supports_helix_upload_and_contextual_table(tmp_path: Path) -> None:
    client = TestClient(create_app(_settings(tmp_path)))
    _upload_nps_march(client)
    helix_fixture = _build_helix_fixture(tmp_path / "helix.xlsx")

    with helix_fixture.open("rb") as handle:
        upload_response = client.post(
            "/api/uploads/helix",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    helix_fixture.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )

    assert upload_response.status_code == 200
    upload_payload = upload_response.json()
    assert upload_payload["status"] == "completed"
    assert upload_payload["row_count"] == 2
    assert upload_payload["dataset"]["available"] is True

    context_response = client.get(
        "/api/dashboard/context",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
        },
    )
    context_payload = context_response.json()
    assert context_payload["helix_dataset"]["available"] is True
    assert context_payload["helix_dataset"]["rows"] == 2
    assert pd.notna(pd.to_datetime(context_payload["helix_dataset"]["updated_at"], errors="coerce"))

    data_response = client.get(
        "/api/dashboard/data/helix",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "limit": 10,
        },
    )
    assert data_response.status_code == 200
    data_payload = data_response.json()
    assert data_payload["dataset_kind"] == "helix"
    assert data_payload["total_rows"] == 2
    assert data_payload["columns"][:8] == [
        "BBVA_SourceServiceCompany",
        "BBVA_SourceServiceN1",
        "BBVA_SourceServiceN2",
        "CreatedDate",
        "Incident Number",
        "Record ID",
        "Detailed Description",
        "Short Description",
    ]
    assert data_payload["rows"][0]["BBVA_SourceServiceN1"] == "Senda"
    assert (
        data_payload["rows"][0]["Incident Number__href"]
        == "https://itsmhelixbbva-smartit.onbmc.com/smartit/app/#/incidentPV/RID-1"
    )

    linking_response = client.get(
        "/api/dashboard/linking",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "nps_group": "Todos",
        },
    )
    assert linking_response.status_code == 200
    linking_payload = linking_response.json()
    assert linking_payload["available"] is True
    assert linking_payload["kpis"]["incidents"] == 2
    assert linking_payload["causal_method"]["value"] == "palanca_touchpoint"
    assert linking_payload["navigation"][1]["label"] == "Touchpoints afectados por Palanca"
    assert "situation" in linking_payload
    assert "narrative" in linking_payload["situation"]
    assert "entity_summary" in linking_payload
    assert "scenarios" in linking_payload
    assert "deep_dive" in linking_payload
    assert linking_payload["navigation"][3]["label"] == "Análisis de Tópicos de NPS afectados"
    assert linking_payload["deep_dive"]["title"] == "Análisis de Tópicos de NPS afectados"
    assert linking_payload["deep_dive"]["topic_filter"]["default"] == "Todos"
    assert isinstance(linking_payload["deep_dive"]["topic_filter"]["options"], list)
    assert linking_payload["deep_dive"]["topic_filter"]["options"][0]["value"] == "Todos"
    assert linking_payload["deep_dive"]["ranking"]["rows"]
    assert linking_payload["deep_dive"]["evidence"]["rows"]
    assert linking_payload["deep_dive"]["trending"]["figure"] is not None
    assert linking_payload["scenarios"]["cards"][0]["anchor_topic"]
    assert linking_payload["entity_summary"]["table"][0]["Tópico NPS ancla"]
    assert [tab["label"] for tab in linking_payload["deep_dive"]["tabs"]] == [
        "Ranking de hipótesis",
        "Evidence wall",
    ]


def test_helix_links_resolve_incident_number_through_record_id() -> None:
    base_url = "https://itsmhelixbbva-smartit.onbmc.com/smartit/app/#/incidentPV/"
    incident_number = "INC000104366753"
    record_id = "IDGH5CDNHIEUEAT2Q5F3T2Q5F3CHLN"
    frame = pd.DataFrame(
        {
            "Incident Number": [incident_number],
            "Record ID": [record_id],
        }
    )

    lookup = build_helix_incident_url_lookup(frame, base_url=base_url)
    expected_url = f"{base_url}{record_id}"

    assert lookup[incident_number] == expected_url
    enriched = enrich_helix_incident_links(frame, base_url=base_url)
    assert enriched.loc[0, "Incident Number__href"] == expected_url


def test_helix_links_build_urls_from_record_id_and_do_not_fallback_to_incident_number() -> None:
    base_url = "https://itsmhelixbbva-smartit.onbmc.com/smartit/app/#/incidentPV/"
    explicit_url = f"{base_url}EXPLICIT_RECORD"
    frame = pd.DataFrame(
        {
            "Incident Number": ["INC-EXPLICIT", "INC-NO-RECORD"],
            "Record ID": ["RID-SHOULD-NOT-WIN", ""],
            "Incident URL": [explicit_url, base_url],
        }
    )

    lookup = build_helix_incident_url_lookup(frame, base_url=base_url)

    assert lookup["INC-EXPLICIT"] == f"{base_url}RID-SHOULD-NOT-WIN"
    assert "INC-NO-RECORD" not in lookup


def test_dashboard_report_endpoint_returns_a_valid_powerpoint(tmp_path: Path) -> None:
    client = TestClient(create_app(_settings(tmp_path)))
    _upload_nps_march(client)
    helix_fixture = _build_helix_fixture(tmp_path / "helix-report.xlsx")

    with helix_fixture.open("rb") as handle:
        upload_response = client.post(
            "/api/uploads/helix",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    helix_fixture.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )

    assert upload_response.status_code == 200

    report_response = client.get(
        "/api/dashboard/report/pptx",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "nps_group": "Todos",
            "min_n": 200,
            "min_similarity": 0.25,
            "max_days_apart": 10,
            "touchpoint_source": "domain_touchpoint",
        },
    )

    assert report_response.status_code == 200
    assert (
        report_response.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert "attachment;" in report_response.headers["content-disposition"]
    assert report_response.headers["x-nps-lens-saved-path"].endswith(".pptx")
    assert Path(report_response.headers["x-nps-lens-saved-path"]).exists()

    presentation = Presentation(BytesIO(report_response.content))
    assert len(presentation.slides) >= 8


def test_dashboard_report_endpoint_respects_selected_period_and_baseline_history(
    tmp_path: Path,
) -> None:
    client = TestClient(create_app(_settings(tmp_path)))
    _upload_nps_jan_feb(client)
    _upload_nps_march(client)
    helix_fixture = _build_helix_fixture(tmp_path / "helix-report-periods.xlsx")

    with helix_fixture.open("rb") as handle:
        upload_response = client.post(
            "/api/uploads/helix",
            data={
                "service_origin": "BBVA México",
                "service_origin_n1": "Senda",
                "service_origin_n2": "",
            },
            files={
                "file": (
                    helix_fixture.name,
                    handle,
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )

    assert upload_response.status_code == 200

    report_response = client.get(
        "/api/dashboard/report/pptx",
        params={
            "service_origin": "BBVA México",
            "service_origin_n1": "Senda",
            "service_origin_n2": "",
            "pop_year": "2026",
            "pop_month": "03",
            "nps_group": "Todos",
            "min_n": 50,
            "min_similarity": 0.25,
            "max_days_apart": 10,
            "touchpoint_source": "domain_touchpoint",
        },
    )

    assert report_response.status_code == 200

    presentation = Presentation(BytesIO(report_response.content))
    assert len(presentation.slides) >= 9

    slide_2_texts: list[str] = []
    for shape in presentation.slides[1].shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                slide_2_texts.append(paragraph.text or "")
    slide_2_text = " ".join(slide_2_texts)

    assert "2026-03" in slide_2_text
    assert "2026-01" not in slide_2_text
    assert "2026-02" not in slide_2_text

    all_texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    all_texts.append(paragraph.text or "")

    assert any("3. Qué ha cambiado en Palanca" in text for text in all_texts)
    assert not any("Qué ha cambiado en Subpalanca" in text for text in all_texts)
