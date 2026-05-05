from __future__ import annotations

import os
import tempfile
import zipfile
from datetime import date
from io import BytesIO
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
from pptx import Presentation

import nps_lens.reports.ppt_template as ppt_template_module
from nps_lens.analytics.incident_attribution import (
    TOUCHPOINT_SOURCE_BROKEN_JOURNEYS,
    TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
)
from nps_lens.design.tokens import DesignTokens, nps_score_color
from nps_lens.reports import executive_ppt
from nps_lens.reports.content_selectors import (
    parse_markdown_strong,
    select_causal_scenarios,
    select_negative_delta_rows,
    select_nonzero_kpis,
)
from nps_lens.reports.executive_ppt import generate_business_review_ppt
from nps_lens.reports.ppt_template import (
    build_presentation,
    find_corporate_template_path,
    resolve_layout,
)
from nps_lens.services.dashboard_service import DashboardService
from nps_lens.settings import Settings
from nps_lens.ui.charts import chart_daily_kpis, chart_incident_risk_recovery
from nps_lens.ui.theme import get_theme


def _sample_payload() -> dict:
    def _nps_group(score: int) -> str:
        if score <= 6:
            return "DETRACTOR"
        if score >= 9:
            return "PROMOTER"
        return "PASSIVE"

    overall_daily = pd.DataFrame(
        {
            "date": pd.date_range("2026-01-01", periods=40, freq="D"),
            "nps_mean": [7.3, 7.2, 7.1, 7.4, 7.5, 7.4, 7.2, 7.0] * 5,
            "focus_rate": [0.21, 0.23, 0.24, 0.22, 0.20, 0.19, 0.22, 0.25] * 5,
            "incidents": [4, 6, 8, 5, 4, 3, 7, 9] * 5,
            "responses": [120, 115, 118, 123, 119, 121, 117, 124] * 5,
        }
    )

    by_topic_daily = pd.DataFrame(
        {
            "date": pd.date_range("2026-01-01", periods=40, freq="D").tolist() * 3,
            "nps_topic": (["Pagos > SPEI"] * 40)
            + (["Acceso > Login"] * 40)
            + (["Tarjetas > Bloqueo"] * 40),
            "responses": ([44, 46, 45, 43, 44] * 8)
            + ([38, 37, 39, 40, 38] * 8)
            + ([26, 27, 29, 28, 27] * 8),
            "focus_count": ([12, 11, 13, 12, 10] * 8)
            + ([9, 8, 10, 9, 8] * 8)
            + ([5, 4, 5, 5, 4] * 8),
            "nps_mean": ([6.1, 6.3, 6.0, 6.2, 6.4] * 8)
            + ([5.8, 6.0, 5.7, 5.9, 6.1] * 8)
            + ([7.2, 7.1, 7.0, 7.3, 7.2] * 8),
            "focus_rate": ([0.24, 0.23, 0.26, 0.27, 0.25] * 8)
            + ([0.20, 0.19, 0.21, 0.22, 0.20] * 8)
            + ([0.18, 0.17, 0.19, 0.20, 0.18] * 8),
            "incidents": ([2, 3, 4, 2, 1] * 8) + ([1, 2, 2, 1, 0] * 8) + ([1, 1, 2, 1, 1] * 8),
        }
    )

    rationale = pd.DataFrame(
        [
            {
                "nps_topic": "Pagos > SPEI",
                "touchpoint": "Pagos",
                "priority": 0.91,
                "confidence": 0.80,
                "focus_probability_with_incident": 0.47,
                "nps_delta_expected": -4.8,
                "total_nps_impact": 1.9,
                "causal_score": 0.84,
                "nps_points_at_risk": 1.9,
                "nps_points_recoverable": 1.2,
                "best_lag_weeks": 1.0,
            },
            {
                "nps_topic": "Acceso > Login",
                "touchpoint": "Acceso",
                "priority": 0.79,
                "confidence": 0.71,
                "focus_probability_with_incident": 0.39,
                "nps_delta_expected": -3.6,
                "total_nps_impact": 1.1,
                "causal_score": 0.77,
                "nps_points_at_risk": 1.1,
                "nps_points_recoverable": 0.7,
                "best_lag_weeks": 1.0,
            },
            {
                "nps_topic": "Tarjetas > Bloqueo",
                "touchpoint": "Tarjetas",
                "priority": 0.68,
                "confidence": 0.64,
                "focus_probability_with_incident": 0.31,
                "nps_delta_expected": -2.8,
                "total_nps_impact": 0.9,
                "causal_score": 0.66,
                "nps_points_at_risk": 0.9,
                "nps_points_recoverable": 0.5,
                "best_lag_weeks": 2.0,
            },
        ]
    )

    lag_days = pd.DataFrame(
        {
            "nps_topic": ["Pagos > SPEI", "Acceso > Login", "Tarjetas > Bloqueo"],
            "best_lag_days": [4, 3, 5],
        }
    )

    incident_evidence = pd.DataFrame(
        {
            "incident_id": ["INC-9001", "INC-9123", "INC-9200"],
            "incident_date": pd.to_datetime(["2026-01-12", "2026-01-20", "2026-01-28"]),
            "nps_topic": ["Pagos > SPEI", "Acceso > Login", "Tarjetas > Bloqueo"],
            "incident_summary": [
                "Falla intermitente en pagos SPEI de banca móvil.",
                "Error de autenticación al iniciar sesión en app.",
                "Bloqueos recurrentes en activación de tarjeta digital.",
            ],
            "detractor_comment": [
                "No pude transferir y nadie resolvió en soporte.",
                "La app no me deja entrar desde ayer.",
                "Se bloquea la tarjeta y me quedo sin poder pagar.",
            ],
            "similarity": [0.92, 0.88, 0.84],
        }
    )

    changepoints = pd.DataFrame(
        {
            "nps_topic": ["Pagos > SPEI", "Acceso > Login", "Tarjetas > Bloqueo"],
            "changepoints": [
                ["2026-01-10", "2026-01-24"],
                ["2026-01-18"],
                ["2026-01-26"],
            ],
        }
    )

    attribution = pd.DataFrame(
        [
            {
                "nps_topic": "Acceso > Login",
                "touchpoint": "Login",
                "palanca": "Acceso",
                "subpalanca": "Login",
                "linked_incidents": 5,
                "linked_comments": 2,
                "linked_pairs": 5,
                "avg_similarity": 0.89,
                "avg_nps": 1.5,
                "detractor_probability": 0.47,
                "nps_delta_expected": -4.8,
                "total_nps_impact": 1.7,
                "nps_points_at_risk": 1.7,
                "nps_points_recoverable": 1.1,
                "priority": 0.91,
                "confidence": 0.82,
                "causal_score": 0.86,
                "delta_focus_rate_pp": 29.0,
                "incident_rate_per_100_responses": 8.5,
                "incidents": 5,
                "responses": 120,
                "action_lane": "Fix estructural",
                "owner_role": "Producto + Tecnologia",
                "eta_weeks": 6.0,
                "incident_records": [
                    {
                        "incident_id": "INC00001",
                        "summary": "problema en el login",
                        "url": "https://helix.example/INC00001",
                    },
                    {"incident_id": "INC00003", "summary": "no puedo acceder", "url": ""},
                    {
                        "incident_id": "INC00025",
                        "summary": "nada mas entras se desloguea",
                        "url": "",
                    },
                    {
                        "incident_id": "INC00040",
                        "summary": "error al autenticar usuario en acceso web",
                        "url": "",
                    },
                    {
                        "incident_id": "INC00041",
                        "summary": "falla de sesion al entrar en portal empresas",
                        "url": "",
                    },
                ],
                "incident_examples": [
                    "problema en el login",
                    "no puedo acceder",
                    "nada mas entras se desloguea",
                    "error al autenticar usuario en acceso web",
                    "falla de sesion al entrar en portal empresas",
                ],
                "comment_examples": [
                    "NPS 1: No hay quien entre a la aplicación",
                    "NPS 2: La web expulsa al usuario al entrar",
                ],
                "chain_story": "5 incidencias Helix degradan el touchpoint Login y se reflejan en 2 comentarios VoC con NPS muy bajo.",
            }
        ]
    )

    current_dates = pd.date_range("2026-01-01", periods=40, freq="D")
    baseline_dates = pd.date_range("2025-11-22", periods=40, freq="D")
    current_records: list[dict[str, object]] = []
    baseline_records: list[dict[str, object]] = []
    specs = [
        ("Acceso", "Login", "Web", [2, 3, 4, 5, 4], "La app no me deja entrar"),
        ("Pagos", "SPEI", "BBVA", [4, 5, 6, 5, 4], "Falla al transferir"),
        ("Tarjetas", "Bloqueo", "Otros", [7, 8, 7, 8, 9], "Se bloquea la tarjeta"),
    ]
    baseline_specs = [
        ("Acceso", "Login", "Web", [6, 7, 7, 8, 8], "Accedo sin problema"),
        ("Pagos", "SPEI", "BBVA", [7, 8, 8, 7, 8], "Transferencia completada"),
        ("Tarjetas", "Bloqueo", "Otros", [8, 8, 9, 9, 8], "Tarjeta operativa"),
    ]
    for idx, dt in enumerate(current_dates):
        for topic_idx, (palanca, subpalanca, canal, pattern, comment_base) in enumerate(specs):
            score = int(pattern[idx % len(pattern)])
            current_records.append(
                {
                    "ID": f"C-{idx}-{topic_idx}",
                    "Fecha": dt,
                    "NPS": score,
                    "NPS Group": _nps_group(score),
                    "Canal": canal,
                    "Palanca": palanca,
                    "Subpalanca": subpalanca,
                    "Comment": f"{comment_base} · {dt.date()}",
                }
            )
    for idx, dt in enumerate(baseline_dates):
        for topic_idx, (palanca, subpalanca, canal, pattern, comment_base) in enumerate(
            baseline_specs
        ):
            score = int(pattern[idx % len(pattern)])
            baseline_records.append(
                {
                    "ID": f"B-{idx}-{topic_idx}",
                    "Fecha": dt,
                    "NPS": score,
                    "NPS Group": _nps_group(score),
                    "Canal": canal,
                    "Palanca": palanca,
                    "Subpalanca": subpalanca,
                    "Comment": f"{comment_base} · {dt.date()}",
                }
            )

    selected_nps = pd.DataFrame(current_records)
    comparison_nps = pd.concat([pd.DataFrame(baseline_records), selected_nps], ignore_index=True)
    broken_journeys = pd.DataFrame(
        [
            {
                "journey_label": "Pagos / Transferencias / No funciona bien / falla",
                "touchpoint": "Transferencias",
                "palanca": "Pagos / Transferencias",
                "subpalanca": "No funciona bien / falla",
                "journey_keywords": "falla, transferencia, pago",
                "linked_pairs": 46,
                "linked_incidents": 26,
                "linked_comments": 32,
                "avg_similarity": 0.88,
                "avg_nps": 4.2,
                "semantic_cohesion": 0.87,
            },
            {
                "journey_label": "Uso / Practicidad / Facilidad de uso",
                "touchpoint": "Uso",
                "palanca": "Uso",
                "subpalanca": "Facilidad de uso",
                "journey_keywords": "facil, usar, practicidad",
                "linked_pairs": 5,
                "linked_incidents": 3,
                "linked_comments": 4,
                "avg_similarity": 0.81,
                "avg_nps": 5.0,
                "semantic_cohesion": 0.85,
            },
        ]
    )

    return {
        "overall_daily": overall_daily,
        "by_topic_daily": by_topic_daily,
        "rationale": rationale,
        "lag_days": lag_days,
        "incident_evidence": incident_evidence,
        "changepoints": changepoints,
        "attribution": attribution,
        "selected_nps": selected_nps,
        "comparison_nps": comparison_nps,
        "broken_journeys": broken_journeys,
    }


def _assert_no_shape_overflow(prs: Presentation) -> None:
    tolerance = 1000
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "left", None) is None:
                continue
            assert shape.left >= -tolerance, f"shape overflow left on slide {slide_index}"
            assert shape.top >= -tolerance, f"shape overflow top on slide {slide_index}"
            assert (
                shape.left + shape.width <= prs.slide_width + tolerance
            ), f"shape overflow right on slide {slide_index}"
            assert (
                shape.top + shape.height <= prs.slide_height + tolerance
            ), f"shape overflow bottom on slide {slide_index}"


def test_generate_business_review_ppt_builds_new_story() -> None:
    payload = _sample_payload()
    business_story = """# Informe de negocio — NPS Lens

## 1) Qué está pasando
- Muestras: 36,872 · Score medio (0-10): 8.53 · Detractores: 12.7% · Promotores: 72.5%
- Zona de fricción: Agregar funcionalidad · Zona fuerte: FAN

## 2) Cambio vs base de comparación
- Periodo actual: Mes actual (Febrero 2026 · 2026-02-01 → 2026-02-22) (n=20,791)
- Periodo base: Base histórica anterior a Febrero 2026 (2025-11-01 → 2026-01-31) (n=16,081)
- Variación: Delta Score -0.18 · Δ detractores +2.5 pp

## 3) Dónde atacar primero (oportunidades)
- Si mejoramos Palanca=Funcionamiento Continuo, el modelo estima un potencial de +57.2 puntos.

## 4) Qué están diciendo (temas de texto)
- Tema #1: fallas de continuidad, caídas y lentitud en procesos críticos.

## 5) Próximos pasos recomendados
- Validar releases, alinear owners y aterrizar quick wins del mes.
"""

    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="Empresas Mobile",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"],
        nps_points_at_risk=3.9,
        nps_points_recoverable=2.4,
        top3_incident_share=0.74,
        median_lag_weeks=1.2,
        story_md=business_story,
        script_8slides_md="",
        attribution_df=payload["attribution"],
        ranking_df=payload["rationale"],
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=payload["comparison_nps"],
        lag_days_by_topic=payload["lag_days"],
        by_topic_weekly=None,
        lag_weeks_by_topic=None,
        logo_path=None,
        incident_evidence_df=payload["incident_evidence"],
        changepoints_by_topic=payload["changepoints"],
        touchpoint_source="domain_touchpoint",
        entity_summary_df=payload["attribution"],
        entity_summary_kpis=[
            {"label": "Subpalancas activas", "value": "1"},
            {"label": "Confianza media", "value": "0.82"},
            {"label": "Links validados", "value": "5"},
        ],
    )

    assert out.content
    assert out.file_name.endswith(".pptx")
    assert out.slide_count == 10

    prs = Presentation(BytesIO(out.content))
    assert len(prs.slides) == out.slide_count
    _assert_no_shape_overflow(prs)

    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    cover_texts = []
    for shape in prs.slides[0].shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                cover_texts.append(paragraph.text or "")

    assert any("Análisis NPS térmico y causalidad" in t for t in texts)
    assert any("NPS térmico" in t for t in texts)
    assert any("% PROMOTORES" in t for t in texts)
    assert any("1. Evolución del NPS clásico del periodo" in t for t in texts)
    assert any("2. Que dicen los detractores" in t for t in texts)
    assert any("3. Qué ha cambiado en Palanca" in t for t in texts)
    assert not any("Qué ha cambiado en Subpalanca" in t for t in texts)
    assert any("4. Dónde duele en la Web · Palanca" in t for t in texts)
    assert not any("Dónde duele en la Web · Subpalanca" in t for t in texts)
    assert any("5. Oportunidades priorizadas · Palanca" in t for t in texts)
    assert not any("Oportunidades priorizadas · Subpalanca" in t for t in texts)
    assert any("Bloque 1 · Analisis VoC" in t for t in texts)
    assert any("Highlights del periodo" in t for t in texts)
    assert any("Análisis causal empleado · Por Subpalanca" in t for t in texts)
    assert any("6. Journeys de detracción" in t for t in texts)
    assert any("7.1 Acceso > Login" in t for t in texts)
    assert any("Análisis causal de Subpalanca: Escenario #1 ·" in t for t in texts)
    assert any("Sumario del análisis del escenario" in t for t in texts)
    assert any("Ejemplos de incidencias en el caso de uso" in t for t in texts)
    assert any("Ejemplos de Comentarios enlazados" in t for t in texts)
    assert any("Delta Score" in t for t in texts)
    assert not any("Lectura ejecutiva" in t for t in texts)
    assert not any("Criterio de recorte" in t for t in texts)
    assert not any("Mapa de dolor Web por Palanca" in t for t in texts)
    assert not any("Mapa de dolor Web por Subpalanca" in t for t in texts)
    assert not any("Impacto estimado en Palanca" in t for t in texts)
    assert not any("Impacto estimado en Subpalanca" in t for t in texts)
    assert not any("Journeys de detracción con mayor evidencia validada" in t for t in texts)
    assert not any("Detalle de evidencias Helix" in t for t in texts)
    assert not any("Qué destaca" in t for t in texts)
    assert not any("2. Cuándo y cómo lo dicen" in t for t in texts)
    assert not any("5. Casos más alejados del promedio" in t for t in texts)
    assert not any("Situación del periodo" in t for t in texts)
    assert not any("10.1 Matriz visual" in t for t in texts)
    assert not any("11.1 Señal temporal" in t for t in texts)
    assert not any("**" in t for t in texts)
    assert not any("Fix estructural" in t for t in texts)
    assert any("problema en el login" in t for t in texts)
    assert any("No hay quien entre a la aplicación" in t for t in texts)
    assert any("La web expulsa al usuario al entrar" in t for t in texts)
    assert not any("Muestras" in t for t in cover_texts)
    assert all(
        shape.text_frame.vertical_anchor != executive_ppt.MSO_VERTICAL_ANCHOR.MIDDLE
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and any((paragraph.text or "").strip() for paragraph in shape.text_frame.paragraphs)
    )
    with zipfile.ZipFile(BytesIO(out.content)) as archive:
        rels = archive.read("ppt/slides/_rels/slide10.xml.rels").decode("utf-8")
    assert "https://helix.example/INC00001" in rels


def test_generate_business_review_ppt_sanitizes_file_name_for_disk_write() -> None:
    payload = _sample_payload()
    out = generate_business_review_ppt(
        service_origin="MX/BU",
        service_origin_n1="Movil:Empresas?",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"],
        nps_points_at_risk=3.9,
        nps_points_recoverable=2.4,
        top3_incident_share=0.74,
        median_lag_weeks=1.2,
        story_md="",
        script_8slides_md="",
        attribution_df=payload["attribution"],
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=payload["comparison_nps"],
        logo_path=None,
    )

    assert "/" not in out.file_name
    assert ":" not in out.file_name
    assert "?" not in out.file_name


def test_generate_business_review_ppt_can_render_executive_journey_slide() -> None:
    payload = _sample_payload()
    attribution = payload["attribution"].copy()
    attribution.loc[:, "nps_topic"] = [
        "Acceso bloqueado",
    ]
    attribution.loc[:, "touchpoint"] = ["Login / autenticación"]
    attribution.loc[:, "palanca"] = ["Acceso"]
    attribution.loc[:, "subpalanca"] = ["Bloqueo / OTP"]
    attribution.loc[:, "journey_expected_evidence"] = [
        "Comentarios sobre login + incidencias de autenticación"
    ]
    attribution.loc[:, "journey_impact_label"] = ["Muy alto"]
    attribution.loc[:, "presentation_mode"] = [TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS]

    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="Empresas Mobile",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"],
        nps_points_at_risk=3.9,
        nps_points_recoverable=2.4,
        top3_incident_share=0.74,
        median_lag_weeks=1.2,
        story_md="",
        script_8slides_md="",
        attribution_df=attribution,
        ranking_df=payload["rationale"],
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=payload["comparison_nps"],
        lag_days_by_topic=payload["lag_days"],
        logo_path=None,
        incident_evidence_df=payload["incident_evidence"],
        changepoints_by_topic=payload["changepoints"],
        touchpoint_source=TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        entity_summary_df=attribution,
        entity_summary_kpis=[
            {"label": "Journeys de detracción", "value": "1"},
            {"label": "Touchpoints cubiertos", "value": "1"},
            {"label": "Links validados", "value": "5"},
        ],
        broken_journeys_df=payload["broken_journeys"],
    )

    prs = Presentation(BytesIO(out.content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    assert any("6. Journeys de detracción" in t for t in texts)
    assert any("Análisis causal de Journey de detracción: Escenario #1 ·" in t for t in texts)
    assert any("Acceso bloqueado" in t for t in texts)


def test_generate_business_review_ppt_merges_three_causal_scenarios_into_15_slides() -> None:
    payload = _sample_payload()
    base = payload["attribution"].iloc[0].to_dict()
    rows = [
        {
            **base,
            "nps_topic": "Acceso bloqueado",
            "touchpoint": "Login / autenticación",
            "palanca": "Acceso",
            "subpalanca": "Bloqueo / OTP",
            "linked_incidents": 5,
            "linked_comments": 3,
            "linked_pairs": 5,
            "detractor_probability": 0.13,
            "confidence": 0.20,
            "priority": 0.91,
            "presentation_mode": TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        },
        {
            **base,
            "nps_topic": "Operativa crítica fallida",
            "touchpoint": "Transferencias / pagos / firma",
            "palanca": "Operativa",
            "subpalanca": "Error funcional / timeout",
            "linked_incidents": 8,
            "linked_comments": 5,
            "linked_pairs": 10,
            "detractor_probability": 0.45,
            "confidence": 0.15,
            "priority": 0.62,
            "incident_records": [
                {
                    "incident_id": "INC000104256298",
                    "summary": "Condición de horario en proceso KNJCR2UC.",
                    "url": "https://helix.example/INC000104256298",
                },
                {
                    "incident_id": "INC000104257175",
                    "summary": "Caída en contratación de seguro cibernético Web.",
                    "url": "https://helix.example/INC000104257175",
                },
                {
                    "incident_id": "",
                    "summary": "Al ingresar al módulo de transferencias en tiempo real seguimiento, el sistema no muestra la opción de firmar las operaciones.",
                    "url": "",
                },
                {
                    "incident_id": "",
                    "summary": "Al querer obtener los comprobantes de pago del 16/01/2026, el sistema se queda cargando.",
                    "url": "",
                },
            ],
            "comment_examples": [
                "NPS 0: no me deja hacer trasnferencias , marca error de horario",
                "NPS 2: no me avisaron de que mi contrato de seguro subio de precio",
            ],
            "presentation_mode": TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        },
        {
            **base,
            "nps_topic": "Rendimiento degradado",
            "touchpoint": "Lentitud / cuelgues",
            "palanca": "Uso recurrente",
            "subpalanca": "Degradación del servicio",
            "linked_incidents": 2,
            "linked_comments": 2,
            "linked_pairs": 2,
            "detractor_probability": float("nan"),
            "confidence": 0.0,
            "priority": 0.10,
            "incident_records": [
                {
                    "incident_id": "",
                    "summary": "Se detecta que al abrir el PDF del EDC del mes de febrero se muestra un error.",
                    "url": "",
                },
                {
                    "incident_id": "",
                    "summary": "Se procede a obtener un estado de cuenta CFDI y el flujo no completa.",
                    "url": "",
                },
            ],
            "presentation_mode": TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        },
    ]
    attribution = pd.DataFrame(rows)

    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="ENTERPRISE WEB",
        service_origin_n2="",
        period_start=date(2026, 3, 1),
        period_end=date(2026, 3, 29),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"],
        nps_points_at_risk=0.0,
        nps_points_recoverable=0.0,
        top3_incident_share=0.0,
        median_lag_weeks=0.0,
        story_md="",
        script_8slides_md="",
        attribution_df=attribution,
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=payload["comparison_nps"],
        touchpoint_source=TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        entity_summary_df=attribution,
        entity_summary_kpis=[
            {"label": "Journeys de detracción", "value": "3"},
            {"label": "Touchpoints cubiertos", "value": "3"},
            {"label": "Links validados", "value": "17"},
        ],
    )

    prs = Presentation(BytesIO(out.content))
    texts = [
        paragraph.text or ""
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    ]

    assert out.slide_count == 12
    assert any("7.1 Operativa crítica fallida" in t for t in texts)
    assert any("7.2 Acceso bloqueado" in t for t in texts)
    assert any("7.3 Rendimiento degradado" in t for t in texts)
    assert any("Sumario del análisis del escenario" in t for t in texts)
    assert any("Ejemplos de incidencias en el caso de uso" in t for t in texts)
    assert not any("14.1" in t or "14.2" in t or "14.3" in t for t in texts)
    slide_11_texts = [
        paragraph.text or ""
        for shape in prs.slides[10].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    ]
    assert not any("NPS EN RIESGO" in t or "NPS RECUPERABLE" in t for t in slide_11_texts)
    assert any("INC000104257175" in t for t in texts)
    assert any("VÍNCULOS VALIDADOS" in t for t in texts)
    with zipfile.ZipFile(BytesIO(out.content)) as archive:
        rels = archive.read("ppt/slides/_rels/slide10.xml.rels").decode("utf-8")
    assert "https://helix.example/INC000104257175" in rels


def test_generate_business_review_ppt_can_render_broken_journey_story() -> None:
    payload = _sample_payload()
    attribution = payload["attribution"].copy()
    attribution.loc[:, "nps_topic"] = ["Acceso / Login"]
    attribution.loc[:, "touchpoint"] = ["Login"]
    attribution.loc[:, "palanca"] = ["Acceso"]
    attribution.loc[:, "subpalanca"] = ["Login"]
    attribution.loc[:, "journey_route"] = [
        "Incidencia -> Login -> Acceso / Login -> comentario VoC -> NPS"
    ]
    attribution.loc[:, "journey_expected_evidence"] = [
        "Keywords semánticas: Login, Otp. Helix Source Service N2 dominante: Auth."
    ]
    attribution.loc[:, "journey_cx_readout"] = ["5 links Helix↔VoC convergen en este journey roto."]
    attribution.loc[:, "presentation_mode"] = [TOUCHPOINT_SOURCE_BROKEN_JOURNEYS]

    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="Empresas Mobile",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"],
        nps_points_at_risk=3.9,
        nps_points_recoverable=2.4,
        top3_incident_share=0.74,
        median_lag_weeks=1.2,
        story_md="",
        script_8slides_md="",
        attribution_df=attribution,
        ranking_df=payload["rationale"],
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=payload["comparison_nps"],
        lag_days_by_topic=payload["lag_days"],
        logo_path=None,
        incident_evidence_df=payload["incident_evidence"],
        changepoints_by_topic=payload["changepoints"],
        touchpoint_source=TOUCHPOINT_SOURCE_BROKEN_JOURNEYS,
        entity_summary_df=attribution,
        entity_summary_kpis=[
            {"label": "Journeys rotos", "value": "1"},
            {"label": "Touchpoints detectados", "value": "1"},
            {"label": "Links validados", "value": "5"},
        ],
        broken_journeys_df=payload["broken_journeys"],
    )

    prs = Presentation(BytesIO(out.content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    assert any("6. Journeys rotos" in t for t in texts)
    assert any("Análisis causal de Journey roto: Escenario #1 ·" in t for t in texts)
    assert any("Acceso / Login" in t for t in texts)


def test_ppt_analytics_helpers_build_dynamic_tables() -> None:
    payload = _sample_payload()
    current = executive_ppt._coerce_nps_records(payload["selected_nps"])
    compare = executive_ppt._coerce_nps_records(payload["comparison_nps"])
    current_period, baseline_period = executive_ppt._split_period_frames(
        compare,
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
    )

    overview = executive_ppt._period_overview(current)
    assert int(overview["comments"]) > 0
    assert float(overview["detractor_rate"]) > 0

    palanca_change = executive_ppt._driver_change_table(
        current_period,
        baseline_period,
        dimension="Palanca",
    )
    assert not palanca_change.empty
    assert "delta_nps" in palanca_change.columns


def test_add_topic_timing_slide_reuses_app_charts_and_handles_empty_state() -> None:
    payload = _sample_payload()
    current = executive_ppt._coerce_nps_records(payload["selected_nps"])

    prs = Presentation()
    executive_ppt._add_topic_timing_slide(
        prs,
        period_label="2026-02-01 -> 2026-02-22",
        period_days=22,
        selected_nps_df=payload["selected_nps"],
    )
    executive_ppt._add_topic_timing_slide(
        prs,
        period_label="2026-02-01 -> 2026-02-22",
        period_days=22,
        selected_nps_df=pd.DataFrame(),
    )

    assert len(prs.slides) == 2

    slide_texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text or "")

    assert any("2. Cuándo y cómo lo dicen" in text for text in slide_texts)
    assert any("Cuándo lo dicen" in text for text in slide_texts)
    assert any("Cómo lo dicen" in text for text in slide_texts)
    assert any(
        "No hay señal suficiente para mostrar el volumen diario del periodo." in text
        for text in slide_texts
    )
    assert any(
        "No hay señal suficiente para la distribución diaria por grupo." in text
        for text in slide_texts
    )

    group_matrix = executive_ppt._group_matrix(current, dimension="Palanca")
    assert not group_matrix.empty
    assert set(group_matrix["band"].tolist()) <= {"Detractor", "Pasivo", "Promotor"}

    opportunities = executive_ppt._opportunities_table(current, dimension="Palanca", min_n=20)
    assert not opportunities.empty
    assert {"dimension", "potential_uplift", "confidence"}.issubset(opportunities.columns)

    gaps = executive_ppt._gap_vs_overall_table(current, top_k=5)
    assert len(gaps) <= 5
    assert not gaps.empty


def test_executive_ppt_legacy_chart_helpers_render_expected_figures() -> None:
    payload = _sample_payload()
    current = executive_ppt._coerce_nps_records(payload["selected_nps"])
    daily_mix = executive_ppt._daily_group_mix(current)
    topic_summary = executive_ppt._topic_summary(payload["by_topic_daily"])

    top_fig = executive_ppt._top_topics_fig(topic_summary, top_k=5)
    assert top_fig is not None
    assert len(top_fig.data) == 1

    heatmap_fig = executive_ppt._topic_heatmap_fig(payload["by_topic_daily"], top_k=3)
    assert heatmap_fig is not None
    assert len(heatmap_fig.data) == 1

    mix_fig = executive_ppt._daily_group_mix_fig(daily_mix)
    assert mix_fig is not None
    assert len(mix_fig.data) == 3

    themed = executive_ppt._apply_ppt_figure_theme(
        go.Figure(
            [
                go.Bar(name="Promotores", x=[1], y=[2]),
                go.Bar(name="Pasivos", x=[1], y=[3]),
                go.Bar(name="Detractores", x=[1], y=[4]),
                go.Bar(name="Incidencias", x=[1], y=[1]),
                go.Scatter(name="NPS clásico", x=[1], y=[2], mode="lines+markers"),
                go.Scatter(name="Incidencias", x=[1], y=[1], mode="lines"),
            ]
        )
    )
    assert themed.layout.legend.orientation == "h"
    assert themed.layout.font.size == 17
    assert themed.layout.legend.yanchor == "bottom"
    assert themed.data[0].marker.color == "#" + executive_ppt.BBVA_COLORS["green"]
    assert themed.data[1].marker.color == "#" + executive_ppt.BBVA_COLORS["yellow"]
    assert themed.data[2].marker.color == "#" + executive_ppt.BBVA_COLORS["red"]
    assert themed.data[3].marker.color == "#" + executive_ppt.BBVA_COLORS["sky"]

    heatmap_themed = executive_ppt._apply_ppt_figure_theme(
        go.Figure(
            [
                go.Heatmap(
                    z=[[0, 1]],
                    x=["2026-02-10", "2026-02-11"],
                    y=["Incidencias"],
                    colorbar=dict(title="Incidencias"),
                )
            ]
        )
    )
    assert heatmap_themed.layout.font.size == 18
    assert heatmap_themed.layout.margin.r >= 84
    assert heatmap_themed.data[0].xgap >= 2
    assert heatmap_themed.data[0].ygap >= 2
    assert heatmap_themed.data[0].colorbar.title.side == "right"


def test_add_opportunity_slide_reuses_app_chart_and_bullets() -> None:
    payload = _sample_payload()
    current = executive_ppt._coerce_nps_records(payload["selected_nps"])
    opportunities = executive_ppt._opportunities_table(current, dimension="Palanca", min_n=20)

    prs = Presentation()
    executive_ppt._add_opportunity_slide(
        prs,
        period_label="2026-02-01 -> 2026-02-22",
        opportunities_df=opportunities,
    )
    executive_ppt._add_opportunity_slide(
        prs,
        period_label="2026-02-01 -> 2026-02-22",
        opportunities_df=pd.DataFrame(columns=opportunities.columns),
    )

    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    assert any("6. Oportunidades a priorizar" in text for text in texts)
    assert any("Ranking por impacto estimado x confianza" in text for text in texts)
    assert any("Si mejoramos" in text for text in texts)
    assert not any("**" in text for text in texts)
    assert any(
        run.font.bold
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if "Palanca=" in run.text or "potencial de +" in run.text
    )
    assert any(
        "No se identificaron oportunidades robustas con el umbral actual." in text for text in texts
    )


def test_executive_ppt_helper_functions_cover_business_formatting_paths() -> None:
    assert executive_ppt._fmt_pct_or_nd(0.25) == "25%"
    assert executive_ppt._fmt_pct_or_nd(float("nan")) == "n/d"
    assert (
        executive_ppt._fmt_count_with_label(1, singular="incidencia", plural="incidencias")
        == "**1** incidencia"
    )
    assert executive_ppt._fmt_signed_or_nd(-2.34, decimals=1) == "-2,3"
    assert executive_ppt._fmt_num_or_nd(7.891, decimals=1) == "7,9"
    assert executive_ppt._clip("abcdefgh", 5) == "abcd…"
    assert executive_ppt._wrap_label("", width=6, max_lines=2) == ""
    assert executive_ppt._wrap_label("uno dos tres cuatro cinco seis", width=6, max_lines=2) != ""
    assert executive_ppt._focus_risk_label("detractores") == "detracción"
    assert executive_ppt._focus_probability_label("promotores") == "Prob. de promoción"
    assert executive_ppt._focus_risk_label("otros") == "otros"
    assert executive_ppt._format_opportunity_scope("Palanca", "Pagos") == "Pagos (palanca)"
    assert executive_ppt._format_opportunity_scope("Subpalanca", "Login") == "Login (subpalanca)"
    assert executive_ppt._format_opportunity_scope("nps_topic", "Tema X") == "Tema X"
    assert executive_ppt._clean_evidence_excerpt("", max_len=20) == ""
    assert executive_ppt._clean_evidence_excerpt("Descripción: texto de prueba", max_len=20) != ""
    assert executive_ppt._is_cover_metric_line("Score medio del periodo")

    source_df = pd.DataFrame(
        {
            "Fecha": pd.to_datetime(["2026-01-01", "2026-01-03", "2025-12-20"]),
            "NPS": [8, 4, 9],
        }
    )
    current, baseline = executive_ppt._split_source_period_frames(
        source_df,
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
    )
    assert len(current) == 2
    assert len(baseline) == 1
    assert str(executive_ppt._coerce_datetime_scalar("09-03-2026").date()) == "2026-03-09"
    assert executive_ppt._coerce_datetime_series(["09-03-2026", "10-03-2026"]).notna().all()
    assert executive_ppt._safe_date("09-03-2026") == "2026-03-09"
    assert executive_ppt._safe_dt("09-03-2026") is not None
    assert executive_ppt._month_label_es(date(2026, 3, 9)) == "marzo 2026"
    assert executive_ppt._slug("") == "na"
    assert executive_ppt._slug("Señal crítica / pagos") == "senal-critica-pagos"
    assert executive_ppt._safe_date(object()) != ""
    assert executive_ppt._safe_dt("no-date") is None
    assert executive_ppt._safe_float("7.5") == 7.5
    assert executive_ppt._safe_int("7.5") == 7
    assert executive_ppt._nps_band(10) == "Promotor"
    assert executive_ppt._nps_band(8) == "Pasivo"
    assert executive_ppt._nps_band(2) == "Detractor"
    assert (
        executive_ppt._normalize_category_value("Funcionamiento Continuo")
        == "Funcionamiento continuo"
    )
    assert (
        executive_ppt._normalize_category_value("Agregar Funcionalidad") == "Agregar funcionalidad"
    )
    assert executive_ppt._normalize_category_value("Fallas en el Login") == "Fallas en el login"


def test_editorial_content_selectors_are_deterministic_and_hide_zero_kpis() -> None:
    delta_df = pd.DataFrame(
        {
            "value": ["Mejora", "Peor A", "Peor B", "Neutro"],
            "delta_nps": [4.0, -8.0, -8.0, 0.0],
            "n_current": [100, 20, 80, 50],
        }
    )

    selected = select_negative_delta_rows(delta_df, max_rows=2)
    assert selected["value"].tolist() == ["Peor B", "Peor A"]
    assert select_negative_delta_rows(
        pd.DataFrame({"value": ["Mejora"], "delta_nps": [1.0], "n_current": [100]}),
        max_rows=2,
    ).empty

    kpis = select_nonzero_kpis(
        [
            ("Cero", "0 pts", "red"),
            ("Sin dato", "n/d", "blue"),
            ("Con valor", "1,5 pts", "green"),
        ],
        max_items=3,
    )
    assert kpis == [("Con valor", "1,5 pts", "green")]

    segments = parse_markdown_strong("Si mejoramos **Palanca=Acceso**, sube")
    assert [(segment.text, segment.bold) for segment in segments] == [
        ("Si mejoramos ", False),
        ("Palanca=Acceso", True),
        (", sube", False),
    ]

    scenarios = select_causal_scenarios(
        pd.DataFrame(
            [
                {
                    "nps_topic": "Acceso bloqueado",
                    "priority": 0.91,
                    "confidence": 0.20,
                    "detractor_probability": 0.13,
                    "linked_pairs": 5,
                    "linked_incidents": 5,
                    "linked_comments": 3,
                },
                {
                    "nps_topic": "Operativa crítica fallida",
                    "priority": 0.62,
                    "confidence": 0.15,
                    "detractor_probability": 0.45,
                    "linked_pairs": 10,
                    "linked_incidents": 8,
                    "linked_comments": 5,
                },
            ]
        ),
        max_rows=2,
    )
    assert scenarios["nps_topic"].tolist() == ["Operativa crítica fallida", "Acceso bloqueado"]


def test_daily_kpis_chart_places_x_axis_labels_at_bottom() -> None:
    payload = _sample_payload()

    fig = chart_daily_kpis(payload["selected_nps"], get_theme("light"), days=31)

    assert fig is not None
    assert fig.layout.xaxis.side == "bottom"
    assert fig.layout.xaxis.ticklabelposition == "outside bottom"

    stacked = executive_ppt.chart_daily_nps_committee_stack(
        payload["selected_nps"], get_theme("light"), days=31
    )
    assert stacked is not None
    assert len(stacked.data) == 5
    assert [trace.name for trace in stacked.data] == [
        "NPS clásico",
        "% detractores",
        "% promotores",
        "% pasivos",
        "% detractores",
    ]


def test_generate_business_review_ppt_handles_selected_period_without_history_or_chains() -> None:
    payload = _sample_payload()
    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="Empresas Mobile",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=payload["rationale"].head(0),
        nps_points_at_risk=0.0,
        nps_points_recoverable=0.0,
        top3_incident_share=0.0,
        median_lag_weeks=0.0,
        story_md="",
        script_8slides_md="",
        attribution_df=pd.DataFrame(),
        ranking_df=pd.DataFrame(),
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=payload["selected_nps"],
        comparison_nps_df=pd.DataFrame(),
        lag_days_by_topic=pd.DataFrame(),
        by_topic_weekly=None,
        lag_weeks_by_topic=None,
        logo_path=None,
        incident_evidence_df=pd.DataFrame(),
        changepoints_by_topic=pd.DataFrame(),
    )

    prs = Presentation(BytesIO(out.content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    assert any("1. Evolución del NPS clásico del periodo" in t for t in texts)
    assert any("6. Journeys de detracción" in t for t in texts)
    assert not any("7.1" in t for t in texts)


def test_ppt_template_fallback_builds_default_presentation() -> None:
    prs = build_presentation(template_path=None)
    layout = resolve_layout(prs, ["layout inexistente"], fallback_index=0)

    assert prs is not None
    assert layout is not None


def test_ppt_template_path_resolution_supports_explicit_and_env_paths() -> None:
    original = os.environ.get("NPS_LENS_PPT_TEMPLATE")
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        pptx_path = tmp_path / "corporate-template.pptx"
        Presentation().save(pptx_path)

        found_explicit = find_corporate_template_path(
            explicit_path=pptx_path, workspace_root=tmp_path
        )
        assert found_explicit == pptx_path

        os.environ["NPS_LENS_PPT_TEMPLATE"] = str(pptx_path)
        found_env = find_corporate_template_path(explicit_path=pptx_path, workspace_root=tmp_path)
        assert found_env == pptx_path

        fallback_prs = build_presentation(template_path=None, workspace_root=tmp_path / "missing")
        assert fallback_prs is not None

    if original is None:
        os.environ.pop("NPS_LENS_PPT_TEMPLATE", None)
    else:
        os.environ["NPS_LENS_PPT_TEMPLATE"] = original


def test_ppt_template_resolution_handles_duplicates_and_no_match() -> None:
    original_env = os.environ.get("NPS_LENS_PPT_TEMPLATE")
    original_names = ppt_template_module._TEMPLATE_FILE_NAMES
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_path = Path(tmpdir)
            missing = tmp_path / "missing-template.pptx"
            os.environ["NPS_LENS_PPT_TEMPLATE"] = str(missing)
            ppt_template_module._TEMPLATE_FILE_NAMES = ()

            assert (
                find_corporate_template_path(explicit_path=missing, workspace_root=tmp_path) is None
            )

            prs = build_presentation(template_path=None, workspace_root=tmp_path)
            assert prs is not None
    finally:
        ppt_template_module._TEMPLATE_FILE_NAMES = original_names
        if original_env is None:
            os.environ.pop("NPS_LENS_PPT_TEMPLATE", None)
        else:
            os.environ["NPS_LENS_PPT_TEMPLATE"] = original_env


def test_generate_business_review_ppt_falls_back_to_aggregate_signals_without_raw_nps() -> None:
    payload = _sample_payload()
    out = generate_business_review_ppt(
        service_origin="BBVA México",
        service_origin_n1="Empresas Mobile",
        service_origin_n2="",
        period_start=date(2026, 1, 1),
        period_end=date(2026, 1, 31),
        focus_name="detractores",
        overall_weekly=payload["overall_daily"],
        rationale_df=pd.DataFrame(),
        nps_points_at_risk=0.0,
        nps_points_recoverable=0.0,
        top3_incident_share=0.0,
        median_lag_weeks=0.0,
        story_md="",
        script_8slides_md="",
        attribution_df=pd.DataFrame(),
        ranking_df=pd.DataFrame(),
        by_topic_daily=payload["by_topic_daily"],
        selected_nps_df=None,
        comparison_nps_df=None,
        lag_days_by_topic=None,
        by_topic_weekly=None,
        lag_weeks_by_topic=None,
        logo_path=None,
        incident_evidence_df=None,
        changepoints_by_topic=None,
    )

    prs = Presentation(BytesIO(out.content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texts.append(paragraph.text or "")

    assert any("1. Evolución del NPS clásico del periodo" in t for t in texts)
    assert any("5. Oportunidades priorizadas · Palanca" in t for t in texts)


def test_history_fig_daily_uses_requested_colors() -> None:
    daily = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-02", "2026-01-03"]),
            "nps_mean": [7.2, 7.4, 7.1],
            "detractor_rate": [0.22, 0.20, 0.24],
            "incidents": [3, 4, 2],
        }
    )
    fig = executive_ppt._history_fig(daily, focus_name="detractores")
    assert fig is not None
    expected_markers = [
        nps_score_color(DesignTokens.default(), "light", value)
        for value in daily["nps_mean"].tolist()
    ]
    assert fig.data[0]["mode"] == "lines+markers"
    assert fig.data[0]["line"]["color"] == "#" + executive_ppt.BBVA_COLORS["blue"]
    assert list(fig.data[0]["marker"]["color"]) == expected_markers
    assert fig.data[2]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["yellow"]


def test_month_overlap_highlights_matched_incidents_with_labels() -> None:
    month = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-02", "2026-01-03"]),
            "nps_mean": [7.2, 7.3, 7.1],
            "detractor_rate": [0.21, 0.22, 0.24],
            "incidents": [4, 5, 2],
        }
    )
    matched = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-03"]),
            "matched_incidents": [2, 1],
        }
    )
    fig = executive_ppt._month_overlap_fig(
        month,
        focus_name="detractores",
        matched_daily=matched,
    )
    assert fig is not None
    assert fig.data[3]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["orange"]
    assert list(fig.data[3]["text"]) == ["2", "", "1"]


def test_executive_ppt_helper_figures_cover_secondary_paths() -> None:
    payload = _sample_payload()

    chain_fig = executive_ppt._chain_portfolio_fig(
        payload["attribution"],
        highlight_topic="Acceso > Login",
    )
    assert chain_fig is not None
    assert chain_fig.data[0]["marker"]["color"][0] == "#" + executive_ppt.BBVA_COLORS["red"]
    assert len(chain_fig.data[0]["x"]) == 1

    evo = executive_ppt._nps_evolution_fig(
        executive_ppt._daily_group_mix(executive_ppt._coerce_nps_records(payload["selected_nps"])),
        payload["overall_daily"],
    )
    assert evo is not None
    assert len(evo.data) == 3
    assert evo.data[0]["name"] == "NPS clásico"
    assert evo.data[1]["name"] == "% detractores"

    change_df = pd.DataFrame(
        {
            "value": ["A", "B", "C"],
            "n_current": [20, 10, 5],
            "delta_nps": [1.2, -2.5, 0.5],
        }
    )
    delta_fig = executive_ppt._delta_bars_fig(
        change_df,
        metric="delta_nps",
        x_title="Cambio NPS",
    )
    assert delta_fig is not None
    assert delta_fig.data[0]["orientation"] == "h"

    matrix_df = pd.DataFrame(
        {
            "Palanca": ["Pagos", "Pagos", "Pagos", "Acceso", "Acceso", "Acceso"],
            "band": ["Detractor", "Pasivo", "Promotor"] * 2,
            "share": [0.5, 0.3, 0.2, 0.2, 0.4, 0.4],
        }
    )
    heatmap = executive_ppt._group_heatmap_fig(matrix_df, dimension="Palanca")
    assert heatmap is not None
    assert len(heatmap.data) == 1

    gaps = pd.DataFrame(
        {
            "value": ["Pagos", "Acceso"],
            "gap_vs_overall": [-10.5, -3.2],
        }
    )
    gap_fig = executive_ppt._gap_vs_overall_fig(gaps)
    assert gap_fig is not None
    assert gap_fig.data[0]["orientation"] == "h"

    opps = pd.DataFrame(
        {
            "dimension": ["Palanca", "Subpalanca", "nps_topic"],
            "value": ["Pagos", "Login", "Transferencias lentas"],
            "confidence": [0.7, 0.4, 0.5],
            "potential_uplift": [4.2, 2.1, 1.3],
            "n": [100, 64, 25],
        }
    )
    opp_fig = executive_ppt._opportunity_bubble_fig(opps)
    assert opp_fig is not None
    assert len(opp_fig.data) == 1


def test_text_topic_slide_uses_all_clusters_for_chart_and_top_three_for_table() -> None:
    topics = pd.DataFrame(
        {
            "cluster_id": [1, 2, 3, 4, 5],
            "n": [500, 400, 300, 200, 100],
            "top_terms": [["uno", "dos"]] * 5,
            "examples": [["ejemplo"]] * 5,
            "label": [""] * 5,
            "top_terms_txt": ["uno, dos"] * 5,
            "example_txt": ["ejemplo"] * 5,
        }
    )

    fig = executive_ppt._build_text_topic_figure(topics)
    assert fig is not None
    assert len(fig.data[0].x) == 5

    prs = Presentation()
    executive_ppt._add_deep_dive_slide(
        prs,
        period_label="2026-03-01 -> 2026-03-29",
        text_topics_df=topics,
        topic_figure=None,
    )
    texts = [
        paragraph.text or ""
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
    ]
    assert "uno, dos" in texts
    assert "ejemplo" in texts
    assert not any(text in {"1", "2", "3", "4", "5"} for text in texts)


def test_chain_helpers_cover_edge_cases() -> None:
    assert executive_ppt._chain_list([" A ", "", "B"]) == ["A", "B"]
    assert executive_ppt._chain_list(None) == []
    assert executive_ppt._chain_list("uno") == ["uno"]
    assert executive_ppt._chain_header("Helix", shown=2, total=5) == "Helix (2 de 5)"
    assert executive_ppt._chain_header("Helix", shown=2, total=2) == "Helix (2)"
    assert executive_ppt._chain_incident_records([{"incident_id": "INC1", "summary": "hola"}]) == [
        {"incident_id": "INC1", "summary": "hola", "url": ""}
    ]
    assert executive_ppt._chain_incident_records(
        [{"incident_id": "INC2", "summary": "hola", "incident_id__href": "https://helix/2"}]
    ) == [{"incident_id": "INC2", "summary": "hola", "url": "https://helix/2"}]
    assert executive_ppt._chain_incident_records(["bad"]) == []


def test_dashboard_service_injects_helix_urls_into_incident_records(tmp_path: Path) -> None:
    settings = Settings(
        data_dir=tmp_path / "data",
        database_path=tmp_path / "data" / "dashboard.sqlite3",
        frontend_dist_dir=tmp_path / "frontend-dist",
        frontend_public_dir=tmp_path / "frontend-public",
        api_host="127.0.0.1",
        api_port=8000,
        default_service_origin="BBVA México",
        default_service_origin_n1="Senda",
        allowed_service_origins=["BBVA México"],
        allowed_service_origin_n1={"BBVA México": ["Senda"]},
        log_level="INFO",
    )
    service = DashboardService(repository=object(), settings=settings)  # type: ignore[arg-type]
    chain_df = pd.DataFrame(
        {
            "incident_records": [
                [{"incident_id": "INC00042", "summary": "Falla de acceso", "url": ""}]
            ]
        }
    )
    helix_df = pd.DataFrame(
        {
            "Incident Number": ["INC00042"],
            "Record ID": ["AGGADG1A2B3C"],
        }
    )

    out = service._inject_incident_record_urls(chain_df, helix_df=helix_df)
    record = out.iloc[0]["incident_records"][0]
    assert record["url"].endswith("/AGGADG1A2B3C")
    assert record["incident_id__href"] == record["url"]


def test_incident_risk_recovery_wraps_labels_for_small_ppt_panels() -> None:
    rationale = pd.DataFrame(
        {
            "nps_topic": ["Pagos / Transferencias / No funciona bien / Error intermitente"],
            "nps_points_at_risk": [0.74],
            "nps_points_recoverable": [0.15],
            "priority": [0.82],
        }
    )

    fig = chart_incident_risk_recovery(rationale, get_theme("light"), top_k=1)
    assert fig is not None
    assert "<br>" in str(fig.data[0]["y"][0]) or "…" in str(fig.data[0]["y"][0])
    assert fig.data[0]["cliponaxis"] is False
    assert fig.data[1]["cliponaxis"] is False


def test_build_incident_timeline_daily_filters_to_matching_hot_terms() -> None:
    timeline = pd.DataFrame(
        {
            "date": [
                "2026-02-10",
                "2026-02-10",
                "2026-02-12",
                "2026-02-13",
                "2026-03-01",
            ],
            "helix_records": [2, 1, 3, 0, 5],
            "nps_comments": [1, 2, 1, 4, 1],
            "hot_term": ["pagos", "login", "pagos", "pagos", "pagos"],
        }
    )
    evidence = pd.DataFrame(
        {
            "hot_term": ["pagos", "login", "otros"],
            "hot_rank": [1, 2, 4],
        }
    )

    out = executive_ppt._hotspot_matches_by_day(
        timeline,
        evidence,
        month_start=pd.Timestamp("2026-02-01"),
        month_end=pd.Timestamp("2026-02-28"),
    )

    assert list(out["matched_incidents"]) == [3, 3]
    assert out["date"].dt.strftime("%Y-%m-%d").tolist() == ["2026-02-10", "2026-02-12"]

    missing = executive_ppt._hotspot_matches_by_day(
        pd.DataFrame({"date": ["2026-02-10"]}),
        evidence,
        month_start=pd.Timestamp("2026-02-01"),
        month_end=pd.Timestamp("2026-02-28"),
    )
    assert missing.empty


def test_topic_metrics_and_placeholder_text_helpers() -> None:
    payload = _sample_payload()
    metrics = executive_ppt._topic_metrics("Pagos > SPEI", payload["rationale"])
    assert metrics["risk"] == 1.9
    assert metrics["recoverable"] == 1.2
    assert executive_ppt._topic_metrics("Inexistente", payload["rationale"]) == {}
    assert executive_ppt._topic_metrics("Pagos > SPEI", pd.DataFrame()) == {}

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    executive_ppt._set_placeholder_text(
        slide,
        0,
        "Titulo de prueba",
        font_name=executive_ppt.BBVA_FONT_HEAD,
        size_pt=24,
    )
    assert slide.placeholders[0].text == "Titulo de prueba"
    executive_ppt._set_placeholder_text(
        slide,
        99,
        "Ignorado",
        font_name=executive_ppt.BBVA_FONT_HEAD,
        size_pt=24,
    )


def test_add_story_card_caps_bullets_by_height() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    executive_ppt._add_story_card(
        slide,
        left=1.0,
        top=1.0,
        width=4.0,
        height=1.5,
        title="Resumen",
        bullets=["uno", "dos", "tres", "cuatro", "cinco"],
    )

    shape = slide.shapes[-1]
    texts = [p.text for p in shape.text_frame.paragraphs]
    assert texts[0] == "Resumen"
    assert len([t for t in texts[1:] if t]) == 3


def test_executive_ppt_helper_guards_return_empty_or_none_cleanly() -> None:
    assert executive_ppt._chain_portfolio_fig(pd.DataFrame(), highlight_topic="x") is None
    assert executive_ppt._nps_evolution_fig(pd.DataFrame(), pd.DataFrame()) is None
    assert executive_ppt._delta_bars_fig(pd.DataFrame(), metric="delta_nps", x_title="x") is None
    assert executive_ppt._group_heatmap_fig(pd.DataFrame(), dimension="Palanca") is None
    assert executive_ppt._gap_vs_overall_fig(pd.DataFrame()) is None
    assert executive_ppt._opportunity_bubble_fig(pd.DataFrame()) is None
    assert executive_ppt._hotspot_matches_by_day(
        None,
        None,
        month_start=pd.Timestamp("2026-02-01"),
        month_end=pd.Timestamp("2026-02-28"),
    ).empty


def test_top_hotspots_fig_uses_top3_colors_and_inbar_labels() -> None:
    evidence = pd.DataFrame(
        {
            "hot_rank": [1, 1, 2, 2, 3, 3],
            "hot_term": [
                "pagos",
                "pagos",
                "movimientos",
                "movimientos",
                "transferencias",
                "transferencias",
            ],
            "mention_incidents": [60, 60, 45, 45, 30, 30],
            "mention_comments": [114, 114, 98, 98, 85, 85],
            "hotspot_comments": [114, 114, 98, 98, 85, 85],
            "hotspot_links": [90, 90, 72, 72, 54, 54],
        }
    )
    timeline = pd.DataFrame(
        {
            "incident_id": ["", "", ""],
            "hot_term": ["pagos", "movimientos", "transferencias"],
            "date": pd.to_datetime(["2026-01-01", "2026-01-01", "2026-01-01"]),
            "helix_records": [60, 45, 30],
            "nps_comments": [114, 98, 85],
        }
    )

    fig = executive_ppt._top_hotspots_fig(evidence, timeline, top_k=3)
    assert fig is not None
    assert len(fig.data) == 1

    colors = list(fig.data[0]["marker"]["color"])
    assert colors == [
        "#" + executive_ppt.BBVA_COLORS["yellow"],
        "#" + executive_ppt.BBVA_COLORS["orange"],
        "#" + executive_ppt.BBVA_COLORS["red"],
    ]
    assert fig.layout.xaxis.visible is False
    assert all(str(t).strip().startswith("#") for t in list(fig.data[0]["text"]))


def test_hotspot_matches_by_day_uses_hot_terms_and_overlap_signal() -> None:
    timeline = pd.DataFrame(
        {
            "incident_id": ["", "", "", ""],
            "hot_term": ["transferencias", "transferencias", "login", "login"],
            "date": pd.to_datetime(["2026-01-05", "2026-01-06", "2026-01-05", "2026-01-07"]),
            "helix_records": [3, 2, 1, 2],
            "nps_comments": [2, 0, 1, 3],
        }
    )
    evidence = pd.DataFrame(
        {
            "hot_term": ["transferencias", "login"],
            "hot_rank": [1, 2],
        }
    )
    out = executive_ppt._hotspot_matches_by_day(
        timeline,
        evidence,
        month_start=pd.Timestamp("2026-01-01"),
        month_end=pd.Timestamp("2026-01-31"),
    )
    assert list(out["date"].dt.strftime("%Y-%m-%d")) == ["2026-01-05", "2026-01-07"]
    assert list(out["matched_incidents"].astype(int)) == [4, 2]


def test_hotspot_stack_fig_uses_requested_color_semantics_and_horizontal_legend() -> None:
    daily = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-02"]),
            "nps_mean": [7.1, 7.2],
            "detractor_rate": [0.2, 0.21],
            "incidents": [5, 6],
        }
    )
    evidence = pd.DataFrame(
        {
            "incident_id": ["INC-1", "INC-2", "INC-3"],
            "incident_date": pd.to_datetime(["2026-01-01", "2026-01-01", "2026-01-02"]),
            "hot_rank": [1, 2, 3],
            "hot_term": ["transferencias", "token", "autenticacion"],
            "similarity": [0.9, 0.8, 0.7],
        }
    )
    fig = executive_ppt._hotspot_stack_fig(daily, evidence)
    assert fig is not None
    assert fig.data[0]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["blue"]
    assert fig.data[1]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["yellow"]
    assert fig.data[2]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["orange"]
    assert fig.data[3]["marker"]["color"] == "#" + executive_ppt.BBVA_COLORS["red"]
    assert fig.layout.legend.orientation == "h"


def test_incident_related_timeline_keeps_only_days_with_related_evidence() -> None:
    timeline = pd.DataFrame(
        {
            "incident_id": ["INC-9001", "INC-9001", "INC-9001", "INC-9002"],
            "date": pd.to_datetime(["2026-01-10", "2026-01-11", "2026-01-12", "2026-01-10"]),
            "helix_records": [1, 0, 2, 1],
            "nps_comments": [0, 3, 0, 1],
        }
    )
    out = executive_ppt._incident_related_timeline(
        incident_id="INC-9001",
        incident_timeline_df=timeline,
    )
    assert not out.empty
    assert list(out["date"].dt.strftime("%Y-%m-%d")) == ["2026-01-10", "2026-01-11", "2026-01-12"]
    assert out["helix_records"].sum() == 3
    assert out["nps_comments"].sum() == 3


def test_prepare_incident_evidence_prefers_detailed_description_rowwise() -> None:
    raw = pd.DataFrame(
        {
            "incident_id": ["INC-1", "INC-2"],
            "summary": ["Resumen muy corto", "Resumen fallback"],
            "Detailed Description": ["Detalle amplio 1", ""],
            "detractor_comment": ["Comentario 1", "Comentario 2"],
            "similarity": [0.81, 0.62],
        }
    )
    out = executive_ppt._prepare_incident_evidence(raw)
    assert out.iloc[0]["incident_summary"] == "Detalle amplio 1"
    assert out.iloc[1]["incident_summary"] == "Resumen fallback"


def test_select_zoom_incidents_groups_hotspots_instead_of_single_1to1() -> None:
    evidence = pd.DataFrame(
        {
            "incident_id": ["INC-A1", "INC-A2", "INC-B1"],
            "incident_date": pd.to_datetime(["2026-01-10", "2026-01-11", "2026-01-13"]),
            "nps_topic": ["Pagos > Transferencias", "Pagos > Transferencias", "Acceso > Login"],
            "incident_summary": ["Fallo 1", "Fallo 2", "Fallo 3"],
            "detractor_comment": ["No transfiere", "Transferencia falla", "No puedo entrar"],
            "similarity": [0.91, 0.87, 0.80],
            "hot_term": ["transferencias", "transferencias", "login"],
            "hot_rank": [1, 1, 2],
            "hotspot_incidents": [8, 8, 3],
            "hotspot_comments": [21, 21, 7],
            "hotspot_links": [13, 13, 5],
        }
    )
    selected = executive_ppt._select_zoom_incidents([], evidence, max_items=2)
    assert len(selected) == 2
    assert selected[0].hot_term == "transferencias"
    assert selected[0].hotspot_incidents == 8
    assert selected[0].hotspot_comments == 21
    assert selected[0].hotspot_links == 13
    assert "INC-A1" in selected[0].sample_incidents
    assert "INC-A2" in selected[0].sample_incidents


def test_incident_related_timeline_can_aggregate_by_hotspot_term() -> None:
    timeline = pd.DataFrame(
        {
            "incident_id": ["INC-A1", "INC-A2", "INC-B1"],
            "hot_term": ["transferencias", "transferencias", "login"],
            "date": pd.to_datetime(["2026-01-10", "2026-01-10", "2026-01-11"]),
            "helix_records": [2, 1, 1],
            "nps_comments": [1, 3, 0],
            "nps_comments_moderate": [1, 1, 0],
            "nps_comments_high": [0, 2, 0],
            "nps_comments_critical": [0, 0, 0],
        }
    )
    out = executive_ppt._incident_related_timeline(
        incident_id="INC-A1",
        hot_term="transferencias",
        incident_timeline_df=timeline,
    )
    assert len(out) == 1
    assert out.iloc[0]["helix_records"] == 3
    assert out.iloc[0]["nps_comments"] == 4
    assert out.iloc[0]["nps_comments_moderate"] == 2
    assert out.iloc[0]["nps_comments_high"] == 2
    assert "INC-A1" in out.iloc[0]["incident_ids"]
    assert "INC-A2" in out.iloc[0]["incident_ids"]


def test_zoom_hotspot_fig_uses_daily_red_comments_blue_points_and_nps_line() -> None:
    rel = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-20"]),
            "helix_records": [1, 1],
            "nps_comments": [2, 0],
            "nps_comments_moderate": [1, 0],
            "nps_comments_high": [1, 0],
            "nps_comments_critical": [0, 0],
            "incident_ids": ["INC-1", "INC-2"],
        }
    )
    incident = executive_ppt.ZoomIncident(
        incident_id="INC-1",
        incident_date=pd.Timestamp("2026-01-20"),
        nps_topic="Pagos > Transferencias",
        incident_summary="Resumen",
        detractor_comment="Comentario",
        similarity=0.9,
        hot_term="transferencias",
    )
    topic_daily = pd.DataFrame(
        {
            "date": pd.to_datetime(["2026-01-01", "2026-01-20"]),
            "nps_mean": [8.2, 6.1],
        }
    )
    fig = executive_ppt._zoom_incident_fig(
        topic_daily=topic_daily,
        related_timeline=rel,
        incident=incident,
        lag_days=4,
        changepoints=[],
        focus_name="detractores",
    )
    assert fig is not None
    assert fig.data[0]["type"] == "bar"
    assert fig.data[0]["name"] == "Comentarios moderados (NPS 5-6)"
    assert fig.data[1]["name"] == "Comentarios altos (NPS 3-4)"
    assert fig.data[2]["name"] == "Comentarios críticos (NPS 0-2)"
    assert fig.data[4]["type"] == "scatter"
    assert fig.data[4]["mode"] == "markers+text"
    assert fig.data[4]["yaxis"] == "y2"
    assert any(str(t) == "INC-1" for t in fig.data[4]["text"])
    assert fig.data[5]["name"] == "Score medio"
    assert fig.data[5]["type"] == "scatter"
    assert fig.data[5]["mode"] == "lines+markers"
    assert fig.data[5]["yaxis"] == "y3"
    assert fig.data[5]["line"]["color"] == "#" + executive_ppt.BBVA_COLORS["blue"]


def test_change_layout_uses_full_width_chart_and_table() -> None:
    layout = executive_ppt.CHANGE_SLIDE_LAYOUT

    assert layout.chart_panel.left < 0.70
    assert layout.chart_panel.width > 12.0
    assert layout.table_panel.top > layout.chart_panel.top + layout.chart_panel.height
    assert layout.max_rows == 4

    df = pd.DataFrame(
        {
            "value": ["B", "A", "C"],
            "delta_nps": [-0.2, -1.1, 0.5],
            "nps_current": [7.0, 6.0, 8.0],
            "nps_baseline": [7.2, 7.1, 7.5],
            "n_current": [100, 80, 30],
            "n_baseline": [120, 95, 40],
        }
    )

    out = select_negative_delta_rows(df, max_rows=2)

    assert out["value"].tolist() == ["A", "B"]


def test_journey_table_exposes_catalog_detail_columns() -> None:
    entity_summary = pd.DataFrame(
        [
            {
                "entity_label": "Operativa crítica fallida",
                "source_nps_topic": "Pagos / Transferencias > Mostrar movimientos actualizados",
                "touchpoint": "Transferencias / pagos / firma",
                "palanca": "Operativa",
                "subpalanca": "Error funcional / timeout",
                "linked_pairs": 16,
                "linked_comments": 13,
                "avg_nps": 2.0,
                "confidence": 0.38,
                "priority": 0.5,
            }
        ]
    )

    table = executive_ppt._build_journey_table(
        touchpoint_source=TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
        entity_summary_df=entity_summary,
        broken_journeys_df=None,
    )

    assert table.loc[0, "journey"] == "Operativa crítica fallida"
    assert table.loc[0, "palanca"] == "Operativa"
    assert table.loc[0, "anchor_topic"].startswith("Pagos / Transferencias")
    assert {"touchpoint", "subpalanca", "links", "confidence"}.issubset(table.columns)
