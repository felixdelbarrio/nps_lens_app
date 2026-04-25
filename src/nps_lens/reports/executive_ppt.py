from __future__ import annotations

import contextlib
import math
import os
import re
import tempfile
import textwrap
import unicodedata
from dataclasses import dataclass
from datetime import date, datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Iterable, Optional

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from nps_lens.analytics.drivers import compute_nps_from_scores, driver_table
from nps_lens.analytics.hotspot_metrics import (
    HOTSPOT_EVIDENCE_COLUMNS,
    summarize_hotspot_counts,
)
from nps_lens.analytics.hotspot_metrics import (
    build_hotspot_daily_breakdown as build_hotspot_daily_breakdown_metrics,
)
from nps_lens.analytics.incident_attribution import (
    EXECUTIVE_JOURNEY_CATALOG,
    TOUCHPOINT_SOURCE_BROKEN_JOURNEYS,
    TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS,
)
from nps_lens.analytics.nps_helix_link import build_nps_topic
from nps_lens.analytics.opportunities import rank_opportunities
from nps_lens.analytics.text_mining import extract_topics
from nps_lens.design.tokens import (
    DesignTokens,
    bbva_typography_tokens,
    executive_report_palette,
    nps_score_color,
    palette,
    plotly_continuous_scale,
    plotly_risk_scale,
)
from nps_lens.domain.causal_methods import get_causal_method_spec
from nps_lens.reports.content_selectors import (
    parse_markdown_strong,
    select_causal_scenarios,
    select_negative_delta_rows,
    select_nonzero_kpis,
    select_opportunities,
    select_text_clusters,
)
from nps_lens.reports.editorial_tokens import EDITORIAL_LIMITS
from nps_lens.reports.ppt_template import (
    CorporatePresentationTheme,
    build_presentation,
    resolve_layout,
)
from nps_lens.reports.presentation_context import (
    CausalScenarioViewModel,
    CausalViewModel,
    DimensionViewModel,
    PresentationContext,
)
from nps_lens.ui.business import driver_delta_table
from nps_lens.ui.charts import (
    _compact_axis_label,
    chart_causal_entity_bar,
    chart_cohort_heatmap,
    chart_daily_kpis,
    chart_daily_mix_business,
    chart_daily_volume,
    chart_driver_bar,
    chart_driver_delta,
    chart_opportunities_bar,
    chart_topic_bars,
)
from nps_lens.ui.narratives import explain_opportunities
from nps_lens.ui.theme import get_theme

BBVA_COLORS = executive_report_palette(DesignTokens.default(), mode="light")
BBVA_TYPOGRAPHY = bbva_typography_tokens()
PPT_THEME = CorporatePresentationTheme(
    display_font=BBVA_TYPOGRAPHY.display,
    heading_font=BBVA_TYPOGRAPHY.heading,
    body_font=BBVA_TYPOGRAPHY.body,
    medium_font=BBVA_TYPOGRAPHY.medium,
)

BBVA_FONT_DISPLAY = PPT_THEME.display_font
BBVA_FONT_HEAD = PPT_THEME.heading_font
BBVA_FONT_BODY = PPT_THEME.body_font
BBVA_FONT_MEDIUM = PPT_THEME.medium_font


def _ppt_nps_marker_colors(values: pd.Series | list[object]) -> list[str]:
    series = values if isinstance(values, pd.Series) else pd.Series(list(values))
    tokens = DesignTokens.default()
    return [nps_score_color(tokens, "light", value) for value in series.tolist()]


@dataclass(frozen=True)
class BusinessPptResult:
    file_name: str
    content: bytes
    slide_count: int
    saved_path: str = ""


@dataclass(frozen=True)
class ZoomIncident:
    incident_id: str
    incident_date: Optional[pd.Timestamp]
    nps_topic: str
    incident_summary: str
    detractor_comment: str
    similarity: float
    hot_term: str
    mention_incidents: int = 0
    mention_comments: int = 0
    hotspot_incidents: int = 0
    hotspot_comments: int = 0
    hotspot_links: int = 0
    sample_incidents: str = ""
    sample_comments: str = ""


def _rgb(hex_code: str) -> RGBColor:
    code = str(hex_code or "").strip().lstrip("#")
    if len(code) != 6:
        code = BBVA_COLORS["ink"]
    return RGBColor(int(code[0:2], 16), int(code[2:4], 16), int(code[4:6], 16))


def _safe_float(v: object, default: float = 0.0) -> float:
    try:
        f = float(v)
    except Exception:
        return float(default)
    if not np.isfinite(f):
        return float(default)
    return float(f)


def _safe_int(v: object, default: int = 0) -> int:
    try:
        i = int(float(v))
    except Exception:
        return int(default)
    return int(i)


def _fmt_locale_number(
    value: object,
    *,
    decimals: int = 0,
    signed: bool = False,
    default: str = "n/d",
) -> str:
    f = _safe_float(value, default=float("nan"))
    if not np.isfinite(f):
        return default
    precision = max(int(decimals), 0)
    rendered = f"{f:+,.{precision}f}" if signed else f"{f:,.{precision}f}"
    return rendered.replace(",", "_").replace(".", ",").replace("_", ".")


def _fmt_count_or_nd(v: object) -> str:
    return _fmt_locale_number(v, decimals=0)


def _fmt_pct_or_nd(v: object, decimals: int = 0) -> str:
    f = _safe_float(v, default=float("nan"))
    return "n/d" if not np.isfinite(f) else f"{_fmt_locale_number(f * 100.0, decimals=decimals)}%"


def _fmt_signed_or_nd(v: object, decimals: int = 1) -> str:
    return _fmt_locale_number(v, decimals=decimals, signed=True)


def _fmt_num_or_nd(v: object, decimals: int = 2) -> str:
    return _fmt_locale_number(v, decimals=decimals)


def _clip(txt: object, max_len: int) -> str:
    s = " ".join(str(txt or "").split())
    if len(s) <= max_len:
        return s
    return s[: max_len - 1].rstrip() + "…"


def _coerce_datetime_scalar(value: object) -> pd.Timestamp:
    try:
        return pd.Timestamp(pd.to_datetime(value, format="mixed", dayfirst=True, errors="coerce"))
    except TypeError:
        return pd.Timestamp(pd.to_datetime(value, dayfirst=True, errors="coerce"))


def _coerce_datetime_series(values: object) -> pd.Series:
    series = values if isinstance(values, pd.Series) else pd.Series(values)
    try:
        return pd.to_datetime(series, format="mixed", dayfirst=True, errors="coerce")
    except TypeError:
        return pd.to_datetime(series, dayfirst=True, errors="coerce")


def _wrap_label(
    txt: object,
    *,
    width: int = 24,
    max_lines: int = 2,
    joiner: str = "<br>",
) -> str:
    clean = " ".join(str(txt or "").split())
    if not clean:
        return ""
    lines = textwrap.wrap(clean, width=max(int(width), 8)) or [clean]
    if len(lines) > max(int(max_lines), 1):
        lines = lines[:max_lines]
        lines[-1] = _clip(lines[-1], max(int(width) - 1, 8))
    return joiner.join(lines)


def _configure_text_frame(tf: object) -> None:
    with contextlib.suppress(Exception):
        tf.word_wrap = True
    with contextlib.suppress(Exception):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    with contextlib.suppress(Exception):
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


def _focus_risk_label(focus_name: str) -> str:
    focus = str(focus_name or "").strip().lower()
    if focus in {"detractores", "detractor", "detraccion", "detracción"}:
        return "detracción"
    if focus in {"promotores", "promotor"}:
        return "promoción"
    return str(focus_name or "impacto").strip()


def _focus_probability_label(focus_name: str) -> str:
    focus_label = _focus_risk_label(focus_name)
    return f"Prob. de {focus_label}"


def _format_opportunity_scope(dimension: object, value: object) -> str:
    dim = str(dimension or "").strip().lower()
    clean_value = _clip(value, 42)
    if dim == "palanca":
        return f"{clean_value} (palanca)"
    if dim == "subpalanca":
        return f"{clean_value} (subpalanca)"
    if dim == "nps_topic":
        return clean_value
    return clean_value


def _clean_evidence_excerpt(text: object, *, max_len: int = 128) -> str:
    clean = " ".join(str(text or "").split())
    if not clean:
        return ""
    for marker in ["Síntoma:", "Sintoma:", "Descripcion:", "Descripción:"]:
        if marker in clean:
            clean = clean.split(marker, 1)[1].strip()
            break
    clean = re.sub(r"^(ACOTAMIENTO IRD|Acotamiento IRD)\s*", "", clean)
    return _clip(clean, max_len)


def _is_cover_metric_line(text: str) -> bool:
    low = str(text or "").strip().lower()
    return any(
        token in low
        for token in [
            "muestras",
            "comentarios analizados",
            "nps medio",
            "nps clásico",
            "nps clasico",
            "detractores",
            "promotores",
        ]
    )


def _slug(value: object, *, max_len: int = 42) -> str:
    raw = str(value or "").strip().lower()
    if not raw:
        return "na"
    norm = unicodedata.normalize("NFKD", raw).encode("ascii", "ignore").decode("ascii")
    norm = re.sub(r"[^a-z0-9]+", "-", norm)
    norm = re.sub(r"-{2,}", "-", norm).strip("-")
    if not norm:
        return "na"
    return norm[: int(max_len)].strip("-") or "na"


def _safe_date(value: object) -> str:
    try:
        ts = _coerce_datetime_scalar(value)
        return str(ts.date())
    except Exception:
        return str(value or "")


def _safe_dt(value: object) -> Optional[pd.Timestamp]:
    try:
        ts = _coerce_datetime_scalar(value)
    except Exception:
        return None
    if pd.isna(ts):
        return None
    return pd.Timestamp(ts)


def _month_label_es(d: date) -> str:
    months = {
        1: "enero",
        2: "febrero",
        3: "marzo",
        4: "abril",
        5: "mayo",
        6: "junio",
        7: "julio",
        8: "agosto",
        9: "septiembre",
        10: "octubre",
        11: "noviembre",
        12: "diciembre",
    }
    return f"{months.get(int(d.month), 'mes')} {int(d.year)}"


def _comment_column(df: pd.DataFrame) -> str:
    for candidate in ["Comment", "Comentario", "comment", "comentario"]:
        if candidate in df.columns:
            return candidate
    return ""


def _nps_band(value: object) -> str:
    score = _safe_float(value, default=float("nan"))
    if not np.isfinite(score):
        return "Sin dato"
    if score <= 6.0:
        return "Detractor"
    if score >= 9.0:
        return "Promotor"
    return "Pasivo"


_CATEGORY_CANONICALS = {
    "funcionamiento continuo": "Funcionamiento continuo",
    "agregar funcionalidad": "Agregar funcionalidad",
    "fallas en el login": "Fallas en el login",
}


def _category_key(value: object) -> str:
    clean = " ".join(str(value or "").split()).strip().casefold()
    if clean in {"", "nan", "none", "null", "<na>"}:
        return ""
    clean = unicodedata.normalize("NFKD", clean).encode("ascii", "ignore").decode("ascii")
    return re.sub(r"\s+", " ", clean).strip()


def _normalize_category_value(value: object) -> str:
    clean = " ".join(str(value or "").split()).strip()
    key = _category_key(clean)
    if not key:
        return ""
    return _CATEGORY_CANONICALS.get(key, clean)


def _normalize_presentation_categories(
    df: pd.DataFrame,
    *,
    columns: Iterable[str] = ("Palanca", "Subpalanca"),
) -> pd.DataFrame:
    if df is None or df.empty:
        return pd.DataFrame(columns=getattr(df, "columns", []))
    out = df.copy()
    for column in columns:
        if column in out.columns:
            out[column] = out[column].map(_normalize_category_value)
    return out


def _coerce_nps_records(nps_df: Optional[pd.DataFrame]) -> pd.DataFrame:
    cols = [
        "date",
        "NPS",
        "comment_txt",
        "nps_topic",
        "Palanca",
        "Subpalanca",
        "band",
    ]
    if nps_df is None or nps_df.empty:
        return pd.DataFrame(columns=cols)

    out = nps_df.copy()
    out["date"] = _coerce_datetime_series(out.get("Fecha")).dt.normalize()
    out["NPS"] = pd.to_numeric(out.get("NPS"), errors="coerce")
    comment_col = _comment_column(out)
    out["comment_txt"] = (
        out.get(comment_col, pd.Series([""] * len(out), index=out.index))
        .astype(str)
        .fillna("")
        .str.strip()
    )
    if "nps_topic" in out.columns:
        out["nps_topic"] = out["nps_topic"].astype(str).fillna("").str.strip()
    else:
        out["nps_topic"] = build_nps_topic(out).astype(str).fillna("").str.strip()
    out["Palanca"] = out.get("Palanca", pd.Series([""] * len(out), index=out.index)).map(
        _normalize_category_value
    )
    out["Subpalanca"] = out.get("Subpalanca", pd.Series([""] * len(out), index=out.index)).map(
        _normalize_category_value
    )
    out["band"] = out["NPS"].map(_nps_band)
    out = out.dropna(subset=["date"]).copy()
    return out[cols].copy()


def _split_period_frames(
    nps_df: pd.DataFrame,
    *,
    period_start: date,
    period_end: date,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    if nps_df.empty:
        return nps_df.copy(), nps_df.copy()
    start_ts = pd.Timestamp(period_start)
    end_ts = pd.Timestamp(period_end)
    current = nps_df[(nps_df["date"] >= start_ts) & (nps_df["date"] <= end_ts)].copy()
    baseline = nps_df[nps_df["date"] < start_ts].copy()
    if baseline.empty:
        baseline = nps_df[(nps_df["date"] < start_ts) | (nps_df["date"] > end_ts)].copy()
    return current, baseline


def _split_source_period_frames(
    nps_df: Optional[pd.DataFrame],
    *,
    period_start: date,
    period_end: date,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    if nps_df is None or nps_df.empty or "Fecha" not in nps_df.columns:
        return pd.DataFrame(), pd.DataFrame()

    out = nps_df.copy()
    out["Fecha"] = _coerce_datetime_series(out["Fecha"])
    out = out.dropna(subset=["Fecha"]).copy()
    if out.empty:
        return pd.DataFrame(), pd.DataFrame()
    out = _normalize_presentation_categories(out)

    start_ts = pd.Timestamp(period_start)
    end_ts = pd.Timestamp(period_end)
    current = out[(out["Fecha"] >= start_ts) & (out["Fecha"] <= end_ts)].copy()
    baseline = out[out["Fecha"] < start_ts].copy()
    if baseline.empty:
        baseline = out[(out["Fecha"] < start_ts) | (out["Fecha"] > end_ts)].copy()
    return current, baseline


def _period_overview(current_nps_df: pd.DataFrame) -> dict[str, object]:
    scores = pd.to_numeric(current_nps_df.get("NPS"), errors="coerce").dropna()
    comment_series = current_nps_df.get("comment_txt", pd.Series(dtype=str)).astype(str).str.strip()
    comment_count = (
        int(comment_series.ne("").sum()) if not comment_series.empty else int(len(scores))
    )
    detr = float((scores <= 6).mean()) if not scores.empty else 0.0
    prom = float((scores >= 9).mean()) if not scores.empty else 0.0
    pas = float(((scores >= 7) & (scores <= 8)).mean()) if not scores.empty else 0.0
    nps_mean = float(scores.mean()) if not scores.empty else float("nan")
    classic_nps = compute_nps_from_scores(scores) if not scores.empty else float("nan")
    daily_mix = _daily_group_mix(current_nps_df)
    start_classic = (
        float(pd.to_numeric(daily_mix["nps_classic"], errors="coerce").iloc[0])
        if not daily_mix.empty
        else float("nan")
    )
    end_classic = (
        float(pd.to_numeric(daily_mix["nps_classic"], errors="coerce").iloc[-1])
        if not daily_mix.empty
        else float("nan")
    )
    start_detr = (
        float(pd.to_numeric(daily_mix["detractor_rate"], errors="coerce").iloc[0])
        if not daily_mix.empty
        else float("nan")
    )
    end_detr = (
        float(pd.to_numeric(daily_mix["detractor_rate"], errors="coerce").iloc[-1])
        if not daily_mix.empty
        else float("nan")
    )
    driver_col = "Subpalanca" if current_nps_df.get("Subpalanca") is not None else "Palanca"
    if driver_col not in current_nps_df.columns:
        driver_col = "Palanca"
    pain_point = ""
    strength_point = ""
    if driver_col in current_nps_df.columns and "NPS" in current_nps_df.columns:
        driver_view = current_nps_df[[driver_col, "NPS"]].copy().dropna(subset=["NPS"])
        if not driver_view.empty:
            driver_view[driver_col] = driver_view[driver_col].astype(str).str.strip()
            driver_view = driver_view[driver_view[driver_col] != ""]
            if not driver_view.empty:
                ranking = (
                    driver_view.groupby(driver_col, dropna=False)
                    .agg(nps_mean=("NPS", "mean"), n=("NPS", "size"))
                    .sort_values(["nps_mean", "n"], ascending=[True, False])
                )
                if not ranking.empty:
                    pain_point = str(ranking.index[0])
                    strength_point = str(ranking.index[-1])
    return {
        "comments": comment_count,
        "nps_mean": nps_mean,
        "detractor_rate": detr,
        "promoter_rate": prom,
        "passive_rate": pas,
        "classic_nps": classic_nps,
        "start_classic": start_classic,
        "end_classic": end_classic,
        "classic_delta": (
            end_classic - start_classic
            if np.isfinite(start_classic) and np.isfinite(end_classic)
            else float("nan")
        ),
        "start_detr": start_detr,
        "end_detr": end_detr,
        "detractor_delta_pp": (
            ((end_detr - start_detr) * 100.0)
            if np.isfinite(start_detr) and np.isfinite(end_detr)
            else float("nan")
        ),
        "pain_point": pain_point,
        "strength_point": strength_point,
    }


def _daily_group_mix(nps_df: pd.DataFrame) -> pd.DataFrame:
    if nps_df is None or nps_df.empty:
        return pd.DataFrame(
            columns=[
                "date",
                "responses",
                "detractor_rate",
                "passive_rate",
                "promoter_rate",
                "nps_classic",
                "nps_mean",
            ]
        )

    work = nps_df.copy()
    work["date"] = _coerce_datetime_series(work["date"])
    work = work.dropna(subset=["date"]).copy()
    if work.empty:
        return pd.DataFrame()

    grouped = (
        work.groupby("date")
        .agg(
            responses=("NPS", "count"),
            detractors=("band", lambda s: int((s == "Detractor").sum())),
            passives=("band", lambda s: int((s == "Pasivo").sum())),
            promoters=("band", lambda s: int((s == "Promotor").sum())),
            nps_mean=("NPS", "mean"),
        )
        .reset_index()
        .sort_values("date")
    )
    grouped["detractor_rate"] = grouped["detractors"] / grouped["responses"].replace({0: np.nan})
    grouped["passive_rate"] = grouped["passives"] / grouped["responses"].replace({0: np.nan})
    grouped["promoter_rate"] = grouped["promoters"] / grouped["responses"].replace({0: np.nan})
    grouped["nps_classic"] = (grouped["promoter_rate"] - grouped["detractor_rate"]) * 100.0
    return grouped


def _merge_daily_incidents(base_df: pd.DataFrame, overall_daily: pd.DataFrame) -> pd.DataFrame:
    base = base_df.copy()
    if base.empty:
        return base
    out = base.copy()
    out["date"] = _coerce_datetime_series(out["date"]).dt.normalize()
    if overall_daily is None or overall_daily.empty or "date" not in overall_daily.columns:
        out["incidents"] = 0.0
        return out
    inc = overall_daily.copy()
    inc["date"] = _coerce_datetime_series(inc["date"]).dt.normalize()
    inc["incidents"] = pd.to_numeric(inc.get("incidents"), errors="coerce").fillna(0.0)
    inc = inc.groupby("date", as_index=False).agg(incidents=("incidents", "sum"))
    out = out.merge(inc, on="date", how="left")
    out["incidents"] = pd.to_numeric(out["incidents"], errors="coerce").fillna(0.0)
    return out


def _topic_summary(by_topic_daily: Optional[pd.DataFrame]) -> pd.DataFrame:
    cols = ["nps_topic", "comments", "focus_comments", "incidents", "nps_mean", "share"]
    if by_topic_daily is None or by_topic_daily.empty or "nps_topic" not in by_topic_daily.columns:
        return pd.DataFrame(columns=cols)

    d = by_topic_daily.copy()
    d["responses"] = pd.to_numeric(d.get("responses"), errors="coerce").fillna(0.0)
    d["focus_count"] = pd.to_numeric(d.get("focus_count"), errors="coerce").fillna(0.0)
    d["incidents"] = pd.to_numeric(d.get("incidents"), errors="coerce").fillna(0.0)
    d["nps_mean"] = pd.to_numeric(d.get("nps_mean"), errors="coerce")
    d["topic_nps_weight"] = d["nps_mean"].fillna(0.0) * d["responses"]
    out = (
        d.groupby("nps_topic", as_index=False)
        .agg(
            comments=("responses", "sum"),
            focus_comments=("focus_count", "sum"),
            incidents=("incidents", "sum"),
            topic_nps_weight=("topic_nps_weight", "sum"),
        )
        .sort_values(["comments", "focus_comments", "incidents"], ascending=False)
    )
    out["nps_mean"] = out["topic_nps_weight"] / out["comments"].replace({0.0: np.nan})
    out["share"] = (
        out["comments"] / out["comments"].sum() if float(out["comments"].sum()) > 0 else 0.0
    )
    out["nps_topic"] = out["nps_topic"].astype(str).str.strip()
    out = out[out["nps_topic"] != ""].copy()
    return out[cols]


def _text_topics_table(current_nps_df: pd.DataFrame, *, top_k: int = 10) -> pd.DataFrame:
    cols = ["cluster_id", "n", "top_terms", "examples", "label", "top_terms_txt", "example_txt"]
    if (
        current_nps_df is None
        or current_nps_df.empty
        or "comment_txt" not in current_nps_df.columns
    ):
        return pd.DataFrame(columns=cols)

    comments = current_nps_df["comment_txt"].astype(str).str.strip()
    comments = comments[comments.ne("")]
    if comments.empty:
        return pd.DataFrame(columns=cols)

    topics = extract_topics(comments, n_clusters=max(int(top_k), 10))
    if not topics:
        return pd.DataFrame(columns=cols)

    d = (
        pd.DataFrame([topic.__dict__ for topic in topics])
        .sort_values("n", ascending=False)
        .head(int(top_k))
        .copy()
    )
    d["label"] = d.apply(
        lambda row: f"#{int(row['cluster_id'])}: {', '.join(list(row['top_terms'])[:3])}",
        axis=1,
    )
    d["top_terms_txt"] = d["top_terms"].apply(
        lambda values: ", ".join([str(v).strip() for v in list(values)[:6] if str(v).strip()])
    )
    d["example_txt"] = d["examples"].apply(
        lambda values: " | ".join([_clip(v, 42) for v in list(values)[:2] if str(v).strip()])
    )
    return d[cols].reset_index(drop=True)


def _driver_change_table(
    current_nps_df: pd.DataFrame,
    baseline_nps_df: pd.DataFrame,
    *,
    dimension: str,
) -> pd.DataFrame:
    cols = [
        "value",
        "n_current",
        "nps_current",
        "nps_baseline",
        "delta_nps",
        "detr_current",
        "detr_baseline",
        "delta_detr_pp",
    ]
    if (
        current_nps_df is None
        or current_nps_df.empty
        or baseline_nps_df is None
        or baseline_nps_df.empty
        or dimension not in current_nps_df.columns
        or dimension not in baseline_nps_df.columns
    ):
        return pd.DataFrame(columns=cols)

    cur = pd.DataFrame([s.__dict__ for s in driver_table(current_nps_df, dimension=dimension)])
    base = pd.DataFrame([s.__dict__ for s in driver_table(baseline_nps_df, dimension=dimension)])
    if cur.empty or base.empty:
        return pd.DataFrame(columns=cols)

    merged = cur.rename(
        columns={
            "n": "n_current",
            "nps": "nps_current",
            "detractor_rate": "detr_current",
        }
    ).merge(
        base[["value", "nps", "detractor_rate"]],
        on="value",
        how="inner",
    )
    if merged.empty:
        return pd.DataFrame(columns=cols)
    merged = merged.rename(
        columns={
            "nps": "nps_baseline",
            "detractor_rate": "detr_baseline",
        }
    )
    merged["delta_nps"] = pd.to_numeric(merged["nps_current"], errors="coerce") - pd.to_numeric(
        merged["nps_baseline"], errors="coerce"
    )
    merged["delta_detr_pp"] = (
        pd.to_numeric(merged["detr_current"], errors="coerce")
        - pd.to_numeric(merged["detr_baseline"], errors="coerce")
    ) * 100.0
    return merged[cols].sort_values(["delta_nps", "n_current"], ascending=[True, False])


def _group_matrix(current_nps_df: pd.DataFrame, *, dimension: str, top_k: int = 8) -> pd.DataFrame:
    cols = [dimension, "band", "share", "count"]
    if current_nps_df is None or current_nps_df.empty or dimension not in current_nps_df.columns:
        return pd.DataFrame(columns=cols)
    d = current_nps_df.copy()
    d[dimension] = d[dimension].astype(str).fillna("").str.strip()
    d = d[d[dimension] != ""].copy()
    if d.empty:
        return pd.DataFrame(columns=cols)

    top_dims = (
        d.groupby(dimension, as_index=False)
        .agg(count=("NPS", "count"), detractors=("band", lambda s: int((s == "Detractor").sum())))
        .sort_values(["detractors", "count"], ascending=False)
        .head(int(top_k))
    )
    scoped = d[d[dimension].isin(top_dims[dimension].tolist())].copy()
    matrix = (
        scoped.groupby([dimension, "band"], as_index=False)
        .agg(count=("NPS", "count"))
        .merge(
            top_dims[[dimension, "count"]].rename(columns={"count": "dimension_total"}),
            on=dimension,
        )
    )
    matrix["share"] = matrix["count"] / matrix["dimension_total"].replace({0: np.nan})
    return matrix[cols]


def _gap_vs_overall_table(current_nps_df: pd.DataFrame, *, top_k: int = 10) -> pd.DataFrame:
    if current_nps_df is None or current_nps_df.empty:
        return pd.DataFrame(columns=["value", "n", "nps", "gap_vs_overall"])
    work = current_nps_df.copy()
    if "nps_topic" not in work.columns or work["nps_topic"].astype(str).str.strip().eq("").all():
        work["nps_topic"] = build_nps_topic(work).astype(str).fillna("").str.strip()
    stats = pd.DataFrame([s.__dict__ for s in driver_table(work, dimension="nps_topic")])
    if stats.empty:
        return pd.DataFrame(columns=["value", "n", "nps", "gap_vs_overall"])
    stats = stats[stats["n"] > 0].copy()
    stats = stats.sort_values(["gap_vs_overall", "n"], ascending=[True, False]).head(int(top_k))
    return stats[["value", "n", "nps", "gap_vs_overall"]]


def _dimension_gap_table(
    current_nps_df: pd.DataFrame,
    *,
    dimension: str,
    top_k: int = 10,
) -> pd.DataFrame:
    cols = ["value", "n", "nps", "gap_vs_overall"]
    if current_nps_df is None or current_nps_df.empty or dimension not in current_nps_df.columns:
        return pd.DataFrame(columns=cols)
    stats = pd.DataFrame([s.__dict__ for s in driver_table(current_nps_df, dimension=dimension)])
    if stats.empty:
        return pd.DataFrame(columns=cols)
    stats = stats[stats["n"] > 0].copy()
    stats = stats.sort_values(["gap_vs_overall", "n"], ascending=[True, False]).head(int(top_k))
    return stats[cols]


def _opportunities_table(
    current_nps_df: pd.DataFrame,
    *,
    dimension: str = "Palanca",
    min_n: int = 200,
) -> pd.DataFrame:
    cols = ["dimension", "value", "n", "current_nps", "potential_uplift", "confidence", "why"]
    if current_nps_df is None or current_nps_df.empty:
        return pd.DataFrame(columns=cols)
    work = _normalize_presentation_categories(current_nps_df, columns=[dimension])
    if "nps_topic" not in work.columns or work["nps_topic"].astype(str).str.strip().eq("").all():
        work["nps_topic"] = build_nps_topic(work).astype(str).fillna("").str.strip()
    dims = [str(dimension)] if str(dimension or "").strip() in work.columns else []
    if not dims:
        return pd.DataFrame(columns=cols)
    work = work[work[dimension].astype(str).str.strip().ne("")].copy()
    if work.empty:
        return pd.DataFrame(columns=cols)
    rows = rank_opportunities(
        work,
        dimensions=dims,
        min_n=max(1, int(min_n)),
    )
    if not rows:
        return pd.DataFrame(columns=cols)
    return pd.DataFrame([row.__dict__ for row in rows])[cols]


def _chain_portfolio_fig(
    chain_df: Optional[pd.DataFrame],
    *,
    highlight_topic: str,
) -> Optional[go.Figure]:
    if chain_df is None or chain_df.empty:
        return None
    d = chain_df.copy()
    d["priority"] = pd.to_numeric(d.get("priority"), errors="coerce").fillna(0.0)
    d["confidence"] = pd.to_numeric(d.get("confidence"), errors="coerce").fillna(0.0)
    d["impact"] = pd.to_numeric(d.get("total_nps_impact"), errors="coerce").fillna(0.0)
    d["links"] = pd.to_numeric(d.get("linked_pairs"), errors="coerce").fillna(0.0)
    if d.empty:
        return None
    d["is_highlight"] = (
        d.get("nps_topic", "").astype(str).str.strip() == str(highlight_topic).strip()
    )
    colors = [
        "#" + (BBVA_COLORS["red"] if is_h else BBVA_COLORS["sky"])
        for is_h in d["is_highlight"].tolist()
    ]
    fig = go.Figure(
        go.Scatter(
            x=d["confidence"],
            y=d["impact"],
            mode="markers+text",
            text=d["nps_topic"].astype(str).map(lambda value: _wrap_label(value, width=14)),
            textposition="top center",
            marker=dict(size=12 + d["links"] * 3.0, color=colors, opacity=0.88),
            customdata=d[["links"]].to_numpy(),
            hovertemplate=(
                "%{text}<br>Solidez=%{x:.2f}<br>Impacto=%{y:.2f} pts"
                "<br>Vínculos=%{customdata[0]:.0f}<extra></extra>"
            ),
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=16, t=20, b=30),
        xaxis=dict(title="Solidez de la evidencia", range=[0, 1]),
        yaxis=dict(title="Impacto total (pts NPS)", rangemode="tozero"),
        showlegend=False,
    )
    return fig


def _nps_evolution_fig(daily_mix: pd.DataFrame, overall_daily: pd.DataFrame) -> Optional[go.Figure]:
    if daily_mix is None or daily_mix.empty:
        return None
    d = _merge_daily_incidents(daily_mix, overall_daily)
    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["nps_classic"],
            name="NPS clásico",
            mode="lines+markers",
            line=dict(color="#" + BBVA_COLORS["blue"], width=3.0),
            marker=dict(size=7, color="#" + BBVA_COLORS["sky"]),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["detractor_rate"] * 100.0,
            name="% detractores",
            yaxis="y2",
            mode="lines",
            line=dict(color="#" + BBVA_COLORS["red"], width=2.4),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["incidents"],
            name="Incidencias",
            yaxis="y3",
            marker=dict(color="#" + BBVA_COLORS["yellow"]),
            opacity=0.70,
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=84, t=20, b=24),
        legend=dict(orientation="h", x=0.0, y=1.08),
        xaxis_title="Día",
        yaxis=dict(title="NPS clásico", rangemode="tozero"),
        yaxis2=dict(
            title="% detractores", overlaying="y", side="right", range=[0, 100], showgrid=False
        ),
        yaxis3=dict(
            title="Incidencias",
            overlaying="y",
            side="right",
            anchor="free",
            position=0.92,
            showgrid=False,
            rangemode="tozero",
        ),
    )
    return fig


def _top_topics_fig(topics_df: pd.DataFrame, *, top_k: int = 10) -> Optional[go.Figure]:
    if topics_df is None or topics_df.empty:
        return None
    d = topics_df.head(int(top_k)).copy().iloc[::-1]
    d["topic_label"] = d["nps_topic"].astype(str).map(lambda value: _wrap_label(value, width=26))
    fig = go.Figure(
        go.Bar(
            x=d["comments"],
            y=d["topic_label"],
            orientation="h",
            marker=dict(color=_ppt_nps_marker_colors(d["nps_mean"])),
            text=[f"{int(v)} · {s*100:.0f}%" for v, s in zip(d["comments"], d["share"])],
            textposition="outside",
            cliponaxis=False,
            customdata=np.column_stack([d["nps_mean"].fillna(np.nan), d["nps_topic"].astype(str)]),
            hovertemplate="%{customdata[1]}<br>Comentarios=%{x:.0f}<br>NPS medio=%{customdata[0]:.2f}<extra></extra>",
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=34, r=28, t=20, b=24),
        xaxis_title="Comentarios",
        yaxis_title="",
        showlegend=False,
    )
    return fig


def _topic_heatmap_fig(
    by_topic_daily: Optional[pd.DataFrame], *, top_k: int = 10
) -> Optional[go.Figure]:
    top = _topic_summary(by_topic_daily).head(int(top_k))
    if top.empty or by_topic_daily is None or by_topic_daily.empty:
        return None
    d = by_topic_daily.copy()
    d["date"] = _coerce_datetime_series(d["date"])
    d["responses"] = pd.to_numeric(d.get("responses"), errors="coerce").fillna(0.0)
    d = d[d["nps_topic"].astype(str).isin(top["nps_topic"].astype(str).tolist())].copy()
    if d.empty:
        return None
    pivot = (
        d.groupby(["nps_topic", "date"], as_index=False)
        .agg(responses=("responses", "sum"))
        .pivot(index="nps_topic", columns="date", values="responses")
        .fillna(0.0)
    )
    pivot = pivot.reindex(top["nps_topic"].tolist()).fillna(0.0)
    pivot.index = [_wrap_label(value, width=24) for value in pivot.index]
    fig = go.Figure(
        data=go.Heatmap(
            z=pivot.to_numpy(),
            x=list(pivot.columns),
            y=list(pivot.index),
            colorscale=plotly_continuous_scale(DesignTokens.default(), "light"),
            showscale=False,
            hovertemplate="%{y}<br>%{x|%d %b}: %{z:.0f} comentarios<extra></extra>",
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=16, t=20, b=30),
        xaxis_title="Día",
        yaxis_title="Top temas",
    )
    return fig


def _daily_group_mix_fig(daily_mix: pd.DataFrame) -> Optional[go.Figure]:
    if daily_mix is None or daily_mix.empty:
        return None
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=daily_mix["date"],
            y=daily_mix["promoter_rate"] * 100.0,
            name="Promotores",
            marker=dict(color="#" + BBVA_COLORS["green"]),
        )
    )
    fig.add_trace(
        go.Bar(
            x=daily_mix["date"],
            y=daily_mix["passive_rate"] * 100.0,
            name="Pasivos",
            marker=dict(color="#" + BBVA_COLORS["yellow"]),
        )
    )
    fig.add_trace(
        go.Bar(
            x=daily_mix["date"],
            y=daily_mix["detractor_rate"] * 100.0,
            name="Detractores",
            marker=dict(color="#" + BBVA_COLORS["red"]),
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=16, t=20, b=24),
        barmode="stack",
        legend=dict(orientation="h", x=0.0, y=1.08),
        xaxis_title="Día",
        yaxis=dict(title="% sobre comentarios del día", range=[0, 100]),
    )
    return fig


def _delta_bars_fig(change_df: pd.DataFrame, *, metric: str, x_title: str) -> Optional[go.Figure]:
    if change_df is None or change_df.empty or metric not in change_df.columns:
        return None
    d = change_df.copy()
    d[metric] = pd.to_numeric(d[metric], errors="coerce")
    d = d.dropna(subset=[metric]).copy()
    if d.empty:
        return None
    d["abs_metric"] = d[metric].abs()
    d = d.sort_values(["abs_metric", "n_current"], ascending=[False, False]).head(8).iloc[::-1]
    d["color"] = np.where(d[metric] < 0, "#" + BBVA_COLORS["red"], "#" + BBVA_COLORS["green"])
    d["axis_label"] = d["value"].astype(str).map(lambda value: _wrap_label(value, width=20))
    fig = go.Figure(
        go.Bar(
            x=d[metric],
            y=d["axis_label"],
            orientation="h",
            marker=dict(color=d["color"].tolist()),
            text=[_fmt_signed_or_nd(v, decimals=1) for v in d[metric].tolist()],
            textposition="outside",
            cliponaxis=False,
            showlegend=False,
        )
    )
    fig.add_vline(x=0, line_color="#" + BBVA_COLORS["line"], line_width=1)
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=24, t=20, b=24),
        xaxis_title=x_title,
        yaxis_title="",
    )
    return fig


def _group_heatmap_fig(matrix_df: pd.DataFrame, *, dimension: str) -> Optional[go.Figure]:
    if matrix_df is None or matrix_df.empty or dimension not in matrix_df.columns:
        return None
    pivot = (
        matrix_df.pivot(index=dimension, columns="band", values="share")
        .fillna(0.0)
        .reindex(columns=["Detractor", "Pasivo", "Promotor"], fill_value=0.0)
    )
    if pivot.empty:
        return None
    pivot.index = [_wrap_label(value, width=18) for value in pivot.index]
    fig = go.Figure(
        data=go.Heatmap(
            z=(pivot.to_numpy() * 100.0),
            x=list(pivot.columns),
            y=list(pivot.index),
            colorscale=plotly_risk_scale(DesignTokens.default(), "light"),
            showscale=False,
            text=np.vectorize(lambda v: f"{v:.0f}%")(pivot.to_numpy() * 100.0),
            texttemplate="%{text}",
            hovertemplate="%{y}<br>%{x}: %{z:.1f}%<extra></extra>",
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=16, t=20, b=24),
        xaxis_title="Grupo NPS",
        yaxis_title="",
    )
    return fig


def _gap_vs_overall_fig(gap_df: pd.DataFrame) -> Optional[go.Figure]:
    if gap_df is None or gap_df.empty:
        return None
    d = gap_df.copy().iloc[::-1]
    d["axis_label"] = d["value"].astype(str).map(lambda value: _wrap_label(value, width=22))
    fig = go.Figure(
        go.Bar(
            x=d["gap_vs_overall"],
            y=d["axis_label"],
            orientation="h",
            marker=dict(color="#" + BBVA_COLORS["red"]),
            text=[_fmt_signed_or_nd(v, decimals=1) for v in d["gap_vs_overall"].tolist()],
            textposition="outside",
            cliponaxis=False,
        )
    )
    fig.add_vline(x=0, line_color="#" + BBVA_COLORS["line"], line_width=1)
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=32, r=28, t=20, b=24),
        xaxis_title="Diferencia frente al NPS medio",
        yaxis_title="",
        showlegend=False,
    )
    return fig


def _opportunity_bubble_fig(opps_df: pd.DataFrame) -> Optional[go.Figure]:
    if opps_df is None or opps_df.empty:
        return None
    d = opps_df.head(10).copy()
    d["size"] = np.clip(np.sqrt(pd.to_numeric(d["n"], errors="coerce").fillna(0.0)) * 2.4, 12, 40)
    d["plot_label"] = d.apply(
        lambda row: _wrap_label(
            _format_opportunity_scope(row.get("dimension"), row.get("value")), width=16
        ),
        axis=1,
    )
    colors = {
        "Palanca": "#" + BBVA_COLORS["blue"],
        "Subpalanca": "#" + BBVA_COLORS["orange"],
        "nps_topic": "#" + BBVA_COLORS["red"],
    }
    fig = go.Figure(
        go.Scatter(
            x=d["confidence"],
            y=d["potential_uplift"],
            mode="markers+text",
            text=d["plot_label"],
            textposition="top center",
            marker=dict(
                size=d["size"],
                color=[
                    colors.get(str(dim), "#" + BBVA_COLORS["sky"])
                    for dim in d["dimension"].tolist()
                ],
                opacity=0.82,
            ),
            hovertemplate=(
                "%{customdata[2]}<br>Dimensión=%{customdata[0]}<br>Potencial=%{y:.1f} pts"
                "<br>Confianza=%{x:.2f}<br>n=%{customdata[1]:.0f}<extra></extra>"
            ),
            customdata=d[["dimension", "n", "value"]].to_numpy(),
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=24, t=20, b=24),
        xaxis=dict(title="Solidez de la evidencia", range=[0, 1]),
        yaxis=dict(title="Impacto potencial (pts NPS)", rangemode="tozero"),
        showlegend=False,
    )
    return fig


def _build_overview_figure(
    selected_nps_df: Optional[pd.DataFrame], *, period_days: int
) -> Optional[go.Figure]:
    fig = chart_daily_kpis(
        selected_nps_df.copy() if selected_nps_df is not None else pd.DataFrame(),
        get_theme("light"),
        days=max(int(period_days), 1),
    )
    if fig is None:
        return None
    fig.update_layout(
        legend=dict(orientation="h", x=0.0, y=1.18, yanchor="bottom", title_text=""),
        margin=dict(l=58, r=58, t=88, b=74),
    )
    fig.update_xaxes(
        side="bottom",
        ticklabelposition="outside bottom",
        automargin=True,
        tickfont=dict(size=20),
        title_font=dict(size=20),
    )
    fig.update_yaxes(tickfont=dict(size=18), title_font=dict(size=18))
    return fig


def _build_text_topic_figure(text_topics_df: pd.DataFrame) -> Optional[go.Figure]:
    topic_fig = chart_topic_bars(
        text_topics_df, get_theme("light"), top_k=EDITORIAL_LIMITS.max_text_clusters
    )
    if topic_fig is None or text_topics_df.empty:
        return topic_fig
    topic_rows = (
        text_topics_df.sort_values("n", ascending=False)
        .head(EDITORIAL_LIMITS.max_text_clusters)
        .copy()
    )
    topic_labels = []
    for row in topic_rows.itertuples():
        terms = [str(term).strip() for term in list(row.top_terms)[:2] if str(term).strip()]
        topic_labels.append(_clip(f"#{int(row.cluster_id)} · {', '.join(terms)}", 24))
    with contextlib.suppress(Exception):
        topic_fig.data[0].y = topic_labels
    counts = pd.to_numeric(text_topics_df.get("n"), errors="coerce").dropna()
    xmax = float(counts.max()) if not counts.empty else 0.0
    topic_fig.update_xaxes(
        title_text="Comentarios",
        nticks=5,
        dtick=2000 if xmax >= 6000 else 1000 if xmax >= 2000 else None,
        tickformat="~s",
        tickfont=dict(size=20),
        title_font=dict(size=20),
    )
    topic_fig.update_yaxes(tickfont=dict(size=28), automargin=True)
    topic_fig.update_layout(margin=dict(l=250, r=24, t=22, b=54), bargap=0.34)
    return topic_fig


def _build_driver_delta_figure(
    delta_df: pd.DataFrame, *, panel_height_in: float
) -> Optional[go.Figure]:
    fig = chart_driver_delta(delta_df, get_theme("light"), top_k=max(len(delta_df), 1))
    if fig is None or delta_df.empty:
        return fig
    plot_df = delta_df.head(EDITORIAL_LIMITS.max_change_rows).copy()
    label_count = len(plot_df)
    max_len = int(plot_df["value"].astype(str).str.len().max() or 0)
    wrap_width = 26 if max_len >= 26 else 22 if max_len >= 18 else 18
    left_margin = 430 if max_len >= 34 else 380 if max_len >= 26 else 330 if max_len >= 18 else 285
    y_font_size = 34 if label_count <= 5 else 30 if label_count <= 7 else 26
    labels = [
        _wrap_label(value, width=wrap_width, max_lines=2, joiner="<br>")
        for value in plot_df["value"].astype(str).tolist()
    ]
    with contextlib.suppress(Exception):
        fig.data[0].y = labels
    fig.update_yaxes(
        tickfont=dict(size=y_font_size, family=BBVA_FONT_MEDIUM),
        automargin=True,
        title_text="",
    )
    fig.update_xaxes(
        title_text="Delta NPS",
        tickfont=dict(size=19),
        title_font=dict(size=20),
        nticks=5,
        zeroline=True,
        zerolinecolor="#" + BBVA_COLORS["line"],
    )
    fig.update_layout(
        margin=dict(l=left_margin, r=48, t=20, b=58 if panel_height_in <= 3.8 else 52),
        bargap=0.30,
    )
    return fig


def _build_web_heatmap_figure(source_df: pd.DataFrame, *, row_dim: str) -> Optional[go.Figure]:
    required = {row_dim, "Canal", "NPS"}
    if source_df is None or source_df.empty or not required.issubset(set(source_df.columns)):
        return None
    chart_df = _normalize_presentation_categories(source_df, columns=[row_dim])
    chart_df = chart_df.dropna(subset=[row_dim, "Canal", "NPS"]).copy()
    if chart_df.empty:
        return None
    chart_df[row_dim] = chart_df[row_dim].astype(str).str.strip()
    chart_df["Canal"] = chart_df["Canal"].astype(str).str.strip()
    chart_df = chart_df[
        chart_df[row_dim].ne("") & chart_df["Canal"].str.casefold().eq("web")
    ].copy()
    if chart_df.empty:
        return None
    fig = chart_cohort_heatmap(
        chart_df, get_theme("light"), row_dim=row_dim, col_dim="Canal", min_n=1
    )
    if fig is None:
        return None
    row_stats = (
        chart_df.groupby(row_dim, as_index=False)
        .agg(n=("NPS", "size"), nps=("NPS", "mean"))
        .sort_values(["nps", "n"], ascending=[True, False])
        .head(EDITORIAL_LIMITS.max_web_rows)
    )
    labels = [
        _compact_axis_label(value, width=24, max_lines=2, max_chars=44).replace("<br>", " ")
        for value in row_stats[row_dim].astype(str).tolist()
    ]
    with contextlib.suppress(Exception):
        fig.data[0].y = labels
    fig.update_yaxes(
        tickfont=dict(size=30 if len(labels) <= 5 else 26, family=BBVA_FONT_MEDIUM),
        automargin=True,
        title_text="",
    )
    fig.update_xaxes(
        side="bottom", tickangle=0, tickfont=dict(size=22), automargin=True, title_text=""
    )
    fig.update_layout(margin=dict(l=390, r=116, t=18, b=54))
    fig.update_coloraxes(
        showscale=True,
        colorbar=dict(
            title="NPS",
            tickmode="array",
            tickvals=[0, 2, 6, 8, 10],
            len=0.74,
            y=0.5,
            thickness=14,
        ),
    )
    return fig


def _build_web_dimension_table(source_df: pd.DataFrame, *, dimension: str) -> pd.DataFrame:
    cols = ["value", "n", "nps", "detractor_rate"]
    required = {dimension, "Canal", "NPS"}
    if source_df is None or source_df.empty or not required.issubset(set(source_df.columns)):
        return pd.DataFrame(columns=cols)
    work = _normalize_presentation_categories(source_df, columns=[dimension])
    work = work.dropna(subset=[dimension, "Canal", "NPS"]).copy()
    work[dimension] = work[dimension].astype(str).str.strip()
    work["Canal"] = work["Canal"].astype(str).str.strip()
    work["NPS"] = pd.to_numeric(work["NPS"], errors="coerce")
    work = (
        work[work[dimension].ne("") & work["Canal"].str.casefold().eq("web")]
        .dropna(subset=["NPS"])
        .copy()
    )
    if work.empty:
        return pd.DataFrame(columns=cols)
    out = (
        work.groupby(dimension, as_index=False)
        .agg(
            n=("NPS", "size"),
            nps=("NPS", "mean"),
            detractor_rate=(
                "NPS",
                lambda s: float((pd.to_numeric(s, errors="coerce") <= 6).mean()),
            ),
        )
        .rename(columns={dimension: "value"})
    )
    return (
        out.sort_values(["nps", "n"], ascending=[True, False])
        .head(EDITORIAL_LIMITS.max_web_rows)[cols]
        .copy()
    )


def _build_opportunity_figure(opp_df: pd.DataFrame) -> Optional[go.Figure]:
    fig = chart_opportunities_bar(opp_df, get_theme("light"), top_k=max(len(opp_df), 1))
    if fig is None or opp_df.empty:
        return fig
    plot_df = select_opportunities(opp_df, max_rows=EDITORIAL_LIMITS.max_opportunities)
    label_count = len(plot_df)
    label_lengths = (
        plot_df["label"].astype(str).str.replace("<br>", " ", regex=False).str.len()
        if "label" in plot_df.columns
        else pd.Series(dtype=float)
    )
    max_len = int(label_lengths.max() or 0)
    y_font_size = 30 if label_count <= 5 else 27 if label_count <= 7 else 24
    left_margin = 340 if max_len >= 28 else 300 if max_len >= 22 else 255
    uplift = pd.to_numeric(plot_df.get("potential_uplift"), errors="coerce")
    text_values = [
        _fmt_signed_or_nd(value, decimals=1) if np.isfinite(value) else ""
        for value in uplift.tolist()
    ]
    with contextlib.suppress(Exception):
        fig.data[0].text = text_values
        fig.data[0].textposition = "outside"
        fig.data[0].cliponaxis = False
        fig.data[0].textfont.size = 20
    fig.update_yaxes(
        title_text="", tickfont=dict(size=y_font_size, family=BBVA_FONT_MEDIUM), automargin=True
    )
    fig.update_xaxes(
        title_text="Impacto estimado", tickfont=dict(size=19), title_font=dict(size=20), nticks=5
    )
    fig.update_layout(margin=dict(l=left_margin, r=72, t=18, b=54), bargap=0.34)
    return fig


def _prepare_opportunity_chart_df(opportunities_df: pd.DataFrame) -> pd.DataFrame:
    opp_chart_df = select_opportunities(
        opportunities_df, max_rows=EDITORIAL_LIMITS.max_opportunities
    )
    if opp_chart_df.empty:
        return opp_chart_df

    def _opp_label(row: pd.Series) -> str:
        value = str(row.get("value", "")).strip()
        base = value
        return _compact_axis_label(
            base, width=22 if len(base) >= 22 else 18, max_lines=2, max_chars=38
        )

    opp_chart_df["label"] = opp_chart_df.apply(_opp_label, axis=1)
    return opp_chart_df


def _build_journey_table(
    *,
    touchpoint_source: str,
    entity_summary_df: pd.DataFrame,
    broken_journeys_df: Optional[pd.DataFrame],
) -> pd.DataFrame:
    cols = ["journey", "touchpoint", "links", "comments", "nps", "confidence"]
    if (
        str(touchpoint_source or "").strip() == TOUCHPOINT_SOURCE_BROKEN_JOURNEYS
        and broken_journeys_df is not None
        and not broken_journeys_df.empty
    ):
        source = broken_journeys_df.copy()
        source["priority_sort"] = pd.to_numeric(source.get("linked_pairs"), errors="coerce").fillna(
            0.0
        )
        source["confidence_sort"] = pd.to_numeric(
            source.get("semantic_cohesion"), errors="coerce"
        ).fillna(0.0)
        source["nps_sort"] = pd.to_numeric(source.get("avg_nps"), errors="coerce").fillna(10.0)
        out = (
            source.sort_values(
                ["priority_sort", "confidence_sort", "nps_sort"], ascending=[False, False, True]
            )
            .head(EDITORIAL_LIMITS.max_journey_rows)
            .copy()
        )
        return pd.DataFrame(
            {
                "journey": out.get("journey_label", pd.Series(dtype=str)).astype(str),
                "touchpoint": out.get("touchpoint", pd.Series(dtype=str)).astype(str),
                "links": pd.to_numeric(out.get("linked_pairs"), errors="coerce")
                .fillna(0)
                .astype(int),
                "comments": pd.to_numeric(out.get("linked_comments"), errors="coerce")
                .fillna(0)
                .astype(int),
                "nps": pd.to_numeric(out.get("avg_nps"), errors="coerce"),
                "confidence": pd.to_numeric(out.get("semantic_cohesion"), errors="coerce").fillna(
                    0.0
                ),
            }
        )
    if entity_summary_df is None or entity_summary_df.empty:
        return pd.DataFrame(columns=cols)
    source = select_causal_scenarios(entity_summary_df, max_rows=EDITORIAL_LIMITS.max_journey_rows)
    return pd.DataFrame(
        {
            "journey": source.get("nps_topic", pd.Series(dtype=str)).astype(str),
            "touchpoint": source.get("touchpoint", pd.Series(dtype=str)).astype(str),
            "links": pd.to_numeric(source.get("linked_pairs"), errors="coerce")
            .fillna(0)
            .astype(int),
            "comments": pd.to_numeric(source.get("linked_comments"), errors="coerce")
            .fillna(0)
            .astype(int),
            "nps": pd.to_numeric(source.get("avg_nps"), errors="coerce"),
            "confidence": pd.to_numeric(source.get("confidence"), errors="coerce").fillna(0.0),
        }
    )


def _build_journey_summary_figure(
    summary_df: pd.DataFrame, *, touchpoint_source: str
) -> Optional[go.Figure]:
    method_spec = get_causal_method_spec(touchpoint_source)
    plot_df = summary_df.copy() if summary_df is not None else pd.DataFrame()
    if plot_df.empty:
        return None
    plot_df["entity_label"] = plot_df.get("nps_topic", "").astype(str).str.strip()
    fig = chart_causal_entity_bar(
        plot_df,
        get_theme("light"),
        entity_label=method_spec.entity_singular,
        top_k=min(EDITORIAL_LIMITS.max_journey_rows, len(plot_df)) if not plot_df.empty else 10,
    )
    if fig is None or plot_df.empty:
        return fig
    y_values = [str(value).replace("<br>", " ") for value in list(getattr(fig.data[0], "y", []))]
    if not y_values:
        return fig
    max_len = max(len(value) for value in y_values)
    label_count = len(y_values)
    wrap_width = 44 if max_len <= 46 else 40 if max_len <= 58 else 36
    max_chars = 92 if max_len <= 72 else 84
    y_font_size = 27 if label_count <= 6 else 24 if label_count <= 8 else 22
    left_margin = 430 if max_len >= 72 else 390 if max_len >= 58 else 350 if max_len >= 44 else 320
    pretty_labels = [
        _compact_axis_label(value, width=wrap_width, max_lines=2, max_chars=max_chars)
        for value in y_values
    ]
    with contextlib.suppress(Exception):
        fig.data[0].y = pretty_labels
    fig.update_yaxes(
        title_text="",
        tickmode="array",
        tickvals=pretty_labels,
        ticktext=pretty_labels,
        tickfont=dict(size=y_font_size, family=BBVA_FONT_MEDIUM, color="#" + BBVA_COLORS["ink"]),
        automargin=True,
    )
    fig.update_xaxes(
        title_text="Links validados Helix↔VoC",
        tickfont=dict(size=17),
        title_font=dict(size=18),
        nticks=6,
        automargin=True,
    )
    fig.update_layout(margin=dict(l=left_margin, r=102, t=14, b=38), bargap=0.22)
    fig.update_coloraxes(
        colorbar=dict(
            title=dict(text="NPS en riesgo", side="right", font=dict(size=15)),
            tickmode="array",
            tickvals=[0, 1, 2, 3, 4],
            tickfont=dict(size=14),
            len=0.82,
            y=0.5,
            thickness=16,
        )
    )
    return fig


def _causal_daily_timeline_fig(
    daily_mix: pd.DataFrame, overall_daily: pd.DataFrame
) -> Optional[go.Figure]:
    if daily_mix is None or daily_mix.empty:
        return None
    d = _merge_daily_incidents(daily_mix, overall_daily)
    if d.empty:
        return None
    tokens = DesignTokens.default()
    pal = palette(tokens, "light")
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["incidents"],
            name="# incidencias",
            yaxis="y2",
            opacity=0.75,
            marker=dict(color=pal["color.primary.accent.value-01.default"]),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["detractor_rate"] * 100.0,
            mode="lines+markers",
            name="% detractores",
            line=dict(color=pal["color.primary.bg.alert"], width=2),
            marker=dict(color=pal["color.primary.bg.alert"], size=6),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["passive_rate"] * 100.0,
            mode="lines+markers",
            name="% pasivos",
            line=dict(color=pal["color.primary.bg.warning"], width=2),
            marker=dict(color=pal["color.primary.bg.warning"], size=6),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["promoter_rate"] * 100.0,
            mode="lines+markers",
            name="% promotores",
            line=dict(color=pal["color.primary.bg.success"], width=2),
            marker=dict(color=pal["color.primary.bg.success"], size=6),
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=48, t=20, b=24),
        legend=dict(orientation="h", x=0.0, y=1.08),
        xaxis_title="Día",
        yaxis=dict(title="Tasa por grupo", tickformat=".0%"),
        yaxis2=dict(
            title="Incidencias",
            overlaying="y",
            side="right",
            rangemode="tozero",
            showgrid=False,
        ),
    )
    return fig


def _patch_kaleido_executable_for_space_paths() -> None:
    """Patch kaleido executable lookup when project path contains spaces."""
    try:
        from kaleido.scopes import base as kaleido_base
    except Exception:
        return

    cls = kaleido_base.BaseScope
    if getattr(cls, "_nps_lens_kaleido_patched", False):
        return

    try:
        default_exec = str(cls.executable_path())
    except Exception:
        return

    if " " not in default_exec or os.name == "nt":
        cls._nps_lens_kaleido_patched = True
        return

    exec_path = Path(default_exec)
    exec_dir = exec_path.parent
    real_bin = exec_dir / "bin" / "kaleido"
    if not real_bin.exists():
        cls._nps_lens_kaleido_patched = True
        return

    shim_dir = Path(tempfile.gettempdir()) / "nps_lens_kaleido"
    shim_dir.mkdir(parents=True, exist_ok=True)
    shim_path = shim_dir / "kaleido-shim"
    shim_path.write_text(
        f'#!/bin/sh\ncd "{exec_dir}" || exit 1\nexec "./bin/kaleido" "$@"\n',
        encoding="utf-8",
    )
    with contextlib.suppress(Exception):
        shim_path.chmod(0o755)

    cls.executable_path = classmethod(lambda scope_cls: str(shim_path))  # type: ignore[assignment]
    cls._nps_lens_kaleido_patched = True


def _apply_ppt_figure_theme(
    fig: go.Figure,
    *,
    panel_width_in: float | None = None,
    panel_height_in: float | None = None,
) -> go.Figure:
    ink = "#" + BBVA_COLORS["ink"]
    grid = "#" + BBVA_COLORS["line"]
    white = "#" + BBVA_COLORS["white"]
    compact_panel = bool(
        (panel_width_in is not None and panel_width_in <= 6.0)
        or (panel_height_in is not None and panel_height_in <= 2.5)
    )
    trace_types = {
        str(getattr(trace, "type", "") or "").strip().lower()
        for trace in fig.data
        if trace is not None
    }
    has_heatmap = "heatmap" in trace_types
    has_legend = (
        sum(
            1
            for trace in fig.data
            if bool(getattr(trace, "showlegend", True))
            and str(getattr(trace, "name", "") or "").strip()
        )
        > 1
    )
    has_scatter_text = any(
        str(getattr(trace, "type", "") or "").strip().lower() == "scatter"
        and "text" in str(getattr(trace, "mode", "") or "").lower()
        for trace in fig.data
    )
    has_colorbar = any(
        bool(getattr(trace, "showscale", False))
        or bool(getattr(getattr(trace, "colorbar", None), "title", None))
        for trace in fig.data
    )

    for trace in fig.data:
        name = str(getattr(trace, "name", "") or "").strip().lower()
        trace_type = str(getattr(trace, "type", "") or "").strip().lower()
        is_incidents = "incid" in name or "helix" in name
        is_detractor = any(token in name for token in ["detrac", "crit", "alto", "foco"])
        is_passive = any(token in name for token in ["pasiv", "moderad"])
        is_promoter = "promot" in name
        is_nps = "nps" in name

        if trace_type == "bar":
            color = None
            if is_incidents:
                color = "#" + BBVA_COLORS["sky"]
            elif is_promoter:
                color = "#" + BBVA_COLORS["green"]
            elif is_passive:
                color = "#" + BBVA_COLORS["yellow"]
            elif is_detractor or is_nps:
                color = "#" + BBVA_COLORS["red"]
            if color:
                with contextlib.suppress(Exception):
                    trace.marker.color = color
        elif trace_type == "scatter":
            if is_nps:
                with contextlib.suppress(Exception):
                    trace.line.color = "#" + BBVA_COLORS["blue"]
                if "markers" in str(getattr(trace, "mode", "") or ""):
                    with contextlib.suppress(Exception):
                        if not isinstance(
                            getattr(trace.marker, "color", None), (list, tuple, np.ndarray)
                        ):
                            trace.marker.color = "#" + BBVA_COLORS["sky"]
                    with contextlib.suppress(Exception):
                        trace.marker.size = max(8, int(getattr(trace.marker, "size", 8) or 8))
            elif is_incidents:
                with contextlib.suppress(Exception):
                    trace.line.color = "#" + BBVA_COLORS["sky"]
                with contextlib.suppress(Exception):
                    trace.marker.color = "#" + BBVA_COLORS["sky"]
            elif is_promoter:
                with contextlib.suppress(Exception):
                    trace.line.color = "#" + BBVA_COLORS["green"]
            elif is_passive:
                with contextlib.suppress(Exception):
                    trace.line.color = "#" + BBVA_COLORS["yellow"]
            elif is_detractor:
                with contextlib.suppress(Exception):
                    trace.line.color = "#" + BBVA_COLORS["red"]

    current_margin = fig.layout.margin.to_plotly_json() if fig.layout.margin else {}
    base_font_size = 18 if has_heatmap or compact_panel else 17
    tick_font_size = 16 if has_heatmap or compact_panel else 15
    axis_title_font_size = 18 if has_heatmap or compact_panel else 17
    legend_font_size = 15 if compact_panel else 14

    def _font_size(value: object, fallback: int) -> int:
        try:
            size = int(float(value))
        except Exception:
            return int(fallback)
        return max(size, int(fallback))

    fig.update_layout(
        template="plotly_white",
        paper_bgcolor=white,
        plot_bgcolor=white,
        font=dict(family=BBVA_FONT_BODY, size=base_font_size, color=ink),
        legend=dict(
            orientation="h",
            x=0.0,
            xanchor="left",
            y=1.12 if has_legend else 1.04,
            yanchor="bottom",
            font=dict(size=legend_font_size, color=ink),
            title_font=dict(size=legend_font_size, color=ink),
            bgcolor="rgba(0,0,0,0)",
        ),
        margin=dict(
            l=max(
                int(current_margin.get("l", 24)),
                78 if has_scatter_text or compact_panel else 58,
            ),
            r=max(int(current_margin.get("r", 24)), 92 if has_colorbar else 44),
            t=max(int(current_margin.get("t", 20)), 40 if has_heatmap or has_legend else 28),
            b=max(
                int(current_margin.get("b", 24)),
                84 if has_heatmap or compact_panel else 58,
            ),
        ),
        hoverlabel=dict(font=dict(family=BBVA_FONT_BODY, size=13, color=ink)),
    )
    fig.for_each_xaxis(
        lambda axis: axis.update(
            tickfont=dict(
                size=_font_size(
                    getattr(getattr(axis, "tickfont", None), "size", None), tick_font_size
                ),
                color=ink,
            ),
            title_font=dict(
                size=_font_size(
                    getattr(getattr(getattr(axis, "title", None), "font", None), "size", None),
                    axis_title_font_size,
                ),
                color=ink,
            ),
            automargin=True,
            gridcolor=grid,
            linecolor=grid,
        )
    )
    fig.for_each_yaxis(
        lambda axis: axis.update(
            tickfont=dict(
                size=_font_size(
                    getattr(getattr(axis, "tickfont", None), "size", None), tick_font_size
                ),
                color=ink,
            ),
            title_font=dict(
                size=_font_size(
                    getattr(getattr(getattr(axis, "title", None), "font", None), "size", None),
                    axis_title_font_size,
                ),
                color=ink,
            ),
            automargin=True,
            gridcolor=grid,
            linecolor=grid,
        )
    )
    for trace in fig.data:
        if "text" not in str(getattr(trace, "mode", "") or "").lower():
            continue
        with contextlib.suppress(Exception):
            trace.textfont.size = max(
                int(getattr(getattr(trace, "textfont", None), "size", 0) or 0),
                16,
            )
    if has_heatmap:
        for trace in fig.data:
            if str(getattr(trace, "type", "") or "").strip().lower() != "heatmap":
                continue
            with contextlib.suppress(Exception):
                trace.xgap = max(int(getattr(trace, "xgap", 0) or 0), 2)
            with contextlib.suppress(Exception):
                trace.ygap = max(int(getattr(trace, "ygap", 0) or 0), 2)
            if has_colorbar:
                with contextlib.suppress(Exception):
                    trace.colorbar.thickness = 14
                with contextlib.suppress(Exception):
                    trace.colorbar.len = 0.74
                with contextlib.suppress(Exception):
                    trace.colorbar.y = 0.48
                with contextlib.suppress(Exception):
                    trace.colorbar.title.side = "right"
                with contextlib.suppress(Exception):
                    trace.colorbar.tickfont.size = 15
                with contextlib.suppress(Exception):
                    trace.colorbar.title.font.size = 16
    for axis_name in fig.layout:
        if not str(axis_name).startswith(("xaxis", "yaxis")):
            continue
        axis = getattr(fig.layout, axis_name, None)
        if axis is None:
            continue
        with contextlib.suppress(Exception):
            if compact_panel and str(axis_name).startswith("xaxis") and axis.nticks is None:
                axis.nticks = 5
        with contextlib.suppress(Exception):
            if compact_panel and str(axis_name).startswith("yaxis") and axis.nticks is None:
                axis.nticks = 8
    return fig


def _pillow_color(
    value: object, default: str = "#42526E"
) -> tuple[int, int, int]:  # pragma: no cover
    try:
        from PIL import ImageColor
    except Exception:
        return (66, 82, 110)

    if isinstance(value, (list, tuple)) and len(value) >= 3:
        with contextlib.suppress(Exception):
            return tuple(int(float(channel)) for channel in value[:3])

    raw = str(value or "").strip()
    if not raw:
        raw = default
    if raw.lower().startswith("rgba("):
        parts = [part.strip() for part in raw[5:-1].split(",")]
        if len(parts) >= 3:
            with contextlib.suppress(Exception):
                return tuple(int(float(channel)) for channel in parts[:3])
    try:
        return ImageColor.getrgb(raw)
    except Exception:
        try:
            return ImageColor.getrgb(default)
        except Exception:
            return (66, 82, 110)


def _pillow_font(size: int, *, bold: bool = False):  # pragma: no cover
    try:
        from PIL import ImageFont
    except Exception:
        return None

    fonts_dir = Path(__file__).resolve().parents[3] / "assets" / "ppt" / "bbva" / "fonts"
    candidates = [
        fonts_dir / ("BentonSansBBVA-Bold.ttf" if bold else "BentonSansBBVA-Book.ttf"),
        fonts_dir / "BentonSansBBVA-Medium.ttf",
    ]
    for candidate in candidates:
        if candidate.exists():
            with contextlib.suppress(Exception):
                return ImageFont.truetype(str(candidate), size=max(int(size), 8))
    with contextlib.suppress(Exception):
        return ImageFont.load_default()
    return None


def _pillow_text_size(draw: object, text: str, font: object) -> tuple[int, int]:  # pragma: no cover
    with contextlib.suppress(Exception):
        bbox = draw.multiline_textbbox((0, 0), str(text or ""), font=font, spacing=4)
        return max(int(bbox[2] - bbox[0]), 0), max(int(bbox[3] - bbox[1]), 0)
    return (0, 0)


def _plotly_tick_label(value: object) -> str:  # pragma: no cover
    if isinstance(value, (pd.Timestamp, datetime)):
        return pd.Timestamp(value).strftime("%d/%m")
    text = str(value or "").strip()
    if not text:
        return ""
    with contextlib.suppress(Exception):
        ts = pd.Timestamp(text)
        if pd.notna(ts):
            return ts.strftime("%d/%m")
    return _clip(text.replace("<br>", " "), 24)


def _plotly_title_text(value: object, default: str = "") -> str:  # pragma: no cover
    if value is None:
        return default
    text = getattr(value, "text", None)
    if text not in (None, ""):
        return str(text)
    raw = str(value or "").strip()
    return raw if raw and raw != "None" else default


def _nice_ticks(
    min_value: float, max_value: float, *, target: int = 5
) -> list[float]:  # pragma: no cover
    if not np.isfinite(min_value) or not np.isfinite(max_value):
        return [0.0, 1.0]
    if math.isclose(min_value, max_value):
        anchor = 0.0 if math.isclose(max_value, 0.0) else max_value
        return [anchor, anchor + 1.0]
    span = max_value - min_value
    raw_step = span / max(int(target), 2)
    magnitude = 10 ** math.floor(math.log10(abs(raw_step))) if raw_step else 1.0
    normalized = raw_step / magnitude if magnitude else raw_step
    if normalized <= 1:
        step = 1 * magnitude
    elif normalized <= 2:
        step = 2 * magnitude
    elif normalized <= 5:
        step = 5 * magnitude
    else:
        step = 10 * magnitude
    start = math.floor(min_value / step) * step
    end = math.ceil(max_value / step) * step
    ticks: list[float] = []
    cursor = start
    while cursor <= end + step * 0.5:
        ticks.append(float(cursor))
        cursor += step
        if len(ticks) > 10:
            break
    return ticks or [min_value, max_value]


def _format_tick_value(value: float) -> str:  # pragma: no cover
    if not np.isfinite(value):
        return ""
    if abs(value) >= 1000:
        return _fmt_count_or_nd(value)
    if math.isclose(value, round(value), abs_tol=1e-9):
        return f"{int(round(value))}"
    return _fmt_num_or_nd(value, decimals=1)


def _plotly_colorscale_color(
    colorscale: object, ratio: float
) -> tuple[int, int, int]:  # pragma: no cover
    with contextlib.suppress(Exception):
        from plotly.colors import sample_colorscale

        sampled = sample_colorscale(colorscale, [min(max(float(ratio), 0.0), 1.0)])[0]
        return _pillow_color(sampled)
    default_scale = [
        [0.0, "#" + BBVA_COLORS["red"]],
        [0.5, "#" + BBVA_COLORS["yellow"]],
        [1.0, "#" + BBVA_COLORS["green"]],
    ]
    with contextlib.suppress(Exception):
        from plotly.colors import sample_colorscale

        sampled = sample_colorscale(default_scale, [min(max(float(ratio), 0.0), 1.0)])[0]
        return _pillow_color(sampled)
    return _pillow_color("#" + BBVA_COLORS["sky"])


def _plotly_bar_colors(
    trace: object, count: int, layout: object
) -> list[tuple[int, int, int]]:  # pragma: no cover
    marker = getattr(trace, "marker", None)
    color = getattr(marker, "color", None)
    if isinstance(color, (list, tuple, np.ndarray, pd.Series)):
        values = list(color)
        if values and all(
            isinstance(item, (int, float, np.integer, np.floating)) for item in values
        ):
            numeric = pd.to_numeric(pd.Series(values), errors="coerce")
            vmin = float(numeric.min()) if not numeric.dropna().empty else 0.0
            vmax = float(numeric.max()) if not numeric.dropna().empty else 1.0
            span = max(vmax - vmin, 1e-9)
            colorscale = getattr(marker, "colorscale", None) or getattr(
                getattr(layout, "coloraxis", None), "colorscale", None
            )
            return [
                (
                    _plotly_colorscale_color(colorscale, (float(value) - vmin) / span)
                    if np.isfinite(float(value))
                    else _pillow_color("#" + BBVA_COLORS["line"])
                )
                for value in numeric.fillna(vmin).tolist()
            ]
        return [_pillow_color(item) for item in values[:count]] + [
            _pillow_color("#" + BBVA_COLORS["sky"])
        ] * max(count - len(values), 0)
    return [_pillow_color(color or "#" + BBVA_COLORS["sky"])] * count


def _pillow_render_heatmap(
    fig: go.Figure, width: int, height: int
) -> Optional[bytes]:  # pragma: no cover
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return None

    traces = [
        trace
        for trace in fig.data
        if str(getattr(trace, "type", "") or "").strip().lower() == "heatmap"
    ]
    if not traces:
        return None
    trace = traces[0]
    z = np.array(getattr(trace, "z", []), dtype=float)
    if z.size == 0:
        return None
    x_values = [str(value).replace("<br>", " ") for value in list(getattr(trace, "x", []))]
    y_values = [str(value).replace("<br>", " ") for value in list(getattr(trace, "y", []))]
    rows = max(len(y_values), z.shape[0])
    cols = max(len(x_values), z.shape[1] if z.ndim > 1 else 0)
    if rows <= 0 or cols <= 0:
        return None

    image = Image.new("RGB", (width, height), _pillow_color("#FFFFFF"))
    draw = ImageDraw.Draw(image)
    title_font = _pillow_font(max(height // 36, 18), bold=True)
    axis_font = _pillow_font(max(height // 44, 14))
    tick_font = _pillow_font(max(height // 46, 12))

    plot_left = int(width * 0.30)
    plot_right = int(width * 0.88)
    plot_top = int(height * 0.10)
    plot_bottom = int(height * 0.86)
    cell_gap = max(int(min(width, height) * 0.003), 2)
    cell_w = max((plot_right - plot_left - cell_gap * (cols - 1)) // max(cols, 1), 1)
    cell_h = max((plot_bottom - plot_top - cell_gap * (rows - 1)) // max(rows, 1), 1)

    valid = z[np.isfinite(z)]
    zmin = float(valid.min()) if valid.size else 0.0
    zmax = float(valid.max()) if valid.size else 1.0
    span = max(zmax - zmin, 1e-9)
    colorscale = getattr(trace, "colorscale", None) or getattr(
        getattr(fig.layout, "coloraxis", None), "colorscale", None
    )

    for row in range(rows):
        for col in range(cols):
            value = float(z[row, col]) if row < z.shape[0] and col < z.shape[1] else float("nan")
            ratio = 0.0 if not np.isfinite(value) else (value - zmin) / span
            fill = _plotly_colorscale_color(colorscale, ratio)
            x0 = plot_left + col * (cell_w + cell_gap)
            y0 = plot_top + row * (cell_h + cell_gap)
            draw.rounded_rectangle(
                [(x0, y0), (x0 + cell_w, y0 + cell_h)],
                radius=max(min(cell_w, cell_h) // 8, 2),
                fill=fill,
                outline=_pillow_color("#FFFFFF"),
                width=1,
            )

    for row, label in enumerate(y_values[:rows]):
        wrapped = _wrap_label(label, width=18, max_lines=2, joiner="\n")
        tw, th = _pillow_text_size(draw, wrapped, tick_font)
        y0 = plot_top + row * (cell_h + cell_gap) + max((cell_h - th) // 2, 0)
        draw.multiline_text(
            (plot_left - tw - 16, y0),
            wrapped,
            fill=_pillow_color("#" + BBVA_COLORS["ink"]),
            font=tick_font,
            spacing=3,
            align="right",
        )

    max_x_ticks = min(cols, 8)
    step = max(int(math.ceil(cols / max(max_x_ticks, 1))), 1)
    for col, label in enumerate(x_values[:cols]):
        if col % step != 0 and col != cols - 1:
            continue
        short = _wrap_label(_plotly_tick_label(label), width=10, max_lines=2, joiner="\n")
        tw, _ = _pillow_text_size(draw, short, tick_font)
        x0 = plot_left + col * (cell_w + cell_gap) + max((cell_w - tw) // 2, 0)
        draw.multiline_text(
            (x0, plot_bottom + 12),
            short,
            fill=_pillow_color("#" + BBVA_COLORS["ink"]),
            font=tick_font,
            spacing=2,
            align="center",
        )

    colorbar_left = int(width * 0.92)
    colorbar_top = plot_top
    colorbar_bottom = plot_bottom
    for idx in range(colorbar_top, colorbar_bottom):
        ratio = 1.0 - ((idx - colorbar_top) / max(colorbar_bottom - colorbar_top, 1))
        draw.line(
            [(colorbar_left, idx), (colorbar_left + 16, idx)],
            fill=_plotly_colorscale_color(colorscale, ratio),
            width=1,
        )
    draw.text(
        (colorbar_left - 4, colorbar_top - 24),
        _plotly_title_text(getattr(getattr(trace, "colorbar", None), "title", None), "NPS"),
        fill=_pillow_color("#" + BBVA_COLORS["ink"]),
        font=title_font,
    )
    for tick in _nice_ticks(zmin, zmax, target=4):
        ratio = (tick - zmin) / span if span else 0.0
        y_tick = colorbar_bottom - int((colorbar_bottom - colorbar_top) * ratio)
        draw.line(
            [(colorbar_left + 18, y_tick), (colorbar_left + 24, y_tick)],
            fill=_pillow_color("#" + BBVA_COLORS["ink"]),
            width=1,
        )
        draw.text(
            (colorbar_left + 28, y_tick - 8),
            _format_tick_value(tick),
            fill=_pillow_color("#" + BBVA_COLORS["muted"]),
            font=axis_font,
        )

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def _pillow_render_xy(
    fig: go.Figure, width: int, height: int
) -> Optional[bytes]:  # pragma: no cover
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return None

    if not fig.data:
        return None

    trace_types = {
        str(getattr(trace, "type", "") or "").strip().lower()
        for trace in fig.data
        if trace is not None
    }
    horizontal = bool(
        trace_types == {"bar"}
        and all(str(getattr(trace, "orientation", "") or "").lower() == "h" for trace in fig.data)
    )

    image = Image.new("RGB", (width, height), _pillow_color("#FFFFFF"))
    draw = ImageDraw.Draw(image)
    axis_font = _pillow_font(max(height // 44, 14))
    tick_font = _pillow_font(max(height // 48, 12))
    legend_font = _pillow_font(max(height // 50, 11))

    legend_items = [
        trace
        for trace in fig.data
        if bool(getattr(trace, "showlegend", True))
        and str(getattr(trace, "name", "") or "").strip()
    ]
    legend_y = 14
    legend_x = 18
    for trace in legend_items[:6]:
        trace_type = str(getattr(trace, "type", "") or "").strip().lower()
        if trace_type == "scatter":
            color = getattr(getattr(trace, "line", None), "color", None) or getattr(
                getattr(trace, "marker", None), "color", None
            )
        else:
            color = getattr(getattr(trace, "marker", None), "color", None)
            if isinstance(color, (list, tuple, np.ndarray, pd.Series)):
                color = list(color)[0] if list(color) else "#" + BBVA_COLORS["sky"]
        draw.rounded_rectangle(
            [(legend_x, legend_y + 4), (legend_x + 16, legend_y + 16)],
            radius=4,
            fill=_pillow_color(color or "#" + BBVA_COLORS["sky"]),
        )
        label = _clip(getattr(trace, "name", "") or "", 22)
        draw.text(
            (legend_x + 22, legend_y),
            label,
            fill=_pillow_color("#" + BBVA_COLORS["ink"]),
            font=legend_font,
        )
        label_w, _ = _pillow_text_size(draw, label, legend_font)
        legend_x += label_w + 52

    plot_left = int(width * (0.34 if horizontal else 0.10))
    plot_right = int(width * 0.92)
    plot_top = int(height * 0.16)
    plot_bottom = int(height * 0.84)
    right_axis_present = any(str(getattr(trace, "yaxis", "y")) == "y2" for trace in fig.data)
    if right_axis_present:
        plot_right = int(width * 0.86)

    draw.rectangle(
        [(plot_left, plot_top), (plot_right, plot_bottom)],
        outline=_pillow_color("#" + BBVA_COLORS["line"]),
        width=1,
    )

    if horizontal:
        trace = next(
            (item for item in fig.data if str(getattr(item, "type", "")).lower() == "bar"), None
        )
        if trace is None:
            return None
        categories = [str(value).replace("<br>", " ") for value in list(getattr(trace, "y", []))]
        values = pd.to_numeric(pd.Series(list(getattr(trace, "x", []))), errors="coerce").fillna(
            0.0
        )
        if values.empty:
            return None
        min_value = min(float(values.min()), 0.0)
        max_value = max(float(values.max()), 0.0)
        span = max(max_value - min_value, 1e-9)
        zero_x = plot_left + int(((0.0 - min_value) / span) * (plot_right - plot_left))
        bar_colors = _plotly_bar_colors(trace, len(values), fig.layout)
        ticks = _nice_ticks(min_value, max_value, target=5)

        for tick in ticks:
            x = plot_left + int(((tick - min_value) / span) * (plot_right - plot_left))
            draw.line(
                [(x, plot_top), (x, plot_bottom)],
                fill=_pillow_color("#" + BBVA_COLORS["line"]),
                width=1,
            )
            label = _format_tick_value(tick)
            tw, _ = _pillow_text_size(draw, label, axis_font)
            draw.text(
                (x - tw // 2, plot_bottom + 12),
                label,
                fill=_pillow_color("#" + BBVA_COLORS["muted"]),
                font=axis_font,
            )

        row_h = max((plot_bottom - plot_top) // max(len(categories), 1), 1)
        for idx, (category, value) in enumerate(zip(categories, values.tolist())):
            center_y = plot_top + idx * row_h + row_h // 2
            label = _wrap_label(category, width=18, max_lines=2, joiner="\n")
            tw, th = _pillow_text_size(draw, label, tick_font)
            draw.multiline_text(
                (plot_left - tw - 16, center_y - th // 2),
                label,
                fill=_pillow_color("#" + BBVA_COLORS["ink"]),
                font=tick_font,
                spacing=2,
                align="right",
            )
            end_x = plot_left + int(((float(value) - min_value) / span) * (plot_right - plot_left))
            x0, x1 = sorted((zero_x, end_x))
            draw.rounded_rectangle(
                [(x0, center_y - max(row_h // 4, 6)), (x1, center_y + max(row_h // 4, 6))],
                radius=8,
                fill=(
                    bar_colors[idx]
                    if idx < len(bar_colors)
                    else _pillow_color("#" + BBVA_COLORS["sky"])
                ),
            )
            value_label = _format_tick_value(float(value))
            draw.text(
                (x1 + 8, center_y - 8),
                value_label,
                fill=_pillow_color("#" + BBVA_COLORS["ink"]),
                font=axis_font,
            )
        draw.line(
            [(zero_x, plot_top), (zero_x, plot_bottom)],
            fill=_pillow_color("#" + BBVA_COLORS["muted"]),
            width=2,
        )
    else:
        category_values: list[object] = []
        for trace in fig.data:
            if str(getattr(trace, "type", "") or "").strip().lower() not in {"bar", "scatter"}:
                continue
            for value in list(getattr(trace, "x", [])):
                if value not in category_values:
                    category_values.append(value)
        if not category_values:
            return None
        n = len(category_values)
        x_positions = {
            value: plot_left + int(idx * (plot_right - plot_left) / max(n - 1, 1))
            for idx, value in enumerate(category_values)
        }

        left_values: list[float] = []
        right_values: list[float] = []
        stacked_primary: dict[object, float] = {value: 0.0 for value in category_values}
        for trace in fig.data:
            trace_type = str(getattr(trace, "type", "") or "").strip().lower()
            series = pd.to_numeric(
                pd.Series(list(getattr(trace, "y", []))), errors="coerce"
            ).fillna(0.0)
            axis_key = "right" if str(getattr(trace, "yaxis", "y")) == "y2" else "left"
            if axis_key == "right":
                right_values.extend(series.tolist())
            else:
                if (
                    trace_type == "bar"
                    and str(getattr(fig.layout, "barmode", "") or "").lower() == "stack"
                ):
                    x_trace = list(getattr(trace, "x", []))
                    for x_value, y_value in zip(x_trace, series.tolist()):
                        stacked_primary[x_value] = stacked_primary.get(x_value, 0.0) + float(
                            y_value
                        )
                else:
                    left_values.extend(series.tolist())
        if stacked_primary:
            left_values.extend(stacked_primary.values())
        left_min = min(0.0, float(min(left_values)) if left_values else 0.0)
        left_max = max(float(max(left_values)) if left_values else 1.0, 0.0)
        right_min = min(0.0, float(min(right_values)) if right_values else 0.0)
        right_max = max(float(max(right_values)) if right_values else 1.0, 0.0)
        left_span = max(left_max - left_min, 1e-9)
        right_span = max(right_max - right_min, 1e-9)

        def left_y(value: float) -> int:
            return plot_bottom - int(((value - left_min) / left_span) * (plot_bottom - plot_top))

        def right_y(value: float) -> int:
            return plot_bottom - int(((value - right_min) / right_span) * (plot_bottom - plot_top))

        for tick in _nice_ticks(left_min, left_max, target=5):
            y = left_y(tick)
            draw.line(
                [(plot_left, y), (plot_right, y)],
                fill=_pillow_color("#" + BBVA_COLORS["line"]),
                width=1,
            )
            label = _format_tick_value(tick)
            tw, _ = _pillow_text_size(draw, label, axis_font)
            draw.text(
                (plot_left - tw - 10, y - 8),
                label,
                fill=_pillow_color("#" + BBVA_COLORS["muted"]),
                font=axis_font,
            )
        if right_axis_present:
            for tick in _nice_ticks(right_min, right_max, target=5):
                y = right_y(tick)
                label = _format_tick_value(tick)
                draw.text(
                    (plot_right + 10, y - 8),
                    label,
                    fill=_pillow_color("#" + BBVA_COLORS["muted"]),
                    font=axis_font,
                )

        if right_axis_present:
            draw.text(
                (plot_right + 8, plot_top - 24),
                _plotly_title_text(
                    getattr(getattr(fig.layout, "yaxis2", None), "title", None), "Incidencias"
                ),
                fill=_pillow_color("#" + BBVA_COLORS["ink"]),
                font=axis_font,
            )
        draw.text(
            (plot_left, plot_top - 24),
            _plotly_title_text(getattr(getattr(fig.layout, "yaxis", None), "title", None), ""),
            fill=_pillow_color("#" + BBVA_COLORS["ink"]),
            font=axis_font,
        )

        bar_width = max(int((plot_right - plot_left) / max(n, 1) * 0.62), 6)
        stacked_cache: dict[tuple[str, object], float] = {}
        for trace in fig.data:
            trace_type = str(getattr(trace, "type", "") or "").strip().lower()
            x_trace = list(getattr(trace, "x", []))
            y_trace = pd.to_numeric(
                pd.Series(list(getattr(trace, "y", []))), errors="coerce"
            ).fillna(0.0)
            axis_key = "right" if str(getattr(trace, "yaxis", "y")) == "y2" else "left"
            if trace_type == "bar":
                bar_colors = _plotly_bar_colors(trace, len(y_trace), fig.layout)
                stacked = (
                    str(getattr(fig.layout, "barmode", "") or "").lower() == "stack"
                    and axis_key == "left"
                )
                for idx, (x_value, y_value) in enumerate(zip(x_trace, y_trace.tolist())):
                    x = x_positions.get(x_value)
                    if x is None:
                        continue
                    if stacked:
                        baseline = stacked_cache.get((axis_key, x_value), 0.0)
                        y0 = left_y(baseline + float(y_value))
                        y1 = left_y(baseline)
                        stacked_cache[(axis_key, x_value)] = baseline + float(y_value)
                    else:
                        y0 = (
                            right_y(float(y_value))
                            if axis_key == "right"
                            else left_y(float(y_value))
                        )
                        y1 = right_y(0.0) if axis_key == "right" else left_y(0.0)
                    draw.rounded_rectangle(
                        [(x - bar_width // 2, min(y0, y1)), (x + bar_width // 2, max(y0, y1))],
                        radius=6,
                        fill=(
                            bar_colors[idx]
                            if idx < len(bar_colors)
                            else _pillow_color("#" + BBVA_COLORS["sky"])
                        ),
                    )
            elif trace_type == "scatter":
                points: list[tuple[int, int]] = []
                line_color = getattr(getattr(trace, "line", None), "color", None) or getattr(
                    getattr(trace, "marker", None), "color", None
                )
                color = _pillow_color(line_color or "#" + BBVA_COLORS["blue"])
                for x_value, y_value in zip(x_trace, y_trace.tolist()):
                    x = x_positions.get(x_value)
                    if x is None or not np.isfinite(float(y_value)):
                        continue
                    y = right_y(float(y_value)) if axis_key == "right" else left_y(float(y_value))
                    points.append((x, y))
                if len(points) >= 2:
                    draw.line(points, fill=color, width=max(width // 500, 3))
                if "markers" in str(getattr(trace, "mode", "") or "").lower():
                    radius = max(width // 260, 4)
                    for x, y in points:
                        draw.ellipse(
                            [(x - radius, y - radius), (x + radius, y + radius)],
                            fill=color,
                            outline=_pillow_color("#FFFFFF"),
                            width=1,
                        )

        max_x_ticks = min(n, 10)
        step = max(int(math.ceil(n / max(max_x_ticks, 1))), 1)
        for idx, value in enumerate(category_values):
            if idx % step != 0 and idx != n - 1:
                continue
            label = _wrap_label(_plotly_tick_label(value), width=10, max_lines=2, joiner="\n")
            tw, _ = _pillow_text_size(draw, label, tick_font)
            x = x_positions[value]
            draw.multiline_text(
                (x - tw // 2, plot_bottom + 10),
                label,
                fill=_pillow_color("#" + BBVA_COLORS["muted"]),
                font=tick_font,
                spacing=2,
                align="center",
            )

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def _pillow_chart_png(
    fig: go.Figure, *, width: int, height: int
) -> Optional[bytes]:  # pragma: no cover
    try:
        themed = _apply_ppt_figure_theme(go.Figure(fig))
    except Exception:
        themed = fig
    if any(
        str(getattr(trace, "type", "") or "").strip().lower() == "heatmap" for trace in themed.data
    ):
        return _pillow_render_heatmap(themed, width, height)
    return _pillow_render_xy(themed, width, height)


def _kaleido_png(
    fig: go.Figure,
    *,
    width: int = 1600,
    height: int = 900,
    panel_width_in: float | None = None,
    panel_height_in: float | None = None,
) -> Optional[bytes]:
    try:
        _patch_kaleido_executable_for_space_paths()
        themed = _apply_ppt_figure_theme(
            go.Figure(fig),
            panel_width_in=panel_width_in,
            panel_height_in=panel_height_in,
        )
        attempts = [
            (themed, max(int(width), 960), max(int(height), 540)),
            (
                themed,
                min(max(int(width), 960), 1800),
                min(max(int(height), 540), 1080),
            ),
            (
                _apply_ppt_figure_theme(
                    go.Figure(fig),
                    panel_width_in=panel_width_in,
                    panel_height_in=panel_height_in,
                ),
                1280,
                720,
            ),
        ]
        for candidate, attempt_width, attempt_height in attempts:
            with contextlib.suppress(Exception):
                return pio.to_image(
                    candidate,
                    format="png",
                    width=attempt_width,
                    height=attempt_height,
                    scale=1,
                )
    except Exception:
        pass
    return _pillow_chart_png(fig, width=max(int(width), 960), height=max(int(height), 540))


def _layout_fallback_index(prs: Presentation, kind: str) -> int:
    count = len(prs.slide_layouts)
    if count <= 0:
        return 0
    if kind == "cover":
        return min(0, count - 1)
    if kind == "section":
        return min(2, count - 1)
    return min(6, count - 1)


def _new_slide(prs: Presentation, *, kind: str = "content") -> object:
    preferred = (
        [PPT_THEME.cover_layout]
        if kind == "cover"
        else ([PPT_THEME.section_layout] if kind == "section" else [PPT_THEME.content_layout])
    )
    layout = resolve_layout(
        prs,
        preferred,
        fallback_index=_layout_fallback_index(prs, kind),
    )
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            continue
    return slide


def _add_bg(slide: object, color: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()


def _add_header(
    slide: object, *, title: str, subtitle: str, dark: bool = False, right_note: str = ""
) -> None:
    title_color = BBVA_COLORS["white"] if dark else BBVA_COLORS["ink"]
    sub_color = "A9B7D2" if dark else BBVA_COLORS["muted"]

    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.28), Inches(9.8), Inches(0.85))
    tf = box.text_frame
    _configure_text_frame(tf)
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = BBVA_FONT_DISPLAY
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = _rgb(title_color)

    sb = slide.shapes.add_textbox(Inches(0.65), Inches(0.95), Inches(11.5), Inches(0.42))
    stf = sb.text_frame
    _configure_text_frame(stf)
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.name = BBVA_FONT_BODY
    sr.font.size = Pt(12.5)
    sr.font.color.rgb = _rgb(sub_color)

    if right_note.strip():
        rb = slide.shapes.add_textbox(Inches(10.4), Inches(0.35), Inches(2.2), Inches(0.30))
        rtf = rb.text_frame
        _configure_text_frame(rtf)
        rtf.clear()
        rp = rtf.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        rr = rp.add_run()
        rr.text = right_note
        rr.font.name = BBVA_FONT_BODY
        rr.font.size = Pt(10)
        rr.font.color.rgb = _rgb(sub_color)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.28), Inches(12.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(BBVA_COLORS["sky"] if dark else BBVA_COLORS["line"])
    line.line.fill.background()


def _panel(
    slide: object,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str = "",
    subtitle: str = "",
    fill: str = "",
    border: str = "",
    title_size: float = 15,
) -> object:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(fill or BBVA_COLORS["white"])
    box.line.color.rgb = _rgb(border or BBVA_COLORS["line"])
    if not title and not subtitle:
        return box

    tf = box.text_frame
    _configure_text_frame(tf)
    tf.clear()
    if title:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.name = BBVA_FONT_HEAD
        r0.font.size = Pt(title_size)
        r0.font.bold = True
        r0.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    if subtitle:
        p1 = tf.add_paragraph()
        p1.space_before = Pt(4)
        r1 = p1.add_run()
        r1.text = subtitle
        r1.font.name = BBVA_FONT_BODY
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = _rgb(BBVA_COLORS["muted"])
    return box


def _figure_in_panel(
    slide: object,
    *,
    figure: Optional[go.Figure],
    left: float,
    top: float,
    width: float,
    height: float,
    empty_note: str,
    target_ppi: Optional[int] = None,
) -> None:
    img = None
    if figure is not None:
        if target_ppi is not None and int(target_ppi) > 0:
            width_px = max(int(width * int(target_ppi)), 640)
            height_px = max(int(height * int(target_ppi)), 360)
            scale = 1.0
        else:
            base_ppi = 180
            width_px = max(int(width * base_ppi), 1400)
            height_px = max(int(height * base_ppi), 480)
            scale = max(1800 / max(width_px, 1), 900 / max(height_px, 1), 1.0)
        img = _kaleido_png(
            figure,
            width=int(width_px * scale),
            height=int(height_px * scale),
            panel_width_in=width,
            panel_height_in=height,
        )
    if img is not None:
        slide.shapes.add_picture(
            BytesIO(img),
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height),
        )
        return

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    _configure_text_frame(tf)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = empty_note
    r.font.name = BBVA_FONT_BODY
    r.font.size = Pt(12)
    r.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _add_stat_card(
    slide: object,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    value: str,
    accent: str,
    hint: str = "",
) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    box.line.color.rgb = _rgb(BBVA_COLORS["line"])

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(accent)
    accent_bar.line.fill.background()

    tf = box.text_frame
    _configure_text_frame(tf)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = label.upper()
    r0.font.name = BBVA_FONT_MEDIUM
    r0.font.size = Pt(10)
    r0.font.bold = True
    r0.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    p1 = tf.add_paragraph()
    p1.space_before = Pt(10)
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = BBVA_FONT_DISPLAY
    value_len = len(str(value or ""))
    r1.font.size = Pt(24 if value_len <= 14 else (18 if value_len <= 28 else 12))
    r1.font.bold = True
    r1.font.color.rgb = _rgb(BBVA_COLORS["ink"])

    if hint:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = hint
        r2.font.name = BBVA_FONT_BODY
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _add_markdown_runs(
    paragraph: object,
    *,
    text: object,
    prefix: str = "",
    max_chars: int = 170,
    font_size_pt: float = 11.0,
    color: str = "",
) -> None:
    remaining = max(int(max_chars), 1)
    if prefix:
        run = paragraph.add_run()
        run.text = prefix
        run.font.name = BBVA_FONT_BODY
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = _rgb(color or BBVA_COLORS["muted"])
    visible_len = 0
    truncated = False
    for segment in parse_markdown_strong(text):
        if remaining <= 0:
            truncated = True
            break
        segment_text = segment.text
        if len(segment_text) > remaining:
            segment_text = segment_text[: max(remaining - 1, 0)].rstrip() + "…"
            truncated = True
        run = paragraph.add_run()
        run.text = segment_text
        run.font.name = BBVA_FONT_BODY if not segment.bold else BBVA_FONT_MEDIUM
        run.font.size = Pt(font_size_pt)
        run.font.bold = bool(segment.bold)
        run.font.color.rgb = _rgb(color or BBVA_COLORS["muted"])
        visible_len += len(segment_text)
        remaining -= len(segment_text)
        if truncated:
            break
    if visible_len == 0 and not prefix:
        run = paragraph.add_run()
        run.text = ""
        run.font.name = BBVA_FONT_BODY
        run.font.size = Pt(font_size_pt)


def _add_bullet_lines(
    slide: object,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    accent: str = "",
    body_font_size_pt: float = 11.0,
) -> None:
    if not str(title or "").strip():
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
        box.line.color.rgb = _rgb(accent or BBVA_COLORS["line"])
        tf = box.text_frame
        _configure_text_frame(tf)
        tf.clear()
        for idx, line in enumerate(lines[:6]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(4 if idx else 0)
            p.level = 0
            _add_markdown_runs(
                p,
                text=line,
                prefix="• ",
                max_chars=145 if width <= 4.0 else 170,
                font_size_pt=body_font_size_pt,
                color=BBVA_COLORS["muted"],
            )
        return

    panel = _panel(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        title=title,
        fill=BBVA_COLORS["white"],
        border=accent or BBVA_COLORS["line"],
    )
    tf = panel.text_frame
    _configure_text_frame(tf)
    for line in lines[:6]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        p.level = 0
        _add_markdown_runs(
            p,
            text=line,
            prefix="• ",
            max_chars=145 if width <= 4.0 else 170,
            font_size_pt=body_font_size_pt,
            color=BBVA_COLORS["muted"],
        )


def _add_compact_table(
    slide: object,
    *,
    left: float,
    top: float,
    width: float,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    row_height: float = 0.34,
    col_width_ratios: Optional[list[float]] = None,
    clip_lengths: Optional[list[int]] = None,
    font_size_pt: float = 9.6,
    max_rows: int = 6,
) -> None:
    title_pad = 0.62 if str(title or "").strip() else 0.24
    visible_rows = max(min(len(rows), max(int(max_rows), 1)), 1)
    height = title_pad + 0.54 + row_height * visible_rows
    panel = _panel(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        title=title,
        fill=BBVA_COLORS["white"],
    )
    base_top = top + (0.46 if str(title or "").strip() else 0.18)
    if col_width_ratios and len(col_width_ratios) == len(headers):
        ratio_sum = sum(col_width_ratios) or 1.0
        column_widths = [((width - 0.16) * ratio / ratio_sum) for ratio in col_width_ratios]
    else:
        column_widths = [(width - 0.16) / max(len(headers), 1)] * len(headers)
    x_positions: list[float] = []
    cursor = left
    for col_width in column_widths:
        x_positions.append(cursor)
        cursor += col_width

    for idx, header in enumerate(headers):
        tb = slide.shapes.add_textbox(
            Inches(x_positions[idx] + 0.05),
            Inches(base_top),
            Inches(column_widths[idx] - 0.08),
            Inches(0.20),
        )
        tf = tb.text_frame
        _configure_text_frame(tf)
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.name = BBVA_FONT_MEDIUM
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.color.rgb = _rgb(BBVA_COLORS["blue"])

    for row_idx, row in enumerate(rows[: max(int(max_rows), 1)], start=1):
        current_top = base_top + 0.16 + row_height * row_idx
        for col_idx, value in enumerate(row[: len(headers)]):
            tb = slide.shapes.add_textbox(
                Inches(x_positions[col_idx] + 0.05),
                Inches(current_top),
                Inches(column_widths[col_idx] - 0.08),
                Inches(0.28),
            )
            tf = tb.text_frame
            _configure_text_frame(tf)
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            clip_len = (
                clip_lengths[col_idx]
                if clip_lengths is not None and col_idx < len(clip_lengths)
                else 48
            )
            r.text = _clip(value, clip_len)
            r.font.name = BBVA_FONT_BODY
            r.font.size = Pt(font_size_pt)
            r.font.color.rgb = _rgb(BBVA_COLORS["muted"] if col_idx else BBVA_COLORS["ink"])

    del panel


def _add_chart_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    figure: Optional[go.Figure],
    rationale_title: str,
    rationale_lines: Iterable[str],
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(slide, title=title, subtitle=subtitle)

    chart_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.66), Inches(1.45), Inches(8.2), Inches(5.4)
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    chart_box.line.color.rgb = _rgb(BBVA_COLORS["line"])

    img = _kaleido_png(figure) if figure is not None else None
    if img is not None:
        slide.shapes.add_picture(
            BytesIO(img), Inches(0.82), Inches(1.60), width=Inches(7.84), height=Inches(5.08)
        )
    else:
        tf = chart_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "No hay información suficiente para este gráfico en el periodo seleccionado."
        r.font.name = BBVA_FONT_BODY
        r.font.size = Pt(13)
        r.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    side = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.45), Inches(3.75), Inches(5.4)
    )
    side.fill.solid()
    side.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    side.line.color.rgb = _rgb(BBVA_COLORS["line"])

    stf = side.text_frame
    stf.clear()

    hp = stf.paragraphs[0]
    hr = hp.add_run()
    hr.text = rationale_title
    hr.font.name = BBVA_FONT_HEAD
    hr.font.bold = True
    hr.font.size = Pt(18)
    hr.font.color.rgb = _rgb(BBVA_COLORS["ink"])

    for line in list(rationale_lines)[:8]:
        p = stf.add_paragraph()
        p.level = 0
        p.space_after = Pt(7)
        p.text = f"• {_clip(line, 175)}"
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _add_period_summary_slide(
    prs: Presentation,
    *,
    service_origin: str,
    service_origin_n1: str,
    service_origin_n2: str,
    period_start: date,
    period_end: date,
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])

    month_txt = _month_label_es(period_end)
    p_label = f"{_safe_date(period_start)} -> {_safe_date(period_end)}"
    cards = [
        ("SERVICE ORIGEN", _clip(service_origin or "N/D", 38), "Ámbito analizado"),
        ("NIVEL N1", _clip(service_origin_n1 or "N/D", 38), "Segmentación principal"),
        ("NIVEL N2", _clip(service_origin_n2 or "N/D", 38), "Segmentación secundaria"),
        ("MES EN CURSO", _clip(month_txt.title(), 38), f"Ventana: {p_label}"),
    ]

    left0 = 0.65
    top = 0.85
    card_w = 3.18
    gap = 0.35
    card_h = 3.05

    for i, (label, value, caption) in enumerate(cards):
        left = Inches(left0 + i * (card_w + gap))
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(BBVA_COLORS["bg_light"])
        card.line.color.rgb = _rgb(BBVA_COLORS["line"])

        tf = card.text_frame
        tf.clear()

        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = label
        r1.font.name = BBVA_FONT_MEDIUM
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = _rgb(BBVA_COLORS["ink"])

        p2 = tf.add_paragraph()
        p2.space_before = Pt(18)
        r2 = p2.add_run()
        r2.text = value
        r2.font.name = BBVA_FONT_HEAD
        r2.font.bold = True
        r2.font.size = Pt(30)
        r2.font.color.rgb = _rgb(BBVA_COLORS["ink"])

        p3 = tf.add_paragraph()
        p3.space_before = Pt(8)
        r3 = p3.add_run()
        r3.text = caption
        r3.font.name = BBVA_FONT_BODY
        r3.font.size = Pt(12)
        r3.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _plain_md(text: object) -> str:
    s = str(text or "").strip()
    if not s:
        return ""
    s = re.sub(r"`([^`]*)`", r"\1", s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"^#+\s*", "", s)
    return " ".join(s.split())


def _parse_story_sections(story_md: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    title = ""
    bullets: list[str] = []
    for raw in str(story_md or "").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("## "):
            if title:
                sections.append((title, bullets))
            title = _plain_md(line[3:])
            bullets = []
            continue
        if line.startswith("# "):
            continue
        if line.startswith("- "):
            bullets.append(_plain_md(line[2:]))
            continue
        if title:
            bullets.append(_plain_md(line))
    if title:
        sections.append((title, bullets))
    return sections


def _cover_summary_lines(overview: dict[str, object], story_md: str) -> list[str]:
    lines: list[str] = []
    pain_point = str(overview.get("pain_point", "") or "").strip()
    strength_point = str(overview.get("strength_point", "") or "").strip()
    classic_delta = _safe_float(overview.get("classic_delta"), default=float("nan"))
    detractor_delta_pp = _safe_float(overview.get("detractor_delta_pp"), default=float("nan"))

    if pain_point:
        lines.append(f"La mayor fricción del periodo se concentra en {pain_point}.")
    if strength_point and strength_point != pain_point:
        lines.append(f"La mejor señal de experiencia se observa en {strength_point}.")
    if np.isfinite(classic_delta) and abs(classic_delta) >= 0.1:
        direction = "sube" if classic_delta > 0 else "cae"
        lines.append(
            "Del inicio al cierre del periodo, el NPS clásico "
            f"{direction} {_fmt_num_or_nd(abs(classic_delta), decimals=1)} puntos."
        )
    if np.isfinite(detractor_delta_pp) and abs(detractor_delta_pp) >= 0.1:
        direction = "sube" if detractor_delta_pp > 0 else "baja"
        lines.append(
            "El peso detractor "
            f"{direction} {_fmt_num_or_nd(abs(detractor_delta_pp), decimals=1)} puntos porcentuales "
            "en la ventana analizada."
        )

    for _, bullets in _parse_story_sections(story_md):
        for bullet in bullets:
            clean = str(bullet or "").strip()
            if not clean or _is_cover_metric_line(clean):
                continue
            if clean.lower().startswith(
                (
                    "zona de fricción",
                    "zona de friccion",
                    "zona fuerte",
                    "periodo actual",
                    "periodo base",
                )
            ):
                continue
            clean = (
                clean.replace("VoC", "comentarios de cliente")
                .replace("quick wins", "acciones rápidas")
                .replace("owners", "equipos responsables")
                .replace("owner", "equipo responsable")
                .replace("Si mejoramos Palanca=", "La oportunidad más clara está en ")
                .replace("Si mejoramos Subpalanca=", "La oportunidad más clara está en ")
                .replace("el modelo estima un potencial de", "con un impacto potencial de")
            )
            if clean not in lines:
                lines.append(clean)
            if len(lines) >= 4:
                return lines[:4]
    return lines[:4]


def _add_story_card(
    slide: object,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    fill: str = "",
) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(fill or BBVA_COLORS["white"])
    box.line.color.rgb = _rgb(BBVA_COLORS["line"])

    tf = box.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = BBVA_FONT_HEAD
    r0.font.bold = True
    r0.font.size = Pt(18)
    r0.font.color.rgb = _rgb(BBVA_COLORS["ink"])

    max_bullets = 3 if height <= 1.7 else 4
    for bullet in bullets[:max_bullets]:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(0)
        p.level = 0
        r = p.add_run()
        r.text = f"• {_clip(bullet, 150 if width > 5.5 else 120)}"
        r.font.name = BBVA_FONT_BODY
        r.font.size = Pt(12)
        r.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _add_business_story_slide(
    prs: Presentation,
    *,
    story_md: str,
    period_label: str,
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="Informe de negocio",
        subtitle="Lectura ejecutiva del mes seleccionado frente al histórico anterior, lista para comité.",
        right_note="Slide 2",
    )

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(1.52),
        Inches(11.95),
        Inches(0.55),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = _rgb(BBVA_COLORS["ink"])
    banner.line.fill.background()
    btf = banner.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = f"Periodo analizado: {period_label}"
    br.font.name = BBVA_FONT_MEDIUM
    br.font.bold = True
    br.font.size = Pt(13)
    br.font.color.rgb = _rgb(BBVA_COLORS["white"])

    sections = _parse_story_sections(story_md)
    section_map = {title: bullets for title, bullets in sections}
    layout = [
        ("1) Qué está pasando", 0.70, 2.22, 5.8, 1.75, BBVA_COLORS["white"]),
        ("2) Cambio vs base de comparación", 6.63, 2.22, 5.8, 1.75, BBVA_COLORS["white"]),
        ("3) Dónde atacar primero (oportunidades)", 0.70, 4.10, 5.8, 1.95, BBVA_COLORS["white"]),
        ("4) Qué están diciendo (temas de texto)", 6.63, 4.10, 5.8, 1.95, BBVA_COLORS["white"]),
    ]
    for title, left, top, width, height, fill in layout:
        _add_story_card(
            slide,
            left=left,
            top=top,
            width=width,
            height=height,
            title=title,
            bullets=section_map.get(title, ["Sin contenido disponible para esta sección."]),
            fill=fill,
        )

    _add_story_card(
        slide,
        left=0.70,
        top=6.20,
        width=11.73,
        height=0.78,
        title="5) Próximos pasos recomendados",
        bullets=section_map.get(
            "5) Próximos pasos recomendados",
            ["Validar releases, definir hipótesis y aterrizar el plan de acción del mes."],
        ),
        fill=BBVA_COLORS["bg_light"],
    )


def _add_impact_chain_slide(
    prs: Presentation,
    *,
    cards: list[object],
    focus_name: str,
    period_label: str,
    presentation_mode: str = "",
    executive_journey_catalog: Optional[list[dict[str, object]]] = None,
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    if str(presentation_mode or "").strip() == TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS:
        _add_executive_journey_summary_slide(
            prs,
            cards=cards,
            focus_name=focus_name,
            period_label=period_label,
            executive_journey_catalog=executive_journey_catalog,
        )
        return

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    is_broken_journey_mode = (
        str(presentation_mode or "").strip() == TOUCHPOINT_SOURCE_BROKEN_JOURNEYS
    )
    _add_header(
        slide,
        title="Marco causal",
        subtitle=(
            "Cómo se atribuye el impacto: incidencia -> embeddings / clustering -> journey roto -> VoC -> NPS "
            f"· periodo {period_label}"
            if is_broken_journey_mode
            else f"Cómo se atribuye el impacto: incidencia -> touchpoint -> VoC -> NPS · periodo {period_label}"
        ),
    )
    steps = (
        [
            ("1. Incidencia", "Helix aporta el INC y la descripción ampliada del fallo real."),
            (
                "2. Embeddings + keywords",
                "La app agrupa señales semánticas similares sin depender de una tabla manual de journeys.",
            ),
            (
                "3. Journey roto",
                "Cada cluster se convierte en un touchpoint roto defendible con su palanca dominante.",
            ),
            (
                "4. Comentario VoC",
                "Se muestran verbatims reales enlazados con el cluster, no frases genéricas ni heurísticas aisladas.",
            ),
            (
                "5. NPS",
                f"El efecto final se expresa en riesgo de {focus_name}, delta NPS e impacto total.",
            ),
        ]
        if is_broken_journey_mode
        else [
            ("1. Incidencia", "Helix aporta el INC y la descripción ampliada del fallo real."),
            (
                "2. Touchpoint",
                "Se identifica el momento del journey afectado, no solo el sistema técnico.",
            ),
            (
                "3. Palanca / subpalanca",
                "La fricción se traduce al lenguaje NPS con el mismo topic usado en la app.",
            ),
            (
                "4. Comentario VoC",
                "Se muestran verbatims reales enlazados con el caso Helix, no frases genéricas.",
            ),
            (
                "5. NPS",
                f"El efecto final se expresa en riesgo de {focus_name}, delta NPS e impacto total.",
            ),
        ]
    )
    left = 0.80
    top = 1.70
    width = 11.7
    gap = 0.15
    box_h = 0.78
    for idx, (title, body) in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top + idx * (box_h + gap)),
            Inches(width),
            Inches(box_h),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
        box.line.color.rgb = _rgb(BBVA_COLORS["line"])
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = BBVA_FONT_MEDIUM
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = _rgb(BBVA_COLORS["blue"])
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = BBVA_FONT_BODY
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    note = slide.shapes.add_textbox(Inches(0.90), Inches(6.45), Inches(11.4), Inches(0.55))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    nr = np.add_run()
    nr.text = (
        "Solo se presentan temas con link explícito entre Helix y VoC. "
        "Se excluyen etiquetas genéricas sin comentario defendible, como 'Sin comentarios'."
    )
    nr.font.name = BBVA_FONT_BODY
    nr.font.size = Pt(11.5)
    nr.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _add_executive_journey_summary_slide(
    prs: Presentation,
    *,
    cards: list[object],
    focus_name: str,
    period_label: str,
    executive_journey_catalog: Optional[list[dict[str, object]]] = None,
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    focus_label = (
        "detracción" if str(focus_name).strip().lower() == "detractores" else str(focus_name)
    )
    _add_header(
        slide,
        title=f"NPS Lens — Journeys que explican la {focus_label}",
        subtitle=f"Resumen ejecutivo 1 página · periodo {period_label}",
    )

    objective = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(1.42),
        Inches(12.0),
        Inches(1.08),
    )
    objective.fill.solid()
    objective.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    objective.line.color.rgb = _rgb(BBVA_COLORS["line"])
    otf = objective.text_frame
    otf.clear()
    op = otf.paragraphs[0]
    or1 = op.add_run()
    or1.text = "Objetivo"
    or1.font.name = BBVA_FONT_HEAD
    or1.font.bold = True
    or1.font.size = Pt(16)
    or1.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    op2 = otf.add_paragraph()
    or2 = op2.add_run()
    or2.text = (
        "Identificar rutas de degradación de experiencia que conectan señales operativas "
        "(incidencias) con la voz del cliente (NPS) para priorizar causas raíz accionables."
    )
    or2.font.name = BBVA_FONT_BODY
    or2.font.size = Pt(11.5)
    or2.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    table_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(2.72),
        Inches(12.0),
        Inches(2.42),
    )
    table_box.fill.solid()
    table_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    table_box.line.color.rgb = _rgb(BBVA_COLORS["line"])

    catalog_rows = list(executive_journey_catalog or EXECUTIVE_JOURNEY_CATALOG)
    row_y = 2.92
    cols = [0.92, 3.15, 6.45, 10.15]
    headers = ["Journey", "Qué ocurre", "Evidencia esperada", "Impacto en NPS"]
    for idx, header in enumerate(headers):
        tb = slide.shapes.add_textbox(Inches(cols[idx]), Inches(row_y), Inches(2.2), Inches(0.26))
        ttf = tb.text_frame
        ttf.clear()
        p = ttf.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.name = BBVA_FONT_MEDIUM
        r.font.bold = True
        r.font.size = Pt(11.5)
        r.font.color.rgb = _rgb(BBVA_COLORS["blue"])

    for row_idx, journey in enumerate(catalog_rows[:3], start=1):
        card = next(
            (
                item
                for item in cards
                if str(item.get("nps_topic", "") if isinstance(item, dict) else "")
                == str(journey["title"])
            ),
            None,
        )
        current_y = row_y + 0.32 + (row_idx - 1) * 0.58
        values = [
            f"{row_idx}. {journey['title']}",
            str(journey["what_occurs"]),
            str(journey["expected_evidence"]),
            str(journey["impact_label"]),
        ]
        if isinstance(card, dict):
            impact_override = str(card.get("journey_impact_label", "")).strip()
            if impact_override:
                values[3] = impact_override
        widths = [2.0, 3.05, 3.45, 1.55]
        for idx, (value, width) in enumerate(zip(values, widths)):
            tb = slide.shapes.add_textbox(
                Inches(cols[idx]),
                Inches(current_y),
                Inches(width),
                Inches(0.42),
            )
            ttf = tb.text_frame
            ttf.clear()
            p = ttf.paragraphs[0]
            r = p.add_run()
            r.text = _clip(value, 84)
            r.font.name = BBVA_FONT_BODY
            r.font.size = Pt(10.8)
            r.font.color.rgb = _rgb(BBVA_COLORS["ink"] if idx == 0 else BBVA_COLORS["muted"])

    left_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(5.42),
        Inches(5.85),
        Inches(1.25),
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    left_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    ltf = left_box.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = "Valor diferencial de NPS Lens"
    lr.font.name = BBVA_FONT_HEAD
    lr.font.bold = True
    lr.font.size = Pt(16)
    lr.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    for line in [
        "Temas mencionados por clientes -> Journeys de caída de experiencia",
        "Comentarios aislados -> Conexión con incidencias operativas",
        "Insights descriptivos -> Hipótesis causales accionables",
    ]:
        p = ltf.add_paragraph()
        p.text = f"• {_clip(line, 85)}"
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(10.8)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    right_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.85),
        Inches(5.42),
        Inches(5.87),
        Inches(1.25),
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    right_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    rtf = right_box.text_frame
    rtf.clear()
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = "Resultado esperado"
    rr.font.name = BBVA_FONT_HEAD
    rr.font.bold = True
    rr.font.size = Pt(16)
    rr.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    for line in [
        "Dónde se rompe la experiencia",
        "Qué incidencias lo provocan",
        f"Cuántos {focus_name} genera",
        "Qué acciones priorizar",
    ]:
        p = rtf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(10.8)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _chain_list(value: object) -> list[str]:
    if isinstance(value, list):
        return [str(v).strip() for v in value if str(v).strip()]
    if value is None:
        return []
    txt = str(value).strip()
    return [txt] if txt else []


def _chain_header(label: str, shown: int, total: int) -> str:
    if shown < total:
        return f"{label} ({shown} de {total})"
    return f"{label} ({shown})"


def _chain_incident_records(value: object) -> list[dict[str, str]]:
    if not isinstance(value, list):
        return []
    out: list[dict[str, str]] = []
    for entry in value:
        if not isinstance(entry, dict):
            continue
        incident_id = str(entry.get("incident_id", "") or "").strip()
        summary = str(entry.get("summary", "") or "").strip()
        url = str(entry.get("url", "") or "").strip()
        if incident_id or summary:
            out.append(
                {
                    "incident_id": incident_id,
                    "summary": summary,
                    "url": url,
                }
            )
    return out


def _chain_priority_summary(  # pragma: no cover - legacy helper for compatibility mode
    chain_row: pd.Series, *, focus_name: str
) -> list[str]:
    owner = str(chain_row.get("owner_role", "") or "").strip()
    lane = str(chain_row.get("action_lane", "") or "").strip()
    eta_weeks = _safe_float(chain_row.get("eta_weeks", np.nan), default=np.nan)
    responses = _safe_float(chain_row.get("responses", np.nan), default=np.nan)
    incidents = _safe_float(chain_row.get("incidents", np.nan), default=np.nan)
    incident_rate = _safe_float(
        chain_row.get("incident_rate_per_100_responses", np.nan),
        default=np.nan,
    )
    delta_focus = _safe_float(chain_row.get("delta_focus_rate_pp", np.nan), default=np.nan)
    risk = _safe_float(chain_row.get("nps_points_at_risk", np.nan), default=np.nan)
    recoverable = _safe_float(chain_row.get("nps_points_recoverable", np.nan), default=np.nan)
    priority = _safe_float(chain_row.get("priority", np.nan), default=np.nan)
    confidence = _safe_float(chain_row.get("confidence", np.nan), default=np.nan)
    parts_top = [
        f"Prioridad {_fmt_num_or_nd(priority)}",
        f"Confianza {_fmt_num_or_nd(confidence)}",
        f"NPS en riesgo {_fmt_num_or_nd(risk)} pts",
        f"NPS recuperable {_fmt_num_or_nd(recoverable)} pts",
    ]
    parts_bottom = [
        f"Delta % {focus_name} {_fmt_signed_or_nd(delta_focus)} pp",
        f"Incidencias/100 resp. {_fmt_num_or_nd(incident_rate)}",
        f"Incidencias {_fmt_num_or_nd(incidents, decimals=0)}",
        f"Respuestas {_fmt_num_or_nd(responses, decimals=0)}",
    ]
    if lane:
        parts_bottom.append(f"Lane {lane}")
    if owner:
        parts_bottom.append(f"Owner {owner}")
    if np.isfinite(eta_weeks):
        parts_bottom.append(f"ETA (semanas) {_fmt_num_or_nd(eta_weeks, decimals=1)}")
    return [
        " · ".join(parts_top),
        " · ".join(parts_bottom),
    ]


def _add_chain_evidence_slide(
    prs: Presentation,
    *,
    chain_row: pd.Series,
    idx: int,
    focus_name: str,
    period_label: str,
) -> None:  # pragma: no cover - legacy slide kept for backwards compatibility
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])

    topic = _clip(chain_row.get("nps_topic", "Tema sin etiqueta"), 72)
    touchpoint = _clip(chain_row.get("touchpoint", "Touchpoint sin etiqueta"), 42)
    palanca = _clip(chain_row.get("palanca", "n/d"), 30)
    subpalanca = _clip(chain_row.get("subpalanca", "n/d"), 36)
    presentation_mode = str(chain_row.get("presentation_mode", "") or "").strip()
    linked_incidents = int(_safe_int(chain_row.get("linked_incidents", 0), default=0))
    linked_comments = int(_safe_int(chain_row.get("linked_comments", 0), default=0))
    helix_records = _chain_incident_records(chain_row.get("incident_records"))[:5]
    helix_lines = (
        [
            str(rec.get("summary", "")).strip()
            for rec in helix_records
            if str(rec.get("summary", "")).strip()
        ]
        if helix_records
        else _chain_list(chain_row.get("incident_examples"))[:5]
    )
    voc_lines = _chain_list(chain_row.get("comment_examples"))[:2]
    shown_incidents = len(helix_lines)
    shown_comments = len(voc_lines)
    header_title = (
        f"Journey {idx}: {topic}"
        if presentation_mode == TOUCHPOINT_SOURCE_EXECUTIVE_JOURNEYS
        else (
            f"Journey roto {idx}: {topic}"
            if presentation_mode == TOUCHPOINT_SOURCE_BROKEN_JOURNEYS
            else f"Tema prioritario {idx}: {touchpoint}"
        )
    )
    _add_header(
        slide,
        title=header_title,
        subtitle=f"{topic} · periodo {period_label}",
    )

    flow_y = 1.55
    step_w = 2.32
    step_gap = 0.18
    step_titles = [
        f"({shown_incidents}) Incidencias",
        touchpoint,
        f"{palanca} / {subpalanca}",
        f"({shown_comments}) Comentarios VoC",
        "NPS",
    ]
    step_colors = ["sky", "blue", "green", "orange", "red"]
    for s_idx, (label, color_key) in enumerate(zip(step_titles, step_colors)):
        x = 0.70 + s_idx * (step_w + step_gap)
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(flow_y),
            Inches(step_w),
            Inches(0.72),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
        box.line.color.rgb = _rgb(BBVA_COLORS[color_key])
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = BBVA_FONT_MEDIUM
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = _rgb(BBVA_COLORS["ink"])
        if s_idx < len(step_titles) - 1:
            arr = slide.shapes.add_textbox(
                Inches(x + step_w),
                Inches(flow_y + 0.18),
                Inches(step_gap),
                Inches(0.30),
            )
            atf = arr.text_frame
            atf.clear()
            ap = atf.paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER
            ar = ap.add_run()
            ar.text = "→"
            ar.font.name = BBVA_FONT_HEAD
            ar.font.size = Pt(16)
            ar.font.bold = True
            ar.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    metric_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(2.45),
        Inches(12.0),
        Inches(0.88),
    )
    metric_box.fill.solid()
    metric_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    metric_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    mtf = metric_box.text_frame
    mtf.clear()
    metrics = [
        f"Prob. {focus_name}: {_fmt_pct_or_nd(chain_row.get('detractor_probability', np.nan))}",
        f"Delta NPS: {_fmt_signed_or_nd(chain_row.get('nps_delta_expected', np.nan))}",
        f"Impacto total: {_fmt_num_or_nd(chain_row.get('total_nps_impact', 0.0))} pts",
        f"Links validados: {int(_safe_int(chain_row.get('linked_pairs', 0), default=0))}",
        f"Confianza: {_fmt_num_or_nd(chain_row.get('confidence', 0.0))}",
    ]
    for m_idx, metric in enumerate(metrics):
        p = mtf.paragraphs[0] if m_idx == 0 else mtf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = metric
        r.font.name = BBVA_FONT_BODY
        r.font.size = Pt(11)
        r.font.color.rgb = _rgb(BBVA_COLORS["ink"])

    helix_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(3.55),
        Inches(5.85),
        Inches(2.45),
    )
    helix_box.fill.solid()
    helix_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    helix_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    htf = helix_box.text_frame
    htf.clear()
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = _chain_header("Evidencia Helix", shown_incidents, linked_incidents)
    hr.font.name = BBVA_FONT_HEAD
    hr.font.size = Pt(18)
    hr.font.bold = True
    hr.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    if not helix_lines:
        helix_lines = ["No hay suficiente evidencia Helix validada para elevar otro caso."]
    for idx_line, line in enumerate(helix_lines):
        p = htf.add_paragraph()
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(11.5)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])
        if idx_line < len(helix_records):
            rec = helix_records[idx_line]
            rec_id = str(rec.get("incident_id", "")).strip()
            rec_summary = _clip(str(rec.get("summary", "")).strip(), 150)
            rec_url = str(rec.get("url", "")).strip()
            bullet = p.add_run()
            bullet.text = "• "
            bullet.font.name = BBVA_FONT_BODY
            bullet.font.size = Pt(11.5)
            bullet.font.color.rgb = _rgb(BBVA_COLORS["muted"])
            if rec_id:
                id_run = p.add_run()
                id_run.text = rec_id
                id_run.font.name = BBVA_FONT_MEDIUM
                id_run.font.size = Pt(11.5)
                id_run.font.bold = True
                id_run.font.color.rgb = _rgb(BBVA_COLORS["blue"])
                if rec_url.startswith(("http://", "https://", "file://")):
                    id_run.hyperlink.address = rec_url
            if rec_summary:
                sep_run = p.add_run()
                sep_run.text = " · " if rec_id else ""
                sep_run.font.name = BBVA_FONT_BODY
                sep_run.font.size = Pt(11.5)
                sep_run.font.color.rgb = _rgb(BBVA_COLORS["muted"])
                text_run = p.add_run()
                text_run.text = rec_summary
                text_run.font.name = BBVA_FONT_BODY
                text_run.font.size = Pt(11.5)
                text_run.font.color.rgb = _rgb(BBVA_COLORS["muted"])
        else:
            p.text = f"• {_clip(line, 170)}"

    voc_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.82),
        Inches(3.55),
        Inches(5.88),
        Inches(2.45),
    )
    voc_box.fill.solid()
    voc_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    voc_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    vtf = voc_box.text_frame
    vtf.clear()
    vp = vtf.paragraphs[0]
    vr = vp.add_run()
    vr.text = _chain_header("Evidencia Voz del Cliente", shown_comments, linked_comments)
    vr.font.name = BBVA_FONT_HEAD
    vr.font.size = Pt(18)
    vr.font.bold = True
    vr.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    if not voc_lines:
        voc_lines = [
            "No hay suficiente evidencia VoC enlazada para construir otro relato defendible."
        ]
    for line in voc_lines:
        p = vtf.add_paragraph()
        p.text = f"• {_clip(line, 170)}"
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(11.5)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    footer_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.70),
        Inches(6.18),
        Inches(12.0),
        Inches(0.84),
    )
    footer_box.fill.solid()
    footer_box.fill.fore_color.rgb = _rgb(BBVA_COLORS["white"])
    footer_box.line.color.rgb = _rgb(BBVA_COLORS["line"])
    ftf = footer_box.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "Priorización del tema"
    fr.font.name = BBVA_FONT_HEAD
    fr.font.size = Pt(13)
    fr.font.bold = True
    fr.font.color.rgb = _rgb(BBVA_COLORS["ink"])
    for line in _chain_priority_summary(chain_row, focus_name=focus_name):
        p = ftf.add_paragraph()
        p.text = _clip(line, 175)
        p.font.name = BBVA_FONT_BODY
        p.font.size = Pt(10.6)
        p.font.color.rgb = _rgb(BBVA_COLORS["muted"])

    concl = slide.shapes.add_textbox(Inches(0.78), Inches(7.04), Inches(11.8), Inches(0.22))
    ctf = concl.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cr = cp.add_run()
    cr.text = _clip(chain_row.get("chain_story", ""), 170)
    cr.font.name = BBVA_FONT_BODY
    cr.font.size = Pt(9.8)
    cr.font.color.rgb = _rgb(BBVA_COLORS["muted"])


def _pick_first_col(df: pd.DataFrame, candidates: list[str]) -> str:
    lower_map = {str(c).strip().lower(): str(c) for c in df.columns}
    for c in candidates:
        hit = lower_map.get(str(c).strip().lower())
        if hit:
            return hit
    return ""


def _prepare_daily_signals(
    overall: pd.DataFrame,
    *,
    period_start: Optional[date],
    period_end: Optional[date],
) -> tuple[pd.DataFrame, bool]:
    """Normalize timeline to daily grain with NPS mean, detractor share and incidents."""
    if overall is None or overall.empty:
        return pd.DataFrame(columns=["date", "nps_mean", "detractor_rate", "incidents"]), False

    d = overall.copy()
    time_col = "date" if "date" in d.columns else ("week" if "week" in d.columns else "")
    if not time_col:
        return pd.DataFrame(columns=["date", "nps_mean", "detractor_rate", "incidents"]), False

    d[time_col] = _coerce_datetime_series(d[time_col])
    d = d.dropna(subset=[time_col]).copy()
    if d.empty:
        return pd.DataFrame(columns=["date", "nps_mean", "detractor_rate", "incidents"]), False

    if time_col == "week":
        d["date"] = d[time_col].dt.normalize()
    else:
        d["date"] = d[time_col].dt.normalize()

    inc_col = _pick_first_col(d, ["incidents", "incidencias", "incident_count"])
    focus_col = _pick_first_col(d, ["focus_rate", "detractor_rate", "rate_detractors"])
    nps_col = _pick_first_col(
        d,
        [
            "nps_mean",
            "nps_avg",
            "nps_media",
            "nps",
            "nps_score",
            "score_mean",
            "nps_current",
        ],
    )
    responses_col = _pick_first_col(d, ["responses", "respuestas", "n"])

    d["incidents"] = pd.to_numeric(d.get(inc_col, 0.0), errors="coerce").fillna(0.0).clip(lower=0.0)
    d["detractor_rate"] = (
        pd.to_numeric(d.get(focus_col, np.nan), errors="coerce").fillna(0.0).clip(0.0, 1.0)
    )

    nps_estimated = False
    if nps_col:
        d["nps_mean"] = pd.to_numeric(d[nps_col], errors="coerce")
    else:
        # Fallback when daily mean NPS is not available in aggregates.
        d["nps_mean"] = (1.0 - d["detractor_rate"]) * 10.0
        nps_estimated = True

    if responses_col:
        d["responses"] = (
            pd.to_numeric(d[responses_col], errors="coerce").fillna(0.0).clip(lower=0.0)
        )
    else:
        d["responses"] = 1.0

    if period_start is not None:
        d = d[d["date"] >= pd.Timestamp(period_start)]
    if period_end is not None:
        d = d[d["date"] <= pd.Timestamp(period_end)]
    if d.empty:
        return (
            pd.DataFrame(columns=["date", "nps_mean", "detractor_rate", "incidents"]),
            nps_estimated,
        )

    d["_w"] = np.maximum(pd.to_numeric(d["responses"], errors="coerce").fillna(0.0), 1.0)
    d["_nps_w"] = pd.to_numeric(d["nps_mean"], errors="coerce").fillna(0.0) * d["_w"]
    d["_det_w"] = pd.to_numeric(d["detractor_rate"], errors="coerce").fillna(0.0) * d["_w"]
    agg = (
        d.groupby("date", as_index=False)
        .agg(
            incidents=("incidents", "sum"),
            responses=("responses", "sum"),
            _w=("_w", "sum"),
            _nps_w=("_nps_w", "sum"),
            _det_w=("_det_w", "sum"),
        )
        .sort_values("date")
    )
    agg["nps_mean"] = agg["_nps_w"] / agg["_w"].replace({0.0: np.nan})
    agg["detractor_rate"] = agg["_det_w"] / agg["_w"].replace({0.0: np.nan})
    agg = agg.drop(columns=["_w", "_nps_w", "_det_w"])

    start_d = pd.Timestamp(period_start) if period_start is not None else agg["date"].min()
    end_d = pd.Timestamp(period_end) if period_end is not None else agg["date"].max()
    if pd.isna(start_d) or pd.isna(end_d) or start_d > end_d:
        return (
            pd.DataFrame(columns=["date", "nps_mean", "detractor_rate", "incidents"]),
            nps_estimated,
        )

    idx = pd.date_range(start=start_d.normalize(), end=end_d.normalize(), freq="D")
    out = agg.set_index("date").reindex(idx).rename_axis("date").reset_index()
    out["incidents"] = pd.to_numeric(out["incidents"], errors="coerce").fillna(0.0).clip(lower=0.0)

    for c in ["nps_mean", "detractor_rate"]:
        vals = pd.to_numeric(out[c], errors="coerce")
        if int(vals.notna().sum()) >= 2:
            vals = vals.interpolate(limit_direction="both")
        elif int(vals.notna().sum()) == 1:
            vals = vals.fillna(float(vals.dropna().iloc[0]))
        else:
            vals = vals.fillna(0.0)
        out[c] = vals

    out["nps_mean"] = out["nps_mean"].clip(0.0, 10.0)
    out["detractor_rate"] = out["detractor_rate"].clip(0.0, 1.0)
    return out[["date", "nps_mean", "detractor_rate", "incidents"]].copy(), nps_estimated


def _history_fig(  # pragma: no cover - legacy chart kept for backwards compatibility
    daily: pd.DataFrame, *, focus_name: str
) -> Optional[go.Figure]:
    if daily is None or daily.empty:
        return None
    d = daily.sort_values("date").copy()

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["nps_mean"],
            name="NPS medio",
            mode="lines+markers",
            line=dict(color="#" + BBVA_COLORS["blue"], width=3.2),
            marker=dict(
                size=8,
                color=_ppt_nps_marker_colors(d["nps_mean"]),
                line=dict(color="#" + BBVA_COLORS["bg_light"], width=1),
            ),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["detractor_rate"] * 100.0,
            name=f"Opiniones {focus_name}",
            mode="lines",
            yaxis="y2",
            line=dict(color="#" + BBVA_COLORS["red"], width=2.6),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["incidents"],
            name="Incidencias",
            yaxis="y3",
            marker=dict(color="#" + BBVA_COLORS["yellow"]),
            opacity=0.78,
        )
    )

    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=84, t=22, b=24),
        legend=dict(orientation="h", x=0.01, y=1.08),
        xaxis=dict(title="Día"),
        yaxis=dict(title="NPS medio (0-10)", range=[0, 10]),
        yaxis2=dict(
            title=f"% {focus_name}",
            overlaying="y",
            side="right",
            range=[0, 100],
            showgrid=False,
        ),
        yaxis3=dict(
            title="Incidencias",
            overlaying="y",
            side="right",
            anchor="free",
            position=0.92,
            showgrid=False,
            rangemode="tozero",
        ),
    )
    return fig


def _hotspot_daily_breakdown(
    daily_signals: pd.DataFrame,
    incident_evidence_df: Optional[pd.DataFrame],
    incident_timeline_df: Optional[pd.DataFrame] = None,
) -> tuple[pd.DataFrame, dict[int, str]]:
    return build_hotspot_daily_breakdown_metrics(
        daily_signals,
        incident_evidence_df,
        incident_timeline_df,
        max_hotspots=3,
    )


def _hotspot_stack_fig(
    daily_signals: pd.DataFrame,
    incident_evidence_df: Optional[pd.DataFrame],
    incident_timeline_df: Optional[pd.DataFrame] = None,
) -> Optional[go.Figure]:  # pragma: no cover - legacy chart kept for backwards compatibility
    d, term_by_rank = _hotspot_daily_breakdown(
        daily_signals, incident_evidence_df, incident_timeline_df
    )
    if d.empty:
        return None

    n_lbl = "No hotspot"
    h1_lbl = (
        f"Hotspot 1: {term_by_rank.get(1, '')}".strip(": ")
        if term_by_rank.get(1, "")
        else "Hotspot 1"
    )
    h2_lbl = (
        f"Hotspot 2: {term_by_rank.get(2, '')}".strip(": ")
        if term_by_rank.get(2, "")
        else "Hotspot 2"
    )
    h3_lbl = (
        f"Hotspot 3: {term_by_rank.get(3, '')}".strip(": ")
        if term_by_rank.get(3, "")
        else "Hotspot 3"
    )

    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["no_hotspot"],
            name=n_lbl,
            marker=dict(color="#" + BBVA_COLORS["blue"]),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["hotspot_3"],
            name=h3_lbl,
            marker=dict(color="#" + BBVA_COLORS["yellow"]),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["hotspot_2"],
            name=h2_lbl,
            marker=dict(color="#" + BBVA_COLORS["orange"]),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["hotspot_1"],
            name=h1_lbl,
            marker=dict(color="#" + BBVA_COLORS["red"]),
        )
    )

    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=24, t=22, b=96),
        barmode="stack",
        xaxis=dict(title="Día"),
        yaxis=dict(title="Incidencias registradas", rangemode="tozero"),
        legend=dict(
            orientation="h",
            x=0.0,
            y=-0.24,
            xanchor="left",
            yanchor="top",
            font=dict(size=12),
        ),
    )
    return fig


def _top_hotspots_fig(
    incident_evidence_df: Optional[pd.DataFrame],
    incident_timeline_df: Optional[pd.DataFrame],
    *,
    top_k: int = 3,
) -> Optional[go.Figure]:  # pragma: no cover - legacy chart kept for backwards compatibility
    hs = summarize_hotspot_counts(
        incident_evidence_df,
        incident_timeline_df,
        max_hotspots=max(1, int(top_k)),
    )
    if hs is None or hs.empty:
        return None

    d = hs.copy()
    d["hot_term"] = d.get("hot_term", "").astype(str).str.strip()
    d = d[d["hot_term"] != ""].copy()
    if d.empty:
        return None

    # Keep strict coherence with centralized hotspot summary:
    # rank and totals come directly from summarize_hotspot_counts.
    d["hot_rank"] = pd.to_numeric(d.get("hot_rank"), errors="coerce").fillna(999).astype(int)
    d["mention_comments"] = pd.to_numeric(d.get("mention_comments"), errors="coerce").fillna(0.0)
    d["hotspot_comments"] = pd.to_numeric(d.get("hotspot_comments"), errors="coerce").fillna(0.0)
    d["chart_nps_comments"] = pd.to_numeric(d.get("chart_nps_comments"), errors="coerce").fillna(
        0.0
    )
    d["impact"] = d["chart_nps_comments"].clip(lower=0.0)
    d.loc[d["impact"] <= 0.0, "impact"] = d["hotspot_comments"].clip(lower=0.0)
    d.loc[d["impact"] <= 0.0, "impact"] = d["mention_comments"].clip(lower=0.0)
    d = d.sort_values(["hot_rank"]).head(int(top_k)).copy()
    if d.empty:
        return None

    d["label"] = d["hot_term"].astype(str).str.upper().str.slice(0, 44)
    rank_color = {
        1: "#" + BBVA_COLORS["red"],
        2: "#" + BBVA_COLORS["orange"],
        3: "#" + BBVA_COLORS["yellow"],
    }
    d["bar_color"] = d["hot_rank"].map(rank_color).fillna("#" + BBVA_COLORS["red"])
    d["inbar"] = [
        f"#{int(r)}  {lbl}  ·  {int(v)}"
        for r, lbl, v in zip(d["hot_rank"], d["label"], d["impact"])
    ]
    d = d.iloc[::-1].copy()

    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=d["impact"],
            y=d["label"],
            orientation="h",
            marker=dict(color=d["bar_color"].tolist()),
            text=d["inbar"].tolist(),
            textposition="inside",
            insidetextanchor="start",
            textfont=dict(size=19, color="#" + BBVA_COLORS["white"]),
            customdata=d[
                [
                    "hot_rank",
                    "mention_incidents",
                    "mention_comments",
                    "hotspot_links",
                    "chart_nps_comments",
                ]
            ].to_numpy(),
            hovertemplate=(
                "Hotspot #%{customdata[0]:.0f}=%{y}<br>"
                "Comentarios del gráfico=%{customdata[4]:.0f}<br>"
                "Comentarios negativos=%{customdata[2]:.0f}<br>"
                "Incidencias Helix=%{customdata[1]:.0f}<br>"
                "Links validados=%{customdata[3]:.0f}<extra></extra>"
            ),
        )
    )
    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=24, t=22, b=24),
        uniformtext=dict(minsize=16, mode="show"),
        xaxis=dict(
            title="",
            showticklabels=False,
            showgrid=False,
            zeroline=False,
            visible=False,
        ),
        yaxis=dict(title=""),
        showlegend=False,
    )
    return fig


def _hotspot_matches_by_day(
    incident_timeline_df: Optional[pd.DataFrame],
    incident_evidence_df: Optional[pd.DataFrame],
    *,
    month_start: pd.Timestamp,
    month_end: pd.Timestamp,
) -> pd.DataFrame:
    cols = ["date", "matched_incidents"]
    if incident_timeline_df is None or incident_timeline_df.empty:
        return pd.DataFrame(columns=cols)

    t = incident_timeline_df.copy()
    req = {"date", "helix_records", "nps_comments"}
    if not req.issubset(set(t.columns)):
        return pd.DataFrame(columns=cols)

    t["date"] = _coerce_datetime_series(t["date"])
    t = t.dropna(subset=["date"])
    t = t[(t["date"] >= month_start) & (t["date"] <= month_end)].copy()
    if t.empty:
        return pd.DataFrame(columns=cols)

    t["helix_records"] = (
        pd.to_numeric(t["helix_records"], errors="coerce").fillna(0.0).clip(lower=0.0)
    )
    t["nps_comments"] = (
        pd.to_numeric(t["nps_comments"], errors="coerce").fillna(0.0).clip(lower=0.0)
    )

    # Prefer hotspot-level rows (hot_term pre-aggregated in the app payload).
    top_terms: list[str] = []
    if (
        incident_evidence_df is not None
        and not incident_evidence_df.empty
        and {"hot_term", "hot_rank"}.issubset(set(incident_evidence_df.columns))
    ):
        e = incident_evidence_df.copy()
        e["hot_term"] = e["hot_term"].astype(str).str.strip()
        e["hot_rank"] = pd.to_numeric(e["hot_rank"], errors="coerce")
        e = e.dropna(subset=["hot_rank"])
        e = e[e["hot_term"] != ""]
        if not e.empty:
            e = e.sort_values(["hot_rank"])
            top_terms = e["hot_term"].drop_duplicates().head(3).tolist()

    if "hot_term" in t.columns:
        t["hot_term"] = t["hot_term"].astype(str).str.strip()
        t = (
            t[t["hot_term"].isin(top_terms)].copy()
            if top_terms
            else t[t["hot_term"].str.strip() != ""].copy()
        )

    if t.empty:
        return pd.DataFrame(columns=cols)

    # Coincidence means there is operational signal and detractor signal on the day.
    t = t[(t["helix_records"] > 0) & (t["nps_comments"] > 0)].copy()
    if t.empty:
        return pd.DataFrame(columns=cols)

    out = (
        t.groupby("date", as_index=False)
        .agg(matched_incidents=("helix_records", "sum"))
        .sort_values("date")
    )
    out["matched_incidents"] = (
        pd.to_numeric(out["matched_incidents"], errors="coerce").fillna(0.0).clip(lower=0.0)
    )
    return out


def _month_overlap_fig(
    month_daily: pd.DataFrame,
    *,
    focus_name: str,
    matched_daily: pd.DataFrame,
    matched_label: str = "Incidencias con match detractor",
) -> Optional[go.Figure]:  # pragma: no cover - legacy chart kept for backwards compatibility
    if month_daily is None or month_daily.empty:
        return None

    d = month_daily.sort_values("date").copy()
    m = (
        matched_daily.copy()
        if matched_daily is not None
        else pd.DataFrame(columns=["date", "matched_incidents"])
    )
    m["date"] = _coerce_datetime_series(m.get("date"))
    m = m.dropna(subset=["date"])
    m = m.groupby("date", as_index=False).agg(matched_incidents=("matched_incidents", "sum"))
    d = d.merge(m, on="date", how="left")
    d["matched_incidents"] = pd.to_numeric(d["matched_incidents"], errors="coerce").fillna(0.0)

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["nps_mean"],
            name="NPS medio",
            mode="lines+markers",
            line=dict(color="#" + BBVA_COLORS["blue"], width=3.2),
            marker=dict(
                size=8,
                color=_ppt_nps_marker_colors(d["nps_mean"]),
                line=dict(color="#" + BBVA_COLORS["bg_light"], width=1),
            ),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=d["date"],
            y=d["detractor_rate"] * 100.0,
            name=f"Opiniones {focus_name}",
            mode="lines",
            yaxis="y2",
            line=dict(color="#" + BBVA_COLORS["red"], width=2.6),
        )
    )
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["incidents"],
            name="Incidencias totales",
            yaxis="y3",
            marker=dict(color="#" + BBVA_COLORS["yellow"]),
            opacity=0.42,
        )
    )

    highlight_text = [str(int(v)) if float(v) > 0 else "" for v in d["matched_incidents"].tolist()]
    fig.add_trace(
        go.Bar(
            x=d["date"],
            y=d["matched_incidents"],
            name=matched_label,
            yaxis="y3",
            marker=dict(color="#" + BBVA_COLORS["orange"]),
            opacity=0.88,
            text=highlight_text,
            textposition="outside",
        )
    )

    fig.update_layout(
        template="plotly_white",
        margin=dict(l=24, r=84, t=22, b=24),
        legend=dict(orientation="h", x=0.01, y=1.10),
        barmode="overlay",
        xaxis=dict(title="Día"),
        yaxis=dict(title="NPS medio (0-10)", range=[0, 10]),
        yaxis2=dict(
            title=f"% {focus_name}",
            overlaying="y",
            side="right",
            range=[0, 100],
            showgrid=False,
        ),
        yaxis3=dict(
            title="Incidencias",
            overlaying="y",
            side="right",
            anchor="free",
            position=0.92,
            showgrid=False,
            rangemode="tozero",
        ),
    )
    return fig


def _prepare_incident_evidence(  # pragma: no cover - legacy helper kept for backwards compatibility
    incident_evidence_df: Optional[pd.DataFrame],
) -> pd.DataFrame:
    cols = list(HOTSPOT_EVIDENCE_COLUMNS)
    if incident_evidence_df is None or incident_evidence_df.empty:
        return pd.DataFrame(columns=cols)

    d = incident_evidence_df.copy()
    id_col = _pick_first_col(
        d, ["incident_id", "incident number", "incident_number", "id de la incidencia"]
    )
    topic_col = _pick_first_col(d, ["nps_topic", "topic", "topico", "tópico"]) or ""
    date_col = _pick_first_col(d, ["incident_date", "fecha", "date", "opened_at"])
    summary_candidates = [
        "incident_summary",
        "Detailed Description",
        "Detailed Decription",
        "bbva_detaileddescription",
        "description",
        "descripcion",
        "descripción",
        "short description",
        "brief description",
        "bbva_shortdescription",
        "summary",
    ]
    comment_col = _pick_first_col(
        d,
        [
            "detractor_comment",
            "nps_comment",
            "comment",
            "comentario",
            "comentario detractor",
        ],
    )
    sim_col = _pick_first_col(d, ["similarity", "sim", "score"])
    hot_term_col = _pick_first_col(d, ["hot_term", "termino_caliente", "term"])
    hot_rank_col = _pick_first_col(d, ["hot_rank", "hotterm_rank", "rank"])
    mention_inc_col = _pick_first_col(
        d,
        ["mention_incidents", "hotspot_incidents_mentions", "mentions_incidents"],
    )
    mention_com_col = _pick_first_col(
        d,
        ["mention_comments", "hotspot_comments_mentions", "mentions_comments"],
    )
    hot_inc_col = _pick_first_col(d, ["hotspot_incidents", "cluster_incidents", "incidents_count"])
    hot_com_col = _pick_first_col(d, ["hotspot_comments", "cluster_comments", "comments_count"])
    hot_lnk_col = _pick_first_col(d, ["hotspot_links", "cluster_links", "links_count"])
    lower_map = {str(c).strip().lower(): str(c) for c in d.columns}
    summary_cols: list[str] = []
    for cand in summary_candidates:
        hit = lower_map.get(str(cand).strip().lower())
        if not hit:
            continue
        if hit not in summary_cols:
            summary_cols.append(hit)

    if summary_cols:
        summary_stack = pd.concat(
            [
                d[col]
                .astype(str)
                .fillna("")
                .str.strip()
                .replace({"nan": "", "NaN": "", "None": "", "NaT": ""})
                for col in summary_cols
            ],
            axis=1,
        )
        arr = summary_stack.to_numpy(dtype=object)
        non_empty = arr != ""
        first_idx = non_empty.argmax(axis=1)
        has_any = non_empty.any(axis=1)
        vals = np.full(len(summary_stack), "", dtype=object)
        rows = np.arange(len(summary_stack))
        vals[has_any] = arr[rows[has_any], first_idx[has_any]]
        incident_summary = pd.Series(vals, index=d.index)
    else:
        incident_summary = pd.Series([""] * len(d), index=d.index)

    out = pd.DataFrame(
        {
            "incident_id": d[id_col].astype(str) if id_col else "",
            "incident_date": _coerce_datetime_series(d[date_col]) if date_col else pd.NaT,
            "nps_topic": d[topic_col].astype(str) if topic_col else "",
            "incident_summary": incident_summary,
            "detractor_comment": d[comment_col].astype(str) if comment_col else "",
            "similarity": (
                pd.to_numeric(d[sim_col], errors="coerce").fillna(0.0) if sim_col else 0.0
            ),
            "hot_term": d[hot_term_col].astype(str) if hot_term_col else "",
            "hot_rank": pd.to_numeric(d[hot_rank_col], errors="coerce") if hot_rank_col else np.nan,
            "mention_incidents": (
                pd.to_numeric(d[mention_inc_col], errors="coerce").fillna(0.0)
                if mention_inc_col
                else 0.0
            ),
            "mention_comments": (
                pd.to_numeric(d[mention_com_col], errors="coerce").fillna(0.0)
                if mention_com_col
                else 0.0
            ),
            "hotspot_incidents": (
                pd.to_numeric(d[hot_inc_col], errors="coerce").fillna(0.0) if hot_inc_col else 0.0
            ),
            "hotspot_comments": (
                pd.to_numeric(d[hot_com_col], errors="coerce").fillna(0.0) if hot_com_col else 0.0
            ),
            "hotspot_links": (
                pd.to_numeric(d[hot_lnk_col], errors="coerce").fillna(0.0) if hot_lnk_col else 0.0
            ),
        }
    )
    out["incident_id"] = out["incident_id"].astype(str).str.strip()
    out["nps_topic"] = out["nps_topic"].astype(str).str.strip()
    out["incident_summary"] = out["incident_summary"].astype(str).str.strip()
    out["detractor_comment"] = out["detractor_comment"].astype(str).str.strip()
    out["hot_term"] = out["hot_term"].astype(str).str.strip()
    out["mention_incidents"] = (
        pd.to_numeric(out["mention_incidents"], errors="coerce").fillna(0.0).astype(int)
    )
    out["mention_comments"] = (
        pd.to_numeric(out["mention_comments"], errors="coerce").fillna(0.0).astype(int)
    )
    out["hotspot_incidents"] = (
        pd.to_numeric(out["hotspot_incidents"], errors="coerce").fillna(0.0).astype(int)
    )
    out["hotspot_comments"] = (
        pd.to_numeric(out["hotspot_comments"], errors="coerce").fillna(0.0).astype(int)
    )
    out["hotspot_links"] = (
        pd.to_numeric(out["hotspot_links"], errors="coerce").fillna(0.0).astype(int)
    )
    out = out[(out["incident_id"] != "") | (out["nps_topic"] != "")].copy()
    out = out.sort_values(
        ["hot_rank", "similarity"], ascending=[True, False], na_position="last"
    ).reset_index(drop=True)
    return out[cols]


def _top_topics_for_zoom(
    rationale_df: pd.DataFrame,
    ranking_df: Optional[pd.DataFrame],
    *,
    max_topics: int = 3,
) -> list[str]:  # pragma: no cover - legacy helper kept for backwards compatibility
    topics: list[str] = []
    if rationale_df is not None and not rationale_df.empty and "nps_topic" in rationale_df.columns:
        src = rationale_df.copy()
        sort_cols = []
        if "priority" in src.columns:
            sort_cols.append("priority")
        if "nps_points_at_risk" in src.columns:
            sort_cols.append("nps_points_at_risk")
        if sort_cols:
            src = src.sort_values(sort_cols, ascending=False)
        topics.extend(src["nps_topic"].astype(str).tolist())

    if ranking_df is not None and not ranking_df.empty and "nps_topic" in ranking_df.columns:
        src = ranking_df.copy()
        metric = "confidence_learned" if "confidence_learned" in src.columns else "confidence"
        if metric in src.columns:
            src[metric] = pd.to_numeric(src[metric], errors="coerce").fillna(0.0)
            src = src.sort_values([metric], ascending=False)
        topics.extend(src["nps_topic"].astype(str).tolist())

    out: list[str] = []
    seen: set[str] = set()
    for t in topics:
        key = str(t).strip()
        if not key or key in seen:
            continue
        out.append(key)
        seen.add(key)
        if len(out) >= int(max_topics):
            break
    return out


def _select_zoom_incidents(
    top_topics: list[str],
    incident_evidence_df: pd.DataFrame,
    *,
    max_items: int = 3,
) -> list[ZoomIncident]:  # pragma: no cover - legacy helper kept for backwards compatibility
    selected: list[ZoomIncident] = []

    if incident_evidence_df is not None and not incident_evidence_df.empty:
        ev = incident_evidence_df.copy()
        ev["hot_rank"] = pd.to_numeric(ev.get("hot_rank"), errors="coerce")
        ev["similarity"] = pd.to_numeric(ev.get("similarity"), errors="coerce").fillna(0.0)
        ev["hot_term"] = ev.get("hot_term", "").astype(str).str.strip()
        ev["incident_id"] = ev.get("incident_id", "").astype(str).str.strip()
        ev["detractor_comment"] = ev.get("detractor_comment", "").astype(str)
        ev["_has_comment"] = ev["detractor_comment"].str.strip().ne("")

        hotspots: list[tuple[float, str, pd.DataFrame]] = []
        ranked = ev.dropna(subset=["hot_rank"])
        if not ranked.empty:
            for (rk, term), g in ranked.groupby(["hot_rank", "hot_term"], dropna=False):
                hotspots.append((float(rk), str(term or "").strip(), g.copy()))
            hotspots.sort(
                key=lambda x: (
                    x[0],
                    -float(pd.to_numeric(x[2]["similarity"], errors="coerce").fillna(0.0).mean()),
                )
            )
        else:
            for topic in top_topics:
                g = ev[ev["nps_topic"].astype(str) == str(topic)].copy()
                if g.empty:
                    continue
                hotspots.append((999.0, "", g))
            if not hotspots:
                hotspots = [(999.0, "", ev.copy())]

        for _rk, term, g in hotspots:
            if g.empty:
                continue
            g = g.sort_values(["_has_comment", "similarity"], ascending=[False, False]).copy()
            rep = g.iloc[0]

            inc_id = str(rep.get("incident_id", "")).strip()
            topic = str(rep.get("nps_topic", "")).strip() or "Tópico sin etiqueta"
            hot_term = _clip(term or rep.get("hot_term", ""), 48)

            mention_inc = _safe_int(g.get("mention_incidents", pd.Series([0])).max(), default=0)
            mention_com = _safe_int(g.get("mention_comments", pd.Series([0])).max(), default=0)
            hotspot_inc = _safe_int(g.get("hotspot_incidents", pd.Series([0])).max(), default=0)
            hotspot_com = _safe_int(g.get("hotspot_comments", pd.Series([0])).max(), default=0)
            hotspot_lnk = _safe_int(g.get("hotspot_links", pd.Series([0])).max(), default=0)
            if mention_inc <= 0:
                mention_inc = int(hotspot_inc)
            if mention_com <= 0:
                mention_com = int(hotspot_com)

            sample_incidents = (
                g["incident_id"]
                .astype(str)
                .str.strip()
                .replace("", np.nan)
                .dropna()
                .drop_duplicates()
                .head(3)
                .tolist()
            )
            sample_comments = (
                g["detractor_comment"]
                .astype(str)
                .str.strip()
                .replace("", np.nan)
                .dropna()
                .drop_duplicates()
                .head(2)
                .tolist()
            )

            selected.append(
                ZoomIncident(
                    incident_id=inc_id or f"N/D-{len(selected) + 1}",
                    incident_date=_safe_dt(rep.get("incident_date")),
                    nps_topic=topic,
                    incident_summary=_clip(rep.get("incident_summary", ""), 190),
                    detractor_comment=_clip(rep.get("detractor_comment", ""), 220),
                    similarity=_safe_float(rep.get("similarity", 0.0), default=0.0),
                    hot_term=hot_term,
                    mention_incidents=mention_inc,
                    mention_comments=mention_com,
                    hotspot_incidents=hotspot_inc,
                    hotspot_comments=hotspot_com,
                    hotspot_links=hotspot_lnk,
                    sample_incidents=", ".join(sample_incidents),
                    sample_comments=" | ".join([_clip(c, 95) for c in sample_comments]),
                )
            )
            if len(selected) >= int(max_items):
                return selected

    while len(selected) < int(max_items):
        idx = len(selected) + 1
        topic = top_topics[idx - 1] if idx - 1 < len(top_topics) else "Tópico sin evidencia"
        selected.append(
            ZoomIncident(
                incident_id=f"N/D-{idx}",
                incident_date=None,
                nps_topic=topic,
                incident_summary="No se encontró descripción de incidencia para este tópico en el periodo.",
                detractor_comment="No se encontró comentario detractor validado con la política estricta activa.",
                similarity=0.0,
                hot_term="",
                mention_incidents=0,
                mention_comments=0,
                hotspot_incidents=0,
                hotspot_comments=0,
                hotspot_links=0,
                sample_incidents="",
                sample_comments="",
            )
        )
    return selected[: int(max_items)]


def _hotspot_summary_row(
    incident: ZoomIncident,
    hotspot_summary: pd.DataFrame,
) -> Optional[pd.Series]:  # pragma: no cover - legacy helper kept for backwards compatibility
    if hotspot_summary is None or hotspot_summary.empty:
        return None

    hot_term_key = str(incident.hot_term or "").strip()
    if hot_term_key and "hot_term" in hotspot_summary.columns:
        hit = hotspot_summary[
            hotspot_summary["hot_term"].astype(str).str.strip() == hot_term_key
        ].head(1)
        if not hit.empty:
            return hit.iloc[0]

    if "hot_rank" in hotspot_summary.columns:
        ranked = pd.to_numeric(hotspot_summary["hot_rank"], errors="coerce")
        if np.isfinite(ranked).any():
            order = hotspot_summary.assign(_rank=ranked).sort_values(["_rank"], na_position="last")
            return order.head(1).iloc[0]
    return None


def _incident_related_timeline(
    *,
    incident_id: str,
    hot_term: str = "",
    incident_timeline_df: Optional[pd.DataFrame],
) -> pd.DataFrame:  # pragma: no cover - legacy helper kept for backwards compatibility
    cols = [
        "date",
        "helix_records",
        "nps_comments",
        "nps_comments_moderate",
        "nps_comments_high",
        "nps_comments_critical",
        "incident_ids",
    ]
    if incident_timeline_df is None or incident_timeline_df.empty:
        return pd.DataFrame(columns=cols)
    req = {"incident_id", "date", "helix_records", "nps_comments"}
    if not req.issubset(set(incident_timeline_df.columns)):
        return pd.DataFrame(columns=cols)

    d = incident_timeline_df.copy()
    hot_term_key = str(hot_term or "").strip()
    if hot_term_key and "hot_term" in d.columns:
        d["hot_term"] = d["hot_term"].astype(str).str.strip()
        d = d[d["hot_term"] == hot_term_key].copy()
    else:
        d["incident_id"] = d["incident_id"].astype(str).str.strip()
        d = d[d["incident_id"] == str(incident_id).strip()].copy()
    if d.empty:
        return pd.DataFrame(columns=cols)

    d["date"] = _coerce_datetime_series(d["date"])
    d = d.dropna(subset=["date"])
    if d.empty:
        return pd.DataFrame(columns=cols)

    d["helix_records"] = (
        pd.to_numeric(d["helix_records"], errors="coerce").fillna(0.0).clip(lower=0.0)
    )
    d["nps_comments"] = (
        pd.to_numeric(d["nps_comments"], errors="coerce").fillna(0.0).clip(lower=0.0)
    )
    sev_mod = (
        d["nps_comments_moderate"]
        if "nps_comments_moderate" in d.columns
        else pd.Series([0.0] * len(d), index=d.index)
    )
    sev_high = (
        d["nps_comments_high"]
        if "nps_comments_high" in d.columns
        else pd.Series([0.0] * len(d), index=d.index)
    )
    sev_critical = (
        d["nps_comments_critical"]
        if "nps_comments_critical" in d.columns
        else pd.Series([0.0] * len(d), index=d.index)
    )
    d["nps_comments_moderate"] = pd.to_numeric(sev_mod, errors="coerce").fillna(0.0).clip(lower=0.0)
    d["nps_comments_high"] = pd.to_numeric(sev_high, errors="coerce").fillna(0.0).clip(lower=0.0)
    d["nps_comments_critical"] = (
        pd.to_numeric(sev_critical, errors="coerce").fillna(0.0).clip(lower=0.0)
    )
    if "incident_ids" not in d.columns:
        d["incident_ids"] = d.get("incident_id", "").astype(str)
    d["incident_ids"] = d["incident_ids"].astype(str).fillna("")

    def _merge_ids(values: pd.Series) -> str:
        raw: list[str] = []
        for v in values.astype(str).tolist():
            for part in str(v).split("|"):
                p = part.strip()
                if p:
                    raw.append(p)
        seen: set[str] = set()
        out: list[str] = []
        for item in raw:
            if item in seen:
                continue
            seen.add(item)
            out.append(item)
        return " | ".join(out)

    d = (
        d.groupby("date", as_index=False)
        .agg(
            helix_records=("helix_records", "sum"),
            nps_comments=("nps_comments", "sum"),
            nps_comments_moderate=("nps_comments_moderate", "sum"),
            nps_comments_high=("nps_comments_high", "sum"),
            nps_comments_critical=("nps_comments_critical", "sum"),
            incident_ids=("incident_ids", _merge_ids),
        )
        .sort_values("date")
    )
    d = d[(d["helix_records"] > 0) | (d["nps_comments"] > 0)].copy()
    return d[cols].copy()


def _zoom_incident_fig(
    *,
    topic_daily: pd.DataFrame,
    related_timeline: pd.DataFrame,
    incident: ZoomIncident,
    lag_days: int,
    changepoints: list[pd.Timestamp],
    focus_name: str,
    period_start: Optional[date] = None,
    period_end: Optional[date] = None,
) -> Optional[go.Figure]:  # pragma: no cover - legacy chart kept for backwards compatibility
    del incident, lag_days, focus_name

    rel = related_timeline.copy() if related_timeline is not None else pd.DataFrame()
    if not rel.empty:
        rel["date"] = _coerce_datetime_series(rel["date"])
        rel = rel.dropna(subset=["date"]).sort_values("date")
        rel["helix_records"] = (
            pd.to_numeric(rel["helix_records"], errors="coerce").fillna(0.0).clip(lower=0.0)
        )
        rel["nps_comments"] = (
            pd.to_numeric(rel["nps_comments"], errors="coerce").fillna(0.0).clip(lower=0.0)
        )
        rel["nps_comments_moderate"] = (
            pd.to_numeric(rel.get("nps_comments_moderate"), errors="coerce")
            .fillna(0.0)
            .clip(lower=0.0)
        )
        rel["nps_comments_high"] = (
            pd.to_numeric(rel.get("nps_comments_high"), errors="coerce").fillna(0.0).clip(lower=0.0)
        )
        rel["nps_comments_critical"] = (
            pd.to_numeric(rel.get("nps_comments_critical"), errors="coerce")
            .fillna(0.0)
            .clip(lower=0.0)
        )
        rel["incident_ids"] = rel.get("incident_ids", "").astype(str).fillna("")

        def _merge_ids(values: pd.Series) -> str:
            raw: list[str] = []
            for v in values.astype(str).tolist():
                for part in str(v).split("|"):
                    p = part.strip()
                    if p:
                        raw.append(p)
            seen: set[str] = set()
            out: list[str] = []
            for item in raw:
                if item in seen:
                    continue
                seen.add(item)
                out.append(item)
            return " | ".join(out)

        rel = (
            rel.groupby("date", as_index=False)
            .agg(
                helix_records=("helix_records", "sum"),
                nps_comments=("nps_comments", "sum"),
                nps_comments_moderate=("nps_comments_moderate", "sum"),
                nps_comments_high=("nps_comments_high", "sum"),
                nps_comments_critical=("nps_comments_critical", "sum"),
                incident_ids=("incident_ids", _merge_ids),
            )
            .sort_values("date")
        )

    if rel.empty:
        return None

    start_ts: Optional[pd.Timestamp] = None
    end_ts: Optional[pd.Timestamp] = None
    if period_start is not None and period_end is not None:
        start_ts = _safe_dt(period_start)
        end_ts = _safe_dt(period_end)
    if start_ts is None or end_ts is None or end_ts < start_ts:
        start_ts = pd.Timestamp(rel["date"].min()).normalize()
        end_ts = pd.Timestamp(rel["date"].max()).normalize()

    scope = pd.DataFrame({"date": pd.date_range(start=start_ts, end=end_ts, freq="D")})
    z = scope.merge(rel, on="date", how="left")
    z["helix_records"] = pd.to_numeric(z.get("helix_records"), errors="coerce").fillna(0.0)
    z["nps_comments"] = pd.to_numeric(z.get("nps_comments"), errors="coerce").fillna(0.0)
    z["nps_comments_moderate"] = pd.to_numeric(
        z.get("nps_comments_moderate"), errors="coerce"
    ).fillna(0.0)
    z["nps_comments_high"] = pd.to_numeric(z.get("nps_comments_high"), errors="coerce").fillna(0.0)
    z["nps_comments_critical"] = pd.to_numeric(
        z.get("nps_comments_critical"), errors="coerce"
    ).fillna(0.0)
    z["incident_ids"] = z.get("incident_ids", "").astype(str).fillna("")
    nps_series = pd.DataFrame(columns=["date", "nps_mean"])
    if (
        topic_daily is not None
        and not topic_daily.empty
        and {"date", "nps_mean"}.issubset(set(topic_daily.columns))
    ):
        nps_series = topic_daily[["date", "nps_mean"]].copy()
        nps_series["date"] = _coerce_datetime_series(nps_series["date"])
        nps_series["nps_mean"] = pd.to_numeric(nps_series["nps_mean"], errors="coerce")
        nps_series = nps_series.dropna(subset=["date", "nps_mean"])
        nps_series = (
            nps_series.groupby("date", as_index=False)
            .agg(nps_mean=("nps_mean", "mean"))
            .sort_values("date")
        )

    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=z["date"],
            y=z["nps_comments_moderate"],
            name="Comentarios moderados (NPS 5-6)",
            marker=dict(color="#FCA5A5"),
            opacity=0.88,
        )
    )
    fig.add_trace(
        go.Bar(
            x=z["date"],
            y=z["nps_comments_high"],
            name="Comentarios altos (NPS 3-4)",
            marker=dict(color="#EF4444"),
            opacity=0.88,
        )
    )
    fig.add_trace(
        go.Bar(
            x=z["date"],
            y=z["nps_comments_critical"],
            name="Comentarios críticos (NPS 0-2)",
            marker=dict(color="#B91C1C"),
            opacity=0.92,
        )
    )
    totals = z["nps_comments_moderate"] + z["nps_comments_high"] + z["nps_comments_critical"]
    total_labels = [str(int(v)) if float(v) > 0 else "" for v in totals.tolist()]
    fig.add_trace(
        go.Scatter(
            x=z["date"],
            y=totals,
            mode="text",
            text=total_labels,
            textposition="top center",
            showlegend=False,
            hoverinfo="skip",
            textfont=dict(size=10, color="#" + BBVA_COLORS["muted"]),
        )
    )

    points = z[z["helix_records"] > 0].copy()

    def _point_label(v: object) -> str:
        items = [p.strip() for p in str(v or "").split("|") if p.strip()]
        if not items:
            return ""
        inc_items = [it for it in items if str(it).upper().startswith("INC")]
        return str((inc_items[0] if inc_items else items[0])).strip()

    point_labels_raw = points["incident_ids"].map(_point_label).tolist()
    point_labels: list[str] = []
    last_labeled_date: Optional[pd.Timestamp] = None
    for dt, lbl in zip(points["date"].tolist(), point_labels_raw):
        label = str(lbl or "").strip()
        if not label:
            point_labels.append("")
            continue
        dt_ts = pd.Timestamp(dt).normalize()
        if last_labeled_date is not None and int((dt_ts - last_labeled_date).days) < 2:
            point_labels.append("")
            continue
        point_labels.append(label)
        last_labeled_date = dt_ts

    fig.add_trace(
        go.Scatter(
            x=points["date"],
            y=points["helix_records"],
            name="Incidencias Helix del hotspot",
            yaxis="y2",
            mode="markers+text",
            marker=dict(size=8),
            text=point_labels,
            textposition="top right",
            textfont=dict(size=9, color="#" + BBVA_COLORS["blue"]),
        )
    )

    if not nps_series.empty:
        fig.add_trace(
            go.Scatter(
                x=nps_series["date"],
                y=nps_series["nps_mean"],
                name="NPS medio",
                yaxis="y3",
                mode="lines+markers",
                line=dict(color="#" + BBVA_COLORS["blue"], width=2.6),
                marker=dict(
                    size=6,
                    color=_ppt_nps_marker_colors(nps_series["nps_mean"]),
                    line=dict(color="#" + BBVA_COLORS["bg_light"], width=1),
                ),
            )
        )

    cp_in_window = []
    min_d = z["date"].min()
    max_d = z["date"].max()
    for cp in changepoints:
        if cp < min_d or cp > max_d:
            continue
        cp_in_window.append(cp)
        fig.add_vline(
            x=cp,
            line_width=1.5,
            line_dash="dot",
            line_color="#" + BBVA_COLORS["sky"],
        )

    fig.update_layout(
        template="plotly_white",
        margin=dict(l=20, r=20, t=20, b=24),
        legend=dict(orientation="h", x=0.01, y=1.08),
        barmode="stack",
        xaxis_title="Día",
        yaxis=dict(title="Comentarios negativos (volumen diario)", rangemode="tozero"),
        yaxis2=dict(
            title="Incidencias Helix del hotspot (puntos)",
            overlaying="y",
            side="right",
            rangemode="tozero",
            showgrid=False,
        ),
        yaxis3=dict(
            title="NPS medio (0-10)",
            overlaying="y",
            side="right",
            anchor="free",
            position=0.92,
            range=[0, 10],
            showgrid=False,
        ),
    )

    if cp_in_window:
        fig.add_annotation(
            xref="paper",
            yref="paper",
            x=0.99,
            y=0.98,
            text=f"Change points visibles: {len(cp_in_window)}",
            showarrow=False,
            font=dict(size=10, color="#" + BBVA_COLORS["muted"]),
            align="right",
        )

    return fig


def _topic_metrics(topic: str, rationale_df: pd.DataFrame) -> dict[str, float]:
    if rationale_df is None or rationale_df.empty or "nps_topic" not in rationale_df.columns:
        return {}
    hit = rationale_df[rationale_df["nps_topic"].astype(str) == str(topic)].head(1)
    if hit.empty:
        return {}
    r = hit.iloc[0]
    return {
        "risk": _safe_float(r.get("nps_points_at_risk", 0.0), default=0.0),
        "recoverable": _safe_float(r.get("nps_points_recoverable", 0.0), default=0.0),
        "priority": _safe_float(r.get("priority", 0.0), default=0.0),
        "confidence": _safe_float(r.get("confidence", 0.0), default=0.0),
        "focus_probability": _safe_float(
            r.get("focus_probability_with_incident", np.nan), default=np.nan
        ),
        "nps_delta_expected": _safe_float(r.get("nps_delta_expected", np.nan), default=np.nan),
        "total_nps_impact": _safe_float(r.get("total_nps_impact", 0.0), default=0.0),
        "causal_score": _safe_float(r.get("causal_score", 0.0), default=0.0),
        "lag_weeks": _safe_float(r.get("best_lag_weeks", np.nan), default=np.nan),
    }


def _set_placeholder_text(
    slide: object, idx: int, text: str, *, font_name: str, size_pt: float
) -> None:
    try:
        placeholder = slide.placeholders[idx]
    except Exception:
        return
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size_pt)
    r.font.bold = True
    r.font.color.rgb = _rgb(BBVA_COLORS["ink"])


def _build_dimension_view_model(
    *,
    dimension: str,
    slide_number: int,
    selected_raw: pd.DataFrame,
    current_source_period: pd.DataFrame,
    baseline_source_period: pd.DataFrame,
    current_label: str,
    baseline_label: str,
) -> DimensionViewModel:
    min_n = max(10, min(50, int(max(len(selected_raw), 1) * 0.02)))
    delta_df = driver_delta_table(
        current_source_period,
        baseline_source_period,
        dimension=dimension,
        min_n=min_n,
    )
    if delta_df.empty:
        delta_df = _driver_change_table(
            selected_raw,
            pd.DataFrame(columns=selected_raw.columns),
            dimension=dimension,
        )
    change_table = select_negative_delta_rows(delta_df, max_rows=EDITORIAL_LIMITS.max_change_rows)
    change_figure = _build_driver_delta_figure(change_table, panel_height_in=4.18)
    source_for_web = current_source_period if not current_source_period.empty else selected_raw
    web_table = _build_web_dimension_table(source_for_web, dimension=dimension)
    opportunities_min_n = max(20, min(200, int(max(len(selected_raw), 1) * 0.02)))
    opportunities = _opportunities_table(
        selected_raw, dimension=dimension, min_n=opportunities_min_n
    )
    opportunities = _prepare_opportunity_chart_df(opportunities)
    return DimensionViewModel(
        dimension=dimension,
        slide_number=slide_number,
        change_df=delta_df,
        change_table_df=change_table,
        change_figure=change_figure,
        web_heatmap_figure=_build_web_heatmap_figure(source_for_web, row_dim=dimension),
        web_table_df=web_table,
        opportunities_df=opportunities,
        opportunities_figure=_build_opportunity_figure(opportunities),
        opportunity_bullets=explain_opportunities(
            opportunities, max_items=EDITORIAL_LIMITS.max_opportunity_bullets
        ),
    )


def _build_causal_scenarios(
    chains: pd.DataFrame,
    *,
    focus_name: str,
) -> list[CausalScenarioViewModel]:
    scenarios: list[CausalScenarioViewModel] = []
    selected = select_causal_scenarios(chains, max_rows=EDITORIAL_LIMITS.max_causal_scenarios)
    for idx, (_, row) in enumerate(selected.iterrows(), start=1):
        raw_kpis = [
            (
                _focus_probability_label(focus_name),
                _fmt_pct_or_nd(row.get("detractor_probability", np.nan)),
                BBVA_COLORS["red"],
            ),
            ("Confianza", _fmt_num_or_nd(row.get("confidence", np.nan)), BBVA_COLORS["green"]),
            (
                "Vínculos validados",
                str(int(_safe_int(row.get("linked_pairs", 0), default=0))),
                BBVA_COLORS["sky"],
            ),
            (
                "Cambio esperado en NPS",
                _fmt_signed_or_nd(row.get("nps_delta_expected", np.nan)),
                BBVA_COLORS["orange"],
            ),
            (
                "Impacto total",
                f"{_fmt_num_or_nd(row.get('total_nps_impact', np.nan))} pts",
                BBVA_COLORS["blue"],
            ),
        ]
        incident_lines = [
            _clean_evidence_excerpt(line, max_len=130)
            for line in _chain_list(row.get("incident_examples"))[
                : EDITORIAL_LIMITS.max_helix_evidence
            ]
        ]
        comment_lines = [
            _clean_evidence_excerpt(line, max_len=135)
            for line in _chain_list(row.get("comment_examples"))[
                : EDITORIAL_LIMITS.max_voc_evidence
            ]
        ]
        helix_records = _chain_incident_records(row.get("incident_records"))
        helix_lines = [
            _clean_evidence_excerpt(
                f"{record.get('incident_id', '')}: {record.get('summary', '')}", max_len=145
            )
            for record in helix_records[: EDITORIAL_LIMITS.max_helix_evidence]
        ]
        if not helix_lines:
            helix_lines = incident_lines[: EDITORIAL_LIMITS.max_helix_evidence]
        scenarios.append(
            CausalScenarioViewModel(
                index=idx,
                row=row,
                kpis=select_nonzero_kpis(
                    raw_kpis, max_items=EDITORIAL_LIMITS.max_visible_causal_kpis
                ),
                incident_lines=incident_lines,
                comment_lines=comment_lines,
                helix_evidence_lines=helix_lines,
            )
        )
    return scenarios


def _build_presentation_context(
    *,
    service_origin: str,
    service_origin_n1: str,
    service_origin_n2: str,
    period_start: date,
    period_end: date,
    focus_name: str,
    overall_weekly: pd.DataFrame,
    story_md: str,
    attribution_df: Optional[pd.DataFrame],
    selected_nps_df: Optional[pd.DataFrame],
    comparison_nps_df: Optional[pd.DataFrame],
    touchpoint_source: str,
    entity_summary_df: Optional[pd.DataFrame],
    entity_summary_kpis: Optional[list[dict[str, str]]],
    broken_journeys_df: Optional[pd.DataFrame],
) -> PresentationContext:
    period_label = f"{_safe_date(period_start)} -> {_safe_date(period_end)}"
    period_days = (pd.Timestamp(period_end) - pd.Timestamp(period_start)).days + 1
    selected_raw = _coerce_nps_records(selected_nps_df)
    compare_raw = _coerce_nps_records(comparison_nps_df)
    current_period, baseline_period = _split_period_frames(
        compare_raw if not compare_raw.empty else selected_raw,
        period_start=period_start,
        period_end=period_end,
    )
    current_source_period, baseline_source_period = _split_source_period_frames(
        comparison_nps_df if comparison_nps_df is not None else selected_nps_df,
        period_start=period_start,
        period_end=period_end,
    )
    if current_source_period.empty and selected_nps_df is not None:
        current_source_period, _ = _split_source_period_frames(
            selected_nps_df,
            period_start=period_start,
            period_end=period_end,
        )
    if selected_raw.empty:
        selected_raw = current_period.copy()
    if selected_raw.empty:
        selected_raw = compare_raw[
            (compare_raw["date"] >= pd.Timestamp(period_start))
            & (compare_raw["date"] <= pd.Timestamp(period_end))
        ].copy()

    daily_signals, _ = _prepare_daily_signals(
        overall_weekly,
        period_start=period_start,
        period_end=period_end,
    )
    daily_mix = _daily_group_mix(selected_raw)
    if daily_mix.empty and not daily_signals.empty:
        daily_mix = daily_signals[["date", "nps_mean", "detractor_rate"]].copy()
        daily_mix["responses"] = 0.0
        daily_mix["passive_rate"] = (1.0 - daily_mix["detractor_rate"]).clip(lower=0.0)
        daily_mix["promoter_rate"] = 0.0
        daily_mix["nps_classic"] = (1.0 - daily_mix["detractor_rate"] * 2.0) * 100.0

    overview = _period_overview(selected_raw)
    text_topics = select_text_clusters(
        _text_topics_table(selected_raw, top_k=10),
        max_clusters=EDITORIAL_LIMITS.max_text_clusters,
    )
    current_label = f"{_safe_date(period_start)} -> {_safe_date(period_end)}"
    baseline_label = (
        f"{_safe_date(baseline_period['date'].min())} -> {_safe_date(baseline_period['date'].max())}"
        if not baseline_period.empty
        else "sin base histórica"
    )
    dimensions = {
        "Palanca": _build_dimension_view_model(
            dimension="Palanca",
            slide_number=5,
            selected_raw=selected_raw,
            current_source_period=current_source_period,
            baseline_source_period=baseline_source_period,
            current_label=current_label,
            baseline_label=baseline_label,
        ),
        "Subpalanca": _build_dimension_view_model(
            dimension="Subpalanca",
            slide_number=6,
            selected_raw=selected_raw,
            current_source_period=current_source_period,
            baseline_source_period=baseline_source_period,
            current_label=current_label,
            baseline_label=baseline_label,
        ),
    }
    del current_label, baseline_label

    chains = attribution_df.copy() if attribution_df is not None else pd.DataFrame()
    causal_entity_summary = (
        entity_summary_df.copy() if entity_summary_df is not None else pd.DataFrame()
    )
    method_spec = get_causal_method_spec(touchpoint_source)
    causal = CausalViewModel(
        touchpoint_source=str(touchpoint_source or "").strip(),
        method_label=method_spec.label,
        method_title=method_spec.navigation_title,
        method_subtitle=method_spec.navigation_subtitle,
        entity_summary_df=causal_entity_summary,
        entity_summary_figure=_build_journey_summary_figure(
            causal_entity_summary, touchpoint_source=touchpoint_source
        ),
        entity_summary_kpis=list(entity_summary_kpis or []),
        journey_table_df=_build_journey_table(
            touchpoint_source=touchpoint_source,
            entity_summary_df=causal_entity_summary,
            broken_journeys_df=broken_journeys_df,
        ),
        scenarios=_build_causal_scenarios(chains, focus_name=focus_name),
    )
    return PresentationContext(
        service_origin=service_origin,
        service_origin_n1=service_origin_n1,
        service_origin_n2=service_origin_n2,
        period_start=period_start,
        period_end=period_end,
        period_label=period_label,
        period_days=period_days,
        focus_name=focus_name,
        overview=overview,
        story_md=story_md,
        selected_raw=selected_raw,
        daily_mix=daily_mix,
        daily_signals=daily_signals,
        overview_figure=_build_overview_figure(selected_nps_df, period_days=period_days),
        text_topics_df=text_topics,
        text_topic_figure=_build_text_topic_figure(text_topics),
        current_label=f"{_safe_date(period_start)} -> {_safe_date(period_end)}",
        baseline_label=(
            f"{_safe_date(baseline_period['date'].min())} -> {_safe_date(baseline_period['date'].max())}"
            if not baseline_period.empty
            else "sin base histórica"
        ),
        dimensions=dimensions,
        causal=causal,
    )


def _add_cover_slide(
    prs: Presentation,
    *,
    service_origin: str,
    service_origin_n1: str,
    service_origin_n2: str,
    period_start: date,
    period_end: date,
    overview: dict[str, object],
    story_md: str,
) -> None:
    del overview, story_md
    slide = _new_slide(prs, kind="cover")
    _add_bg(slide, BBVA_COLORS["bg_dark"])
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.24), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(BBVA_COLORS["sky"])
    accent.line.fill.background()

    title = "Análisis NPS térmico y causalidad"
    subtitle = f"{service_origin} · {service_origin_n1}".strip(" ·")
    if service_origin_n2:
        subtitle = f"{subtitle} · {service_origin_n2}".strip(" ·")
    period_label = f"{_safe_date(period_start)} -> {_safe_date(period_end)}"
    subtitle = f"{subtitle} · {period_label}".strip(" ·")

    hero = slide.shapes.add_textbox(Inches(0.78), Inches(3.04), Inches(11.85), Inches(0.86))
    htf = hero.text_frame
    _configure_text_frame(htf)
    htf.clear()
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    hr.font.name = BBVA_FONT_DISPLAY
    hr.font.size = Pt(38)
    hr.font.bold = True
    hr.font.color.rgb = _rgb(BBVA_COLORS["white"])

    sub = slide.shapes.add_textbox(Inches(0.82), Inches(3.94), Inches(10.80), Inches(0.50))
    stf = sub.text_frame
    _configure_text_frame(stf)
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.name = BBVA_FONT_BODY
    sr.font.size = Pt(15)
    sr.font.color.rgb = _rgb("C7D3EA")


def _add_nps_section_cover_slide(prs: Presentation, *, context: PresentationContext) -> None:
    slide = _new_slide(prs, kind="cover")
    _add_bg(slide, BBVA_COLORS["bg_dark"])
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.24), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(BBVA_COLORS["sky"])
    accent.line.fill.background()

    eyebrow = slide.shapes.add_textbox(Inches(0.78), Inches(0.72), Inches(4.8), Inches(0.36))
    tf = eyebrow.text_frame
    _configure_text_frame(tf)
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "Bloque 1 · Experiencia térmica"
    r.font.name = BBVA_FONT_MEDIUM
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = _rgb(BBVA_COLORS["sky"])

    title = slide.shapes.add_textbox(Inches(0.78), Inches(1.18), Inches(7.4), Inches(1.25))
    ttf = title.text_frame
    _configure_text_frame(ttf)
    ttf.clear()
    tr = ttf.paragraphs[0].add_run()
    tr.text = "NPS térmico"
    tr.font.name = BBVA_FONT_DISPLAY
    tr.font.size = Pt(48)
    tr.font.bold = True
    tr.font.color.rgb = _rgb(BBVA_COLORS["white"])

    subtitle = slide.shapes.add_textbox(Inches(0.82), Inches(2.46), Inches(7.2), Inches(0.72))
    stf = subtitle.text_frame
    _configure_text_frame(stf)
    stf.clear()
    sr = stf.paragraphs[0].add_run()
    sr.text = f"{context.service_origin} · {context.service_origin_n1} · {context.period_label}"
    sr.font.name = BBVA_FONT_BODY
    sr.font.size = Pt(16)
    sr.font.color.rgb = _rgb("C7D3EA")

    summary = _cover_summary_lines(context.overview, context.story_md)[:3] or [
        "El bloque ordena la señal del periodo desde evolución, comentario y deterioros por dimensión.",
    ]
    _add_bullet_lines(
        slide,
        left=0.82,
        top=3.42,
        width=7.2,
        height=2.50,
        title="Lectura ejecutiva",
        lines=summary,
        accent=BBVA_COLORS["sky"],
        body_font_size_pt=13.8,
    )
    _add_stat_card(
        slide,
        left=8.70,
        top=1.10,
        width=3.40,
        height=1.12,
        label="NPS medio",
        value=_fmt_num_or_nd(context.overview.get("nps_mean", np.nan), decimals=1),
        accent=BBVA_COLORS["green"],
        hint="Escala 0-10",
    )
    _add_stat_card(
        slide,
        left=8.70,
        top=2.52,
        width=3.40,
        height=1.12,
        label="Detractores",
        value=_fmt_pct_or_nd(context.overview.get("detractor_rate", np.nan)),
        accent=BBVA_COLORS["red"],
        hint="Peso sobre respuestas",
    )
    _add_stat_card(
        slide,
        left=8.70,
        top=3.94,
        width=3.40,
        height=1.12,
        label="% promotores",
        value=_fmt_pct_or_nd(context.overview.get("promoter_rate", np.nan)),
        accent=BBVA_COLORS["green"],
        hint="Valoraciones >= 9",
    )
    _add_stat_card(
        slide,
        left=8.70,
        top=5.36,
        width=3.40,
        height=1.12,
        label="Comentarios",
        value=_fmt_count_or_nd(context.overview.get("comments", 0)),
        accent=BBVA_COLORS["blue"],
        hint="Base útil analizada",
    )


def _add_dimension_change_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    view_model: DimensionViewModel,
) -> None:
    dimension = view_model.dimension
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title=f"{view_model.slide_number}. Qué ha cambiado en {dimension}",
        subtitle=(
            f"Periodo actual frente a la base histórica anterior · actual {context.current_label} · "
            f"base {context.baseline_label}"
        ),
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=7.78,
        height=5.42,
        title=f"Deterioros principales en {dimension}",
    )
    _figure_in_panel(
        slide,
        figure=view_model.change_figure,
        left=0.88,
        top=1.84,
        width=7.24,
        height=4.76,
        empty_note=f"No hay base histórica suficiente para comparar {dimension.lower()}.",
        target_ppi=178,
    )
    rows = [
        [
            _clip(row.value, 34),
            _fmt_count_or_nd(row.n_current),
            _fmt_num_or_nd(row.nps_current, decimals=1),
            _fmt_signed_or_nd(row.delta_nps, decimals=1),
        ]
        for row in view_model.change_table_df.head(EDITORIAL_LIMITS.max_change_rows).itertuples()
    ]
    _add_compact_table(
        slide,
        left=8.70,
        top=1.48,
        width=3.98,
        title="Mayor deterioro",
        headers=[dimension, "n", "NPS", "Δ"],
        rows=rows or [["Sin datos", "-", "-", "-"]],
        row_height=0.48,
        col_width_ratios=[2.4, 0.7, 0.8, 0.7],
        clip_lengths=[34, 8, 8, 8],
        font_size_pt=10.8,
        max_rows=EDITORIAL_LIMITS.max_change_rows,
    )
    _add_bullet_lines(
        slide,
        left=8.70,
        top=5.64,
        width=3.98,
        height=1.26,
        title="Criterio de recorte",
        lines=["Se muestran los mayores deterioros de NPS con volumen defendible para comité."],
        accent=BBVA_COLORS["line"],
        body_font_size_pt=10.8,
    )


def _add_web_pain_dimension_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    view_model: DimensionViewModel,
    slide_number: int,
) -> None:
    dimension = view_model.dimension
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title=f"{slide_number}. Dónde duele en la Web · {dimension}",
        subtitle=f"NPS del canal Web por {dimension.lower()} dentro del periodo analizado · {context.period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=8.10,
        height=5.42,
        title=f"Mapa de dolor Web por {dimension}",
    )
    _figure_in_panel(
        slide,
        figure=view_model.web_heatmap_figure,
        left=0.86,
        top=1.86,
        width=7.66,
        height=4.76,
        empty_note=f"No hay señal suficiente para mostrar {dimension.lower()} en el canal Web.",
        target_ppi=178,
    )
    rows = [
        [
            _clip(row.value, 36),
            _fmt_count_or_nd(row.n),
            _fmt_num_or_nd(row.nps, decimals=1),
            _fmt_pct_or_nd(row.detractor_rate),
        ]
        for row in view_model.web_table_df.head(EDITORIAL_LIMITS.max_web_rows).itertuples()
    ]
    _add_compact_table(
        slide,
        left=9.02,
        top=1.48,
        width=3.66,
        title="Focos Web",
        headers=[dimension, "n", "NPS", "% det."],
        rows=rows or [["Sin datos", "-", "-", "-"]],
        row_height=0.43,
        col_width_ratios=[2.3, 0.7, 0.8, 0.9],
        clip_lengths=[34, 8, 8, 8],
        font_size_pt=10.2,
        max_rows=EDITORIAL_LIMITS.max_web_rows,
    )


def _add_opportunity_dimension_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    view_model: DimensionViewModel,
    slide_number: int,
) -> None:
    dimension = view_model.dimension
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title=f"{slide_number}. Oportunidades priorizadas · {dimension}",
        subtitle=f"Ranking por impacto potencial y solidez de evidencia · {context.period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=12.02,
        height=4.52,
        title=f"Impacto estimado en {dimension}",
    )
    _figure_in_panel(
        slide,
        figure=view_model.opportunities_figure,
        left=0.86,
        top=1.84,
        width=11.62,
        height=3.82,
        empty_note=f"No se identificaron oportunidades robustas para {dimension.lower()} con el umbral actual.",
        target_ppi=176,
    )
    _add_bullet_lines(
        slide,
        left=0.66,
        top=6.12,
        width=12.02,
        height=0.88,
        title="",
        lines=view_model.opportunity_bullets,
        accent=BBVA_COLORS["line"],
        body_font_size_pt=12.0,
    )


def _add_overview_slide(
    prs: Presentation,
    *,
    service_origin: str,
    service_origin_n1: str,
    period_label: str,
    period_end: date,
    overview: dict[str, object],
    selected_nps_df: Optional[pd.DataFrame],
    period_days: int,
    overview_figure: Optional[go.Figure] = None,
) -> None:
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="1. Evolución del NPS del periodo",
        subtitle=f"{service_origin} · {service_origin_n1} · {period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=9.00,
        height=5.42,
        title="NPS clásico y peso detractor",
    )
    _figure_in_panel(
        slide,
        figure=(
            overview_figure
            if overview_figure is not None
            else _build_overview_figure(selected_nps_df, period_days=period_days)
        ),
        left=0.82,
        top=1.84,
        width=8.54,
        height=4.84,
        empty_note="No hay suficiente señal diaria para construir la evolución del periodo.",
        target_ppi=180,
    )

    month_label = _month_label_es(period_end).title()
    trend_lines = [
        f"El periodo arranca con NPS clásico {_fmt_num_or_nd(overview.get('start_classic', np.nan))} y termina en {_fmt_num_or_nd(overview.get('end_classic', np.nan))}.",
        f"El peso detractor pasa de {_fmt_pct_or_nd(overview.get('start_detr', np.nan))} a {_fmt_pct_or_nd(overview.get('end_detr', np.nan))}.",
        "NPS clásico = promotores menos detractores; se usa para seguir la señal neta del periodo.",
    ]
    _add_bullet_lines(
        slide,
        left=9.92,
        top=1.48,
        width=2.76,
        height=5.42,
        title=f"As-Is {month_label}",
        lines=trend_lines,
        accent=BBVA_COLORS["orange"],
        body_font_size_pt=13.0,
    )


def _add_deep_dive_slide(
    prs: Presentation,
    *,
    period_label: str,
    text_topics_df: pd.DataFrame,
    topic_figure: Optional[go.Figure] = None,
) -> None:
    topic_fig = (
        topic_figure if topic_figure is not None else _build_text_topic_figure(text_topics_df)
    )

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="2. Qué han dicho los clientes",
        subtitle=f"Temas más repetidos en los comentarios del periodo · {period_label}",
    )
    _panel(slide, left=0.66, top=1.48, width=12.02, height=5.42, title="Top temas del periodo")
    _figure_in_panel(
        slide,
        figure=topic_fig,
        left=0.82,
        top=1.82,
        width=11.62,
        height=3.20,
        empty_note="No hay suficiente volumen textual para construir el top 10.",
        target_ppi=176,
    )

    table_rows = [
        [
            str(row.cluster_id),
            _fmt_count_or_nd(row.n),
            str(row.top_terms_txt),
            str(row.example_txt),
        ]
        for row in text_topics_df.head(3).itertuples()
    ]
    _add_compact_table(
        slide,
        left=0.82,
        top=5.02,
        width=11.62,
        title="",
        headers=["cluster_id", "n", "top_terms", "examples"],
        rows=table_rows or [["-", "-", "Sin datos", "Sin ejemplos"]],
        row_height=0.36,
        col_width_ratios=[0.8, 0.8, 4.5, 4.3],
        clip_lengths=[8, 8, 58, 58],
        font_size_pt=10.4,
        max_rows=3,
    )


def _add_topic_timing_slide(
    prs: Presentation,
    *,
    period_label: str,
    period_days: int,
    selected_nps_df: Optional[pd.DataFrame],
) -> None:
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    source_df = selected_nps_df.copy() if selected_nps_df is not None else pd.DataFrame()
    _add_header(
        slide,
        title="2. Cuándo y cómo lo dicen",
        subtitle=f"Volumen diario de respuestas y reparto de promotores, pasivos y detractores · {period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=12.02,
        height=2.38,
        title="Cuándo lo dicen",
    )
    _figure_in_panel(
        slide,
        figure=chart_daily_volume(source_df, get_theme("light"), days=max(int(period_days), 1)),
        left=0.86,
        top=1.80,
        width=11.62,
        height=1.94,
        empty_note="No hay señal suficiente para mostrar el volumen diario del periodo.",
        target_ppi=180,
    )
    _panel(slide, left=0.66, top=4.02, width=12.02, height=2.88, title="Cómo lo dicen")
    _figure_in_panel(
        slide,
        figure=chart_daily_mix_business(
            source_df, get_theme("light"), days=max(int(period_days), 1)
        ),
        left=0.86,
        top=4.34,
        width=11.62,
        height=2.30,
        empty_note="No hay señal suficiente para la distribución diaria por grupo.",
        target_ppi=180,
    )


def _add_change_vs_past_slide(
    prs: Presentation,
    *,
    period_label: str,
    current_label: str,
    baseline_label: str,
    current_source_df: pd.DataFrame,
    baseline_source_df: pd.DataFrame,
) -> None:
    def _delta_figure(
        *,
        dimension: str,
        panel_height_in: float,
    ) -> Optional[go.Figure]:
        delta_df = driver_delta_table(
            current_source_df,
            baseline_source_df,
            dimension=dimension,
            min_n=50,
        )
        fig = chart_driver_delta(delta_df, get_theme("light"))
        if fig is None or delta_df.empty:
            return fig

        plot_df = delta_df.head(12).copy()
        label_count = len(plot_df)
        max_len = int(plot_df["value"].astype(str).str.len().max() or 0)
        wrap_width = 18 if max_len >= 18 else 14
        left_margin = 300 if max_len >= 18 else 240
        y_font_size = 44 if label_count <= 6 else 40 if label_count <= 8 else 34

        labels = [
            _wrap_label(value, width=wrap_width, max_lines=2, joiner="<br>")
            for value in plot_df["value"].astype(str).tolist()
        ]
        with contextlib.suppress(Exception):
            fig.data[0].y = labels
        fig.update_yaxes(
            tickfont=dict(size=y_font_size),
            automargin=True,
        )
        fig.update_xaxes(
            title_text="Delta NPS",
            tickfont=dict(size=18),
            title_font=dict(size=20),
            nticks=5,
        )
        fig.update_layout(
            margin=dict(
                l=left_margin,
                r=40,
                t=24,
                b=52 if panel_height_in <= 2.0 else 48,
            )
        )
        return fig

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="3. Qué ha cambiado respecto al pasado",
        subtitle=f"Periodo actual frente a la base histórica anterior · actual {current_label} · base {baseline_label}",
    )
    _panel(slide, left=0.66, top=1.48, width=12.02, height=2.38, title="Palanca")
    _figure_in_panel(
        slide,
        figure=_delta_figure(dimension="Palanca", panel_height_in=1.92),
        left=0.86,
        top=1.80,
        width=11.40,
        height=1.92,
        empty_note="No hay base histórica suficiente para comparar por palanca.",
        target_ppi=180,
    )
    _panel(slide, left=0.66, top=4.02, width=12.02, height=2.88, title="Subpalanca")
    _figure_in_panel(
        slide,
        figure=_delta_figure(dimension="Subpalanca", panel_height_in=2.30),
        left=0.86,
        top=4.34,
        width=11.40,
        height=2.30,
        empty_note="No hay base histórica suficiente para comparar por subpalanca.",
        target_ppi=180,
    )


def _add_pain_by_group_slide(
    prs: Presentation,
    *,
    period_label: str,
    selected_nps_df: Optional[pd.DataFrame],
) -> None:
    source_df = selected_nps_df.copy() if selected_nps_df is not None else pd.DataFrame()

    def _web_heatmap_figure(
        *,
        row_dim: str,
        show_scale: bool,
        x_tick_label: str,
        include_all_rows: bool = False,
    ) -> Optional[go.Figure]:
        required = {row_dim, "Canal", "NPS"}
        if source_df.empty or not required.issubset(set(source_df.columns)):
            return None

        chart_df = source_df.dropna(subset=[row_dim, "Canal", "NPS"]).copy()
        if chart_df.empty:
            return None
        chart_df[row_dim] = chart_df[row_dim].astype(str).str.strip()
        chart_df["Canal"] = chart_df["Canal"].astype(str).str.strip()
        chart_df = chart_df[
            chart_df[row_dim].ne("") & chart_df["Canal"].str.casefold().eq("web")
        ].copy()
        if chart_df.empty:
            return None

        fig = chart_cohort_heatmap(
            chart_df,
            get_theme("light"),
            row_dim=row_dim,
            col_dim="Canal",
            min_n=1 if include_all_rows else 30,
        )
        if fig is None:
            return None

        agg = chart_df.groupby([row_dim, "Canal"], as_index=False).agg(n=("NPS", "size")).copy()
        if agg.empty:
            return None

        row_values = (
            agg[row_dim].astype(str).drop_duplicates().tolist()
            if include_all_rows
            else agg.loc[agg["n"] >= 30, row_dim].astype(str).drop_duplicates().tolist()
        )
        if not row_values:
            return None

        label_count = len(row_values)
        max_len = int(pd.Series(row_values).astype(str).str.len().max() or 0)
        y_font_size = (
            44 if label_count <= 5 else 40 if label_count <= 7 else 36 if label_count <= 10 else 31
        )
        left_margin = (
            372 if max_len >= 40 else 352 if max_len >= 34 else 334 if max_len >= 28 else 316
        )
        label_map = {
            _compact_axis_label(value, width=16, max_lines=2, max_chars=30): str(value)
            for value in row_values
        }
        with contextlib.suppress(Exception):
            rendered_y_values = [
                label_map.get(str(value), str(value).replace("<br>", " "))
                for value in list(fig.data[0].y)
            ]
            fig.data[0].y = rendered_y_values
        rendered_y_values = list(
            dict.fromkeys([str(v) for v in list(getattr(fig.data[0], "y", []))])
        )

        if include_all_rows:
            with contextlib.suppress(Exception):
                n_by_row = (
                    agg.groupby(row_dim, as_index=False)["n"]
                    .sum()
                    .set_index(row_dim)["n"]
                    .to_dict()
                )
                z_values = np.array(fig.data[0].z, dtype=float)
                for idx, y_value in enumerate(list(fig.data[0].y)):
                    if int(n_by_row.get(str(y_value), 0)) < 30:
                        z_values[idx, :] = np.nan
                fig.data[0].z = z_values
        fig.update_yaxes(
            title_text="",
            categoryorder="array",
            categoryarray=rendered_y_values,
            tickmode="array",
            tickvals=rendered_y_values,
            ticktext=rendered_y_values,
            tickfont=dict(size=y_font_size, family=BBVA_FONT_BODY, color="#" + BBVA_COLORS["ink"]),
            automargin=True,
        )
        fig.update_xaxes(
            title_text="",
            tickangle=0,
            tickmode="array",
            tickvals=["Web"],
            ticktext=[x_tick_label],
            tickfont=dict(size=14),
            automargin=True,
        )
        fig.update_layout(
            margin=dict(
                l=left_margin,
                r=88 if show_scale else 20,
                t=14,
                b=24,
            )
        )
        fig.update_coloraxes(showscale=show_scale)
        if show_scale:
            fig.update_coloraxes(
                colorbar=dict(
                    title="NPS",
                    tickmode="array",
                    tickvals=[0, 2, 6, 8, 10],
                    len=0.72,
                    y=0.5,
                    thickness=12,
                )
            )
        return fig

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="4. Dónde duele según el tipo de cliente",
        subtitle=f"NPS del canal Web por eje de experiencia dentro del periodo analizado · {period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=5.88,
        height=5.42,
        title="Palanca · canal Web",
    )
    _figure_in_panel(
        slide,
        figure=_web_heatmap_figure(
            row_dim="Palanca",
            show_scale=False,
            x_tick_label="Web: Por palanca",
            include_all_rows=True,
        ),
        left=0.82,
        top=1.86,
        width=5.14,
        height=4.92,
        empty_note="No hay señal suficiente para mostrar la matriz Palanca del canal Web.",
        target_ppi=176,
    )
    _panel(
        slide,
        left=6.80,
        top=1.48,
        width=5.88,
        height=5.42,
        title="Subpalanca · canal Web",
    )
    _figure_in_panel(
        slide,
        figure=_web_heatmap_figure(
            row_dim="Subpalanca",
            show_scale=True,
            x_tick_label="Web: Por subpalanca",
        ),
        left=6.96,
        top=1.86,
        width=5.48,
        height=4.92,
        empty_note="No hay señal suficiente para mostrar la matriz Subpalanca del canal Web.",
        target_ppi=176,
    )


def _add_gap_slide(
    prs: Presentation,
    *,
    period_label: str,
    palanca_gap_df: pd.DataFrame,
    subpalanca_gap_df: pd.DataFrame,
) -> None:
    def _gap_figure(gap_df: pd.DataFrame, *, panel_height_in: float) -> Optional[go.Figure]:
        fig = chart_driver_bar(gap_df, get_theme("light"), top_k=10)
        if fig is None or gap_df.empty:
            return fig

        plot_df = gap_df.head(10).copy()
        max_len = int(plot_df["value"].astype(str).str.len().max() or 0)
        wrap_width = 24 if max_len >= 24 else 22 if max_len >= 18 else 20
        max_chars = 44 if max_len >= 24 else 38 if max_len >= 18 else 34
        y_font_size = 34 if len(plot_df) <= 6 else 31 if len(plot_df) <= 8 else 28
        left_margin = (
            340 if max_len >= 30 else 314 if max_len >= 24 else 288 if max_len >= 18 else 262
        )
        label_map = {
            value: _compact_axis_label(value, width=wrap_width, max_lines=2, max_chars=max_chars)
            for value in plot_df["value"].astype(str).tolist()
        }
        with contextlib.suppress(Exception):
            fig.data[0].y = [label_map.get(str(value), str(value)) for value in list(fig.data[0].y)]
        fig.update_yaxes(
            title_text="",
            tickfont=dict(
                size=y_font_size, family=BBVA_FONT_MEDIUM, color="#" + BBVA_COLORS["ink"]
            ),
            automargin=True,
        )
        fig.update_xaxes(
            title_text="Gap NPS",
            tickfont=dict(size=16),
            title_font=dict(size=17),
            nticks=4,
        )
        fig.update_layout(
            margin=dict(
                l=left_margin,
                r=30,
                t=18,
                b=42 if panel_height_in <= 2.0 else 38,
            )
        )
        return fig

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="5. Casos más alejados del promedio",
        subtitle=f"Top de casos con peor diferencia frente al NPS medio general · {period_label}",
    )
    _panel(slide, left=0.66, top=1.48, width=8.70, height=2.38, title="Palanca")
    _figure_in_panel(
        slide,
        figure=_gap_figure(palanca_gap_df, panel_height_in=2.00),
        left=0.82,
        top=1.72,
        width=8.08,
        height=2.10,
        empty_note="No hay suficiente señal para el ranking de brechas por palanca.",
        target_ppi=174,
    )

    palanca_lines = [
        f"{idx + 1}. {_clip(row.value, 30)} · n={_fmt_count_or_nd(row.n)} · "
        f"NPS {_fmt_num_or_nd(row.nps)} · gap {_fmt_signed_or_nd(row.gap_vs_overall, decimals=1)}"
        for idx, row in enumerate(palanca_gap_df.head(5).itertuples())
    ]
    _add_bullet_lines(
        slide,
        left=9.52,
        top=1.48,
        width=3.16,
        height=2.38,
        title="",
        lines=palanca_lines,
        body_font_size_pt=11.0,
    )

    _panel(slide, left=0.66, top=4.02, width=8.70, height=2.88, title="Subpalanca")
    _figure_in_panel(
        slide,
        figure=_gap_figure(subpalanca_gap_df, panel_height_in=2.44),
        left=0.82,
        top=4.22,
        width=8.08,
        height=2.56,
        empty_note="No hay suficiente señal para el ranking de brechas por subpalanca.",
        target_ppi=174,
    )
    subpalanca_lines = [
        f"{idx + 1}. {_clip(row.value, 30)} · n={_fmt_count_or_nd(row.n)} · "
        f"NPS {_fmt_num_or_nd(row.nps)} · gap {_fmt_signed_or_nd(row.gap_vs_overall, decimals=1)}"
        for idx, row in enumerate(subpalanca_gap_df.head(5).itertuples())
    ]
    _add_bullet_lines(
        slide,
        left=9.52,
        top=4.02,
        width=3.16,
        height=2.88,
        title="",
        lines=subpalanca_lines,
        body_font_size_pt=11.0,
    )


def _add_opportunity_slide(
    prs: Presentation,
    *,
    period_label: str,
    opportunities_df: pd.DataFrame,
) -> None:
    def _opportunity_figure(opp_df: pd.DataFrame) -> Optional[go.Figure]:
        fig = chart_opportunities_bar(opp_df, get_theme("light"), top_k=10)
        if fig is None or opp_df.empty:
            return fig

        plot_df = (
            opp_df.sort_values(["potential_uplift", "confidence"], ascending=[False, False])
            .head(10)
            .copy()
        )
        label_count = len(plot_df)
        label_lengths = (
            plot_df["label"].astype(str).str.replace("<br>", " ", regex=False).str.len()
            if "label" in plot_df.columns
            else pd.Series(dtype=float)
        )
        max_len = int(label_lengths.max() or 0)
        y_font_size = 28 if label_count <= 5 else 26 if label_count <= 7 else 23
        left_margin = 248 if max_len >= 26 else 220 if max_len >= 20 else 188

        uplift = pd.to_numeric(plot_df.get("potential_uplift"), errors="coerce")
        text_values = [
            _fmt_signed_or_nd(value, decimals=1) if np.isfinite(value) else ""
            for value in uplift.tolist()
        ]
        with contextlib.suppress(Exception):
            fig.data[0].text = text_values
        with contextlib.suppress(Exception):
            fig.data[0].textposition = "outside"
        with contextlib.suppress(Exception):
            fig.data[0].cliponaxis = False
        with contextlib.suppress(Exception):
            fig.data[0].textfont.size = 18

        fig.update_yaxes(
            title_text="",
            tickfont=dict(size=y_font_size),
            automargin=True,
        )
        fig.update_xaxes(
            title_text="Impacto estimado",
            tickfont=dict(size=18),
            title_font=dict(size=19),
            nticks=5,
        )
        fig.update_layout(
            margin=dict(
                l=left_margin,
                r=54,
                t=18,
                b=46,
            ),
            bargap=0.34,
        )
        return fig

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title="6. Oportunidades a priorizar",
        subtitle=f"Ranking de oportunidades por impacto potencial y solidez de evidencia · {period_label}",
    )
    opp_chart_df = opportunities_df.copy()
    if not opp_chart_df.empty:
        dimensions = [
            str(value).strip()
            for value in opp_chart_df.get("dimension", pd.Series(dtype=str)).dropna().tolist()
            if str(value).strip()
        ]
        shared_dimension = dimensions[0] if len(set(dimensions)) == 1 else ""

        def _opp_label(row: pd.Series) -> str:
            value = str(row.get("value", "")).strip()
            dimension = str(row.get("dimension", "")).strip()
            base = (
                value
                if shared_dimension and dimension == shared_dimension
                else f"{dimension} · {value}".strip(" ·")
            )
            width = 18 if len(base) >= 20 else 16
            max_chars = 28 if shared_dimension else 32
            return _compact_axis_label(base, width=width, max_lines=2, max_chars=max_chars)

        opp_chart_df["label"] = opp_chart_df.apply(_opp_label, axis=1)
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=12.02,
        height=4.38,
        title="Ranking por impacto estimado x confianza",
    )
    _figure_in_panel(
        slide,
        figure=_opportunity_figure(opp_chart_df),
        left=0.86,
        top=1.84,
        width=11.62,
        height=3.60,
        empty_note="No se identificaron oportunidades robustas con el umbral actual.",
    )
    lines = explain_opportunities(opp_chart_df, max_items=5)
    _add_bullet_lines(
        slide,
        left=0.66,
        top=6.02,
        width=12.02,
        height=1.10,
        title="",
        lines=lines,
        accent=BBVA_COLORS["line"],
        body_font_size_pt=12.5,
    )


def _add_causal_timeline_slide(
    prs: Presentation,
    *,
    period_label: str,
    daily_mix: pd.DataFrame,
    overall_daily: pd.DataFrame,
    nps_points_at_risk: float,
    nps_points_recoverable: float,
    top3_incident_share: float,
    touchpoint_source: str,
) -> None:
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    method_spec = get_causal_method_spec(touchpoint_source)
    incidents_total = int(
        pd.to_numeric(overall_daily.get("incidents", 0.0), errors="coerce").fillna(0.0).sum()
    )
    detractor_avg = float(
        pd.to_numeric(
            overall_daily.get(
                "focus_rate",
                overall_daily.get("detractor_rate", daily_mix.get("detractor_rate", 0.0)),
            ),
            errors="coerce",
        )
        .fillna(0.0)
        .mean()
    )
    _add_header(
        slide,
        title="7. Situación del periodo",
        subtitle=f"{method_spec.situation_subtitle} · {period_label}",
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=8.15,
        height=5.42,
        title="Evolución diaria de incidencias y grupos NPS",
    )
    _figure_in_panel(
        slide,
        figure=_causal_daily_timeline_fig(daily_mix, overall_daily),
        left=0.82,
        top=1.86,
        width=7.83,
        height=4.92,
        empty_note="No hay cobertura diaria suficiente para el timeline causal.",
        target_ppi=180,
    )
    _panel(
        slide,
        left=9.02,
        top=1.48,
        width=3.66,
        height=5.42,
        title="Cómo leerlo",
        border=BBVA_COLORS["red"],
    )
    text_box = slide.shapes.add_textbox(Inches(9.22), Inches(1.98), Inches(3.26), Inches(1.34))
    text_tf = text_box.text_frame
    _configure_text_frame(text_tf)
    text_tf.word_wrap = True
    text_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_tf.clear()
    text_p = text_tf.paragraphs[0]
    text_p.alignment = PP_ALIGN.LEFT
    text_r = text_p.add_run()
    text_r.text = method_spec.situation_note
    text_r.font.name = BBVA_FONT_BODY
    text_r.font.size = Pt(12.0)
    text_r.font.color.rgb = _rgb(BBVA_COLORS["muted"])
    _add_stat_card(
        slide,
        left=9.18,
        top=3.62,
        width=1.60,
        height=1.12,
        label="Incidencias",
        value=_fmt_count_or_nd(incidents_total),
        accent=BBVA_COLORS["blue"],
    )
    _add_stat_card(
        slide,
        left=10.92,
        top=3.62,
        width=1.60,
        height=1.12,
        label="% detractores",
        value=_fmt_pct_or_nd(detractor_avg, decimals=1),
        accent=BBVA_COLORS["red"],
    )
    _add_stat_card(
        slide,
        left=9.18,
        top=4.94,
        width=1.60,
        height=1.12,
        label="NPS en riesgo",
        value=f"{_fmt_num_or_nd(nps_points_at_risk)} pts",
        accent=BBVA_COLORS["orange"],
    )
    _add_stat_card(
        slide,
        left=10.92,
        top=4.94,
        width=1.60,
        height=1.12,
        label="NPS recuperable",
        value=f"{_fmt_num_or_nd(nps_points_recoverable)} pts",
        accent=BBVA_COLORS["green"],
    )


def _add_journeys_summary_slide(
    prs: Presentation,
    *,
    period_label: str,
    touchpoint_source: str,
    entity_summary_df: pd.DataFrame,
    entity_summary_kpis: list[dict[str, str]],
    entity_summary_figure: Optional[go.Figure] = None,
    journey_table_df: Optional[pd.DataFrame] = None,
) -> None:
    method_spec = get_causal_method_spec(touchpoint_source)

    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    journey_title = (
        "12. Journeys rotos"
        if str(touchpoint_source or "").strip() == TOUCHPOINT_SOURCE_BROKEN_JOURNEYS
        else "12. Journeys de detracción"
    )
    _add_header(
        slide,
        title=journey_title,
        subtitle=f"{method_spec.navigation_subtitle} · {period_label}",
    )
    for index, metric in enumerate(entity_summary_kpis[:3]):
        _add_stat_card(
            slide,
            left=0.66 + index * 3.96,
            top=1.48,
            width=3.78,
            height=1.10,
            label=str(metric.get("label", "")).strip(),
            value=str(metric.get("value", "")).strip(),
            accent=(
                BBVA_COLORS["blue"]
                if index == 0
                else BBVA_COLORS["orange"] if index == 1 else BBVA_COLORS["green"]
            ),
        )
    _panel(slide, left=0.66, top=2.72, width=5.58, height=4.18, title=method_spec.chart_title)
    _figure_in_panel(
        slide,
        figure=(
            entity_summary_figure
            if entity_summary_figure is not None
            else _build_journey_summary_figure(
                entity_summary_df, touchpoint_source=touchpoint_source
            )
        ),
        left=0.86,
        top=3.08,
        width=5.14,
        height=3.48,
        empty_note=method_spec.table_empty_message,
        target_ppi=170,
    )
    table_df = journey_table_df if journey_table_df is not None else pd.DataFrame()
    rows = [
        [
            _clip(row.journey, 58),
            _clip(row.touchpoint, 22),
            _fmt_count_or_nd(row.links),
            _fmt_num_or_nd(row.confidence, decimals=2),
        ]
        for row in table_df.head(EDITORIAL_LIMITS.max_journey_rows).itertuples()
    ]
    _add_compact_table(
        slide,
        left=6.52,
        top=2.72,
        width=6.16,
        title=method_spec.table_title,
        headers=["Journey", "Touchpoint", "Links", "Conf."],
        rows=rows or [["Sin evidencia suficiente", "-", "-", "-"]],
        row_height=0.46,
        col_width_ratios=[3.3, 1.2, 0.7, 0.7],
        clip_lengths=[58, 22, 8, 8],
        font_size_pt=10.0,
        max_rows=EDITORIAL_LIMITS.max_journey_rows,
    )


def _add_causal_section_cover_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    nps_points_at_risk: float,
    nps_points_recoverable: float,
    top3_incident_share: float,
) -> None:
    del nps_points_at_risk, nps_points_recoverable, top3_incident_share
    slide = _new_slide(prs, kind="cover")
    _add_bg(slide, BBVA_COLORS["bg_dark"])
    method = context.causal
    eyebrow = slide.shapes.add_textbox(Inches(0.78), Inches(0.72), Inches(5.4), Inches(0.34))
    tf = eyebrow.text_frame
    _configure_text_frame(tf)
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "Bloque 2 · Narrativa causal seleccionada"
    r.font.name = BBVA_FONT_MEDIUM
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = _rgb(BBVA_COLORS["sky"])

    title = slide.shapes.add_textbox(Inches(0.78), Inches(1.18), Inches(8.0), Inches(1.28))
    ttf = title.text_frame
    _configure_text_frame(ttf)
    ttf.clear()
    tr = ttf.paragraphs[0].add_run()
    tr.text = f"11. Narrativa causal · {method.method_label}"
    tr.font.name = BBVA_FONT_DISPLAY
    tr.font.size = Pt(38)
    tr.font.bold = True
    tr.font.color.rgb = _rgb(BBVA_COLORS["white"])

    subtitle = slide.shapes.add_textbox(Inches(0.82), Inches(2.54), Inches(7.8), Inches(0.84))
    stf = subtitle.text_frame
    _configure_text_frame(stf)
    stf.clear()
    sr = stf.paragraphs[0].add_run()
    sr.text = f"{method.method_subtitle} · {context.period_label}"
    sr.font.name = BBVA_FONT_BODY
    sr.font.size = Pt(15)
    sr.font.color.rgb = _rgb("C7D3EA")

    flow_box = _panel(
        slide,
        left=0.82,
        top=3.62,
        width=7.48,
        height=1.26,
        title="Hipótesis causal defendible",
        subtitle=get_causal_method_spec(method.touchpoint_source).flow,
        fill=BBVA_COLORS["white"],
        border=BBVA_COLORS["sky"],
        title_size=13,
    )
    del flow_box
    _add_stat_card(
        slide,
        left=8.70,
        top=2.82,
        width=3.40,
        height=1.34,
        label="Escenarios",
        value=_fmt_count_or_nd(len(method.scenarios)),
        accent=BBVA_COLORS["blue"],
        hint="Casos priorizados",
    )


def _add_causal_analysis_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    scenario: CausalScenarioViewModel,
) -> None:
    row = scenario.row
    method_spec = get_causal_method_spec(
        str(row.get("presentation_mode", "") or context.causal.touchpoint_source).strip()
    )
    title = _clip(row.get("nps_topic", f"Escenario {scenario.index}"), 72)
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title=f"13.{scenario.index} {title}",
        subtitle=(
            f"Análisis causal de {method_spec.entity_singular}: "
            f"Escenario #{scenario.index} · {context.period_label}"
        ),
    )
    chain_statement = (
        f"{_fmt_count_or_nd(row.get('linked_incidents', 0))} incidencias Helix y "
        f"{_fmt_count_or_nd(row.get('linked_comments', 0))} comentarios convergen en "
        f"{title} con lectura {method_spec.label}."
    )
    _panel(
        slide,
        left=0.66,
        top=1.48,
        width=12.02,
        height=0.86,
        title="Lectura ejecutiva",
        subtitle=chain_statement,
        fill=BBVA_COLORS["white"],
        border=BBVA_COLORS["sky"],
        title_size=13,
    )

    evidence = (
        scenario.helix_evidence_lines[: EDITORIAL_LIMITS.max_helix_evidence]
        or scenario.incident_lines[: EDITORIAL_LIMITS.max_helix_evidence]
        or ["No se han encontrado evidencias Helix adicionales defendibles para este escenario."]
    )
    card_width = 5.76
    card_height = 1.34 if len(evidence) > 2 else 1.46
    positions = [(0.66, 2.52), (6.92, 2.52), (0.66, 4.12), (6.92, 4.12)]
    for idx, line in enumerate(evidence[:4]):
        left, top = positions[idx]
        _panel(
            slide,
            left=left,
            top=top,
            width=card_width,
            height=card_height,
            title=f"Evidencia {idx + 1}",
            fill=BBVA_COLORS["white"],
            border=BBVA_COLORS["line"],
            title_size=12,
        )
        tb = slide.shapes.add_textbox(
            Inches(left + 0.18),
            Inches(top + 0.44),
            Inches(card_width - 0.36),
            Inches(card_height - 0.58),
        )
        tf = tb.text_frame
        _configure_text_frame(tf)
        tf.clear()
        p = tf.paragraphs[0]
        _add_markdown_runs(
            p,
            text=line,
            max_chars=152,
            font_size_pt=10.8,
            color=BBVA_COLORS["muted"],
        )

    visible_kpis = scenario.kpis[:4]
    kpi_left = 0.66
    kpi_top = 6.02
    comment_left = 9.08
    comment_width = 3.60
    available_width = comment_left - kpi_left - 0.20
    kpi_count = max(len(visible_kpis), 1)
    kpi_width = min(2.80, available_width / kpi_count - 0.10)
    for pos, (label, value, accent) in enumerate(visible_kpis):
        _add_stat_card(
            slide,
            left=kpi_left + pos * (kpi_width + 0.15),
            top=kpi_top,
            width=kpi_width,
            height=1.16,
            label=label,
            value=value,
            accent=accent,
        )
    _add_bullet_lines(
        slide,
        left=comment_left,
        top=kpi_top,
        width=comment_width,
        height=1.16,
        title="Comentarios enlazados",
        lines=scenario.comment_lines
        or ["No se han encontrado verbatims adicionales para este escenario."],
        accent=BBVA_COLORS["red"],
        body_font_size_pt=9.6,
    )


def _add_causal_evidence_slide(
    prs: Presentation,
    *,
    context: PresentationContext,
    scenario: CausalScenarioViewModel,
) -> None:
    row = scenario.row
    title = _clip(row.get("nps_topic", f"Escenario {scenario.index}"), 72)
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    _add_header(
        slide,
        title=f"14.{scenario.index} Detalle de evidencias Helix",
        subtitle=f"Evidencias con mayor vinculación causal para {title} · {context.period_label}",
    )
    evidence = scenario.helix_evidence_lines[: EDITORIAL_LIMITS.max_helix_evidence]
    if not evidence:
        evidence = [
            "No se han encontrado evidencias Helix adicionales defendibles para este escenario."
        ]
    card_width = 5.76
    card_height = 1.44
    positions = [(0.66, 1.62), (6.92, 1.62), (0.66, 3.42), (6.92, 3.42)]
    for idx, line in enumerate(evidence[:4]):
        left, top = positions[idx]
        _panel(
            slide,
            left=left,
            top=top,
            width=card_width,
            height=card_height,
            title=f"Evidencia {idx + 1}",
            fill=BBVA_COLORS["white"],
            border=BBVA_COLORS["line"],
            title_size=12,
        )
        tb = slide.shapes.add_textbox(
            Inches(left + 0.18), Inches(top + 0.46), Inches(card_width - 0.36), Inches(0.80)
        )
        tf = tb.text_frame
        _configure_text_frame(tf)
        tf.clear()
        p = tf.paragraphs[0]
        _add_markdown_runs(
            p,
            text=line,
            max_chars=150,
            font_size_pt=11.2,
            color=BBVA_COLORS["muted"],
        )
    _add_bullet_lines(
        slide,
        left=0.66,
        top=5.46,
        width=12.02,
        height=1.10,
        title="Criterio de selección",
        lines=[
            "Se priorizan evidencias enlazadas al escenario causal por ranking de prioridad, confianza, score causal y volumen de links."
        ],
        accent=BBVA_COLORS["sky"],
        body_font_size_pt=11.6,
    )


def _add_chain_scenario_slide(
    prs: Presentation,
    *,
    chain_row: pd.Series,
    idx: int,
    focus_name: str,
    period_label: str,
) -> None:
    slide = _new_slide(prs)
    _add_bg(slide, BBVA_COLORS["bg_light"])
    method_spec = get_causal_method_spec(str(chain_row.get("presentation_mode", "") or "").strip())
    title = _clip(chain_row.get("nps_topic", f"Cadena {idx}"), 70)
    linked_incidents = int(_safe_int(chain_row.get("linked_incidents", 0), default=0))
    linked_comments = int(_safe_int(chain_row.get("linked_comments", 0), default=0))
    touchpoint = _clip(chain_row.get("touchpoint", "Touchpoint"), 36)
    anchor_topic = _clip(chain_row.get("anchor_topic", ""), 52)
    owner_role = _clip(chain_row.get("owner_role", "n/d"), 24)
    focus_label = _focus_risk_label(focus_name)
    _add_header(
        slide,
        title=f"9.{idx} Análisis causal",
        subtitle=f"{method_spec.entity_singular}: {title} · {period_label}",
    )
    chain_statement = (
        f"{linked_incidents} incidencias de Helix y {linked_comments} comentarios de cliente "
        f"convergen en '{title}' bajo la lectura causal {method_spec.label}. "
        f"Touchpoint: {touchpoint}. "
        f"{f'Tópico NPS ancla: {anchor_topic}. ' if anchor_topic else ''}"
        f"Riesgo de {focus_label}. Equipo: {owner_role}."
    )
    banner = _panel(
        slide,
        left=0.66,
        top=1.48,
        width=12.0,
        height=0.80,
        title="Cadena de análisis",
        subtitle=chain_statement,
        fill=BBVA_COLORS["white"],
        border=BBVA_COLORS["sky"],
        title_size=13,
    )
    del banner

    metrics = [
        (
            _focus_probability_label(focus_name),
            _fmt_pct_or_nd(chain_row.get("detractor_probability", np.nan)),
            BBVA_COLORS["red"],
        ),
        (
            "Cambio esperado en NPS",
            _fmt_signed_or_nd(chain_row.get("nps_delta_expected", np.nan)),
            BBVA_COLORS["orange"],
        ),
        (
            "Impacto total",
            f"{_fmt_num_or_nd(chain_row.get('total_nps_impact', np.nan))} pts",
            BBVA_COLORS["blue"],
        ),
        (
            "Solidez de la evidencia",
            _fmt_num_or_nd(chain_row.get("confidence", np.nan)),
            BBVA_COLORS["green"],
        ),
        (
            "Vínculos validados",
            str(int(_safe_int(chain_row.get("linked_pairs", 0), default=0))),
            BBVA_COLORS["sky"],
        ),
        (
            "Prioridad sugerida",
            _fmt_num_or_nd(chain_row.get("priority", np.nan)),
            BBVA_COLORS["red"],
        ),
        (
            "NPS en riesgo",
            f"{_fmt_num_or_nd(chain_row.get('nps_points_at_risk', np.nan))} pts",
            BBVA_COLORS["red"],
        ),
        (
            "NPS recuperable",
            f"{_fmt_num_or_nd(chain_row.get('nps_points_recoverable', np.nan))} pts",
            BBVA_COLORS["green"],
        ),
    ]
    for pos, (label, value, accent) in enumerate(metrics):
        row = pos // 4
        col = pos % 4
        _add_stat_card(
            slide,
            left=0.66 + col * 3.02,
            top=2.52 + row * 1.24,
            width=2.82,
            height=1.02,
            label=label,
            value=value,
            accent=accent,
        )

    _add_bullet_lines(
        slide,
        left=0.66,
        top=5.12,
        width=5.75,
        height=1.78,
        title="Incidencias relacionadas",
        lines=[
            _clean_evidence_excerpt(line, max_len=112)
            for line in _chain_list(chain_row.get("incident_examples"))[:3]
        ]
        or ["No se han encontrado evidencias Helix adicionales para este escenario."],
        accent=BBVA_COLORS["orange"],
    )
    _add_bullet_lines(
        slide,
        left=6.63,
        top=5.12,
        width=6.03,
        height=1.78,
        title="Comentarios de cliente",
        lines=[
            _clean_evidence_excerpt(line, max_len=116)
            for line in _chain_list(chain_row.get("comment_examples"))[:2]
        ]
        or ["No se han encontrado verbatims adicionales para este escenario."],
        accent=BBVA_COLORS["red"],
    )


def generate_business_review_ppt(
    *,
    service_origin: str,
    service_origin_n1: str,
    service_origin_n2: str,
    period_start: date,
    period_end: date,
    focus_name: str,
    overall_weekly: pd.DataFrame,
    rationale_df: pd.DataFrame,
    nps_points_at_risk: float,
    nps_points_recoverable: float,
    top3_incident_share: float,
    median_lag_weeks: float,
    story_md: str,
    script_8slides_md: str,
    attribution_df: Optional[pd.DataFrame] = None,
    ranking_df: Optional[pd.DataFrame] = None,
    by_topic_daily: Optional[pd.DataFrame] = None,
    lag_days_by_topic: Optional[pd.DataFrame] = None,
    by_topic_weekly: Optional[pd.DataFrame] = None,
    lag_weeks_by_topic: Optional[pd.DataFrame] = None,
    template_name: str = "Plantilla corporativa fija v1",
    corporate_fixed: bool = True,
    logo_path: Optional[Path] = None,
    selected_nps_df: Optional[pd.DataFrame] = None,
    comparison_nps_df: Optional[pd.DataFrame] = None,
    template_path: Optional[Path] = None,
    incident_evidence_df: Optional[pd.DataFrame] = None,
    changepoints_by_topic: Optional[pd.DataFrame] = None,
    incident_timeline_df: Optional[pd.DataFrame] = None,
    hotspot_focus_note: str = "",
    touchpoint_source: str = "",
    entity_summary_df: Optional[pd.DataFrame] = None,
    entity_summary_kpis: Optional[list[dict[str, str]]] = None,
    executive_journey_catalog: Optional[list[dict[str, object]]] = None,
    broken_journeys_df: Optional[pd.DataFrame] = None,
) -> BusinessPptResult:
    """Build a business deck aligned to the selected period and BBVA corporate template."""
    del (
        script_8slides_md,
        template_name,
        corporate_fixed,
        logo_path,
        incident_evidence_df,
        incident_timeline_df,
        hotspot_focus_note,
        median_lag_weeks,
        rationale_df,
        ranking_df,
        by_topic_daily,
        lag_days_by_topic,
        by_topic_weekly,
        lag_weeks_by_topic,
        changepoints_by_topic,
        executive_journey_catalog,
    )

    prs = build_presentation(template_path=template_path, workspace_root=Path.cwd())
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    context = _build_presentation_context(
        service_origin=service_origin,
        service_origin_n1=service_origin_n1,
        service_origin_n2=service_origin_n2,
        period_start=period_start,
        period_end=period_end,
        focus_name=focus_name,
        overall_weekly=overall_weekly,
        story_md=story_md,
        attribution_df=attribution_df,
        selected_nps_df=selected_nps_df,
        comparison_nps_df=comparison_nps_df,
        touchpoint_source=touchpoint_source,
        entity_summary_df=entity_summary_df,
        entity_summary_kpis=entity_summary_kpis,
        broken_journeys_df=broken_journeys_df,
    )

    _add_cover_slide(
        prs,
        service_origin=context.service_origin,
        service_origin_n1=context.service_origin_n1,
        service_origin_n2=context.service_origin_n2,
        period_start=context.period_start,
        period_end=context.period_end,
        overview=context.overview,
        story_md=context.story_md,
    )
    _add_nps_section_cover_slide(prs, context=context)
    _add_overview_slide(
        prs,
        service_origin=context.service_origin,
        service_origin_n1=context.service_origin_n1,
        period_label=context.period_label,
        period_end=context.period_end,
        overview=context.overview,
        selected_nps_df=selected_nps_df,
        period_days=context.period_days,
        overview_figure=context.overview_figure,
    )
    _add_deep_dive_slide(
        prs,
        period_label=context.period_label,
        text_topics_df=context.text_topics_df,
        topic_figure=context.text_topic_figure,
    )
    _add_dimension_change_slide(prs, context=context, view_model=context.dimensions["Palanca"])
    _add_dimension_change_slide(prs, context=context, view_model=context.dimensions["Subpalanca"])
    _add_web_pain_dimension_slide(
        prs, context=context, view_model=context.dimensions["Palanca"], slide_number=7
    )
    _add_web_pain_dimension_slide(
        prs, context=context, view_model=context.dimensions["Subpalanca"], slide_number=8
    )
    _add_opportunity_dimension_slide(
        prs, context=context, view_model=context.dimensions["Palanca"], slide_number=9
    )
    _add_opportunity_dimension_slide(
        prs, context=context, view_model=context.dimensions["Subpalanca"], slide_number=10
    )
    _add_causal_section_cover_slide(
        prs,
        context=context,
        nps_points_at_risk=nps_points_at_risk,
        nps_points_recoverable=nps_points_recoverable,
        top3_incident_share=top3_incident_share,
    )
    _add_journeys_summary_slide(
        prs,
        period_label=context.period_label,
        touchpoint_source=context.causal.touchpoint_source,
        entity_summary_df=context.causal.entity_summary_df,
        entity_summary_kpis=context.causal.entity_summary_kpis,
        entity_summary_figure=context.causal.entity_summary_figure,
        journey_table_df=context.causal.journey_table_df,
    )
    for scenario in context.causal.scenarios:
        _add_causal_analysis_slide(prs, context=context, scenario=scenario)

    stamp = datetime.now(timezone.utc).strftime("%Y%m%d-%H%M")
    file_name = f"nps-incidencias-{_slug(service_origin)}-{_slug(service_origin_n1)}-{stamp}.pptx"

    buff = BytesIO()
    prs.save(buff)
    return BusinessPptResult(
        file_name=file_name, content=buff.getvalue(), slide_count=len(prs.slides)
    )
