from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable, Optional
import os
import zipfile

from pptx import Presentation


@dataclass(frozen=True)
class CorporatePresentationTheme:
    display_font: str = "Tiempos Headline"
    heading_font: str = "BentonSansBBVA Bold"
    body_font: str = "BentonSansBBVA Book"
    medium_font: str = "BentonSansBBVA Medium"
    cover_layout: str = "Portada Sand"
    content_layout: str = "sin logo y sin mosca"
    section_layout: str = "Portadilla 1"


_TEMPLATE_ENV_VARS = (
    "NPS_LENS_PPT_TEMPLATE",
    "NPS_LENS_BBVA_PPT_TEMPLATE",
)

_TEMPLATE_FILE_NAMES = (
    "SPHERICA-Plantilla-BBVA-Banca-Empresas-e-Instituciones-16-9.potx",
    "SPHERICA-Plantilla-BBVA-Banca-Empresas-e-Instituciones-16-9.pptx",
)


def find_corporate_template_path(
    explicit_path: Optional[Path] = None,
    *,
    workspace_root: Optional[Path] = None,
) -> Optional[Path]:
    candidates: list[Path] = []

    if explicit_path is not None:
        candidates.append(Path(explicit_path).expanduser())

    for env_name in _TEMPLATE_ENV_VARS:
        env_value = os.environ.get(env_name, "").strip()
        if env_value:
            candidates.append(Path(env_value).expanduser())

    roots = [
        workspace_root or Path.cwd(),
        Path.home() / "Downloads",
        Path.home() / "Documents",
    ]
    for root in roots:
        for name in _TEMPLATE_FILE_NAMES:
            candidates.append(root / name)
            candidates.append(root / "assets" / "ppt" / "bbva" / name)
            candidates.append(root / "assets" / "ppt" / "templates" / name)

    seen: set[str] = set()
    for candidate in candidates:
        key = str(candidate)
        if key in seen:
            continue
        seen.add(key)
        if candidate.exists() and candidate.is_file():
            return candidate
    return None


def build_presentation(
    *,
    template_path: Optional[Path] = None,
    workspace_root: Optional[Path] = None,
) -> Presentation:
    resolved = find_corporate_template_path(template_path, workspace_root=workspace_root)
    if resolved is None:
        return Presentation()

    payload = _normalized_template_payload(resolved)
    prs = Presentation(payload)
    clear_all_slides(prs)
    return prs


def clear_all_slides(prs: Presentation) -> None:
    for sld_id in list(prs.slides._sldIdLst):  # pyright: ignore[reportPrivateUsage]
        prs.part.drop_rel(sld_id.rId)
        prs.slides._sldIdLst.remove(sld_id)  # pyright: ignore[reportPrivateUsage]


def resolve_layout(
    prs: Presentation,
    preferred_names: Iterable[str],
    *,
    fallback_index: int = 0,
):
    wanted = [str(name).strip().lower() for name in preferred_names if str(name).strip()]
    for layout in prs.slide_layouts:
        if str(layout.name or "").strip().lower() in wanted:
            return layout
    return prs.slide_layouts[int(fallback_index)]


def _normalized_template_payload(path: Path) -> BytesIO:
    src = Path(path).expanduser()
    data = BytesIO()
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(data, "w") as zout:
        for item in zin.infolist():
            raw = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                raw = raw.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                )
            zout.writestr(item, raw)
    data.seek(0)
    return data
